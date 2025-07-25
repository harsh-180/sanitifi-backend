import warnings
warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl")
from .log_utils import log_user_action
from django.core.mail import send_mail
import re
import os
import subprocess
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import plotly
import plotly.express as px
import plotly.graph_objects as go
import scipy
import sklearn
import statsmodels
import multiprocessing
from multiprocessing import Pool, cpu_count
from tqdm import tqdm
from pyspark.sql import SparkSession
from pyspark.sql.functions import col, lower, regexp_replace
from pyspark.sql.utils import AnalysisException
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import openpyxl
import pandas as pd
import uuid
import time
from io import BytesIO
from datetime import datetime  # Added datetime import
from django.utils import timezone
from rest_framework.response import Response
from rest_framework.views import APIView
from rest_framework_simplejwt.tokens import RefreshToken
from django.contrib.auth import authenticate
from api.models import User,Projects, SavedScript, SavedPlot, SavedPivot, SavedPivotPlot, ProjectShare
from rest_framework.parsers import MultiPartParser, FormParser
from django.shortcuts import get_object_or_404

from django.conf import settings
from django.core.files.storage import default_storage
from django.core.files.base import ContentFile
from .models import Projects, User

from .auth import CustomAuthBackend
import shutil

from openpyxl import load_workbook

import pyarrow

from concurrent.futures import ThreadPoolExecutor, as_completed

from django.core.files.storage import default_storage
from .spark_utils import get_spark_session
import threading
import json
from django.http import FileResponse
from django.http import HttpResponse
# Add this at the top of views.py. a global dictionary to hold cleaned data for each session or file:
# cleaned_files_cache = {}
from django.views.decorators.csrf import csrf_exempt
from django.utils.decorators import method_decorator
from django.http import StreamingHttpResponse
import subprocess
import traceback
import stat
import csv
import io
   
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference, PieChart, LineChart
from openpyxl.utils.dataframe import dataframe_to_rows
from rest_framework.parsers import JSONParser

from django.http import HttpResponse
import tempfile
import math

import math

from google.oauth2 import service_account
from googleapiclient.discovery import build

GOOGLE_SCOPES = [
    'https://www.googleapis.com/auth/spreadsheets',
    'https://www.googleapis.com/auth/drive'
]

def get_gsheet_service():
    # Get service account info from environment variable
    service_account_info = os.getenv('GOOGLE_SERVICE_ACCOUNT_INFO')
    if not service_account_info:
        raise ValueError("GOOGLE_SERVICE_ACCOUNT_INFO environment variable not set")
    
    try:
        # Parse the JSON string from environment variable
        service_account_dict = json.loads(service_account_info)
        creds = service_account.Credentials.from_service_account_info(
            service_account_dict, scopes=GOOGLE_SCOPES)
        sheets_service = build('sheets', 'v4', credentials=creds)
        drive_service = build('drive', 'v3', credentials=creds)
        return sheets_service, drive_service
    except json.JSONDecodeError:
        raise ValueError("Invalid JSON in GOOGLE_SERVICE_ACCOUNT_INFO environment variable")
    except Exception as e:
        raise ValueError(f"Error initializing Google Sheets service: {str(e)}")

def make_json_safe(obj):
    if isinstance(obj, dict):
        return {k: make_json_safe(v) for k, v in obj.items()}
    elif isinstance(obj, list):
        return [make_json_safe(v) for v in obj]
    elif isinstance(obj, float):
        if math.isnan(obj) or math.isinf(obj):
            return None
        return obj
    elif isinstance(obj, (np.floating, np.integer)):
        # Handle numpy float/int types
        if np.isnan(obj) or np.isinf(obj):
            return None
        return obj.item() if hasattr(obj, 'item') else float(obj)
    else:
        return obj
    
def check_project_access(user_id, project_id, file_type=None, file_name=None, sheet_name=None):
    """
    Utility function to check if a user has access to a project or specific file.
    
    Args:
        user_id: ID of the user requesting access
        project_id: ID of the project
        file_type: Optional file type ('kpi' or 'media') for file-specific access
        file_name: Optional file name for file-specific access
        sheet_name: Optional sheet name for file-specific access
    
    Returns:
        tuple: (has_access, share_object, permission_level)
        - has_access: Boolean indicating if user has access
        - share_object: ProjectShare object if access exists, None otherwise
        - permission_level: Permission level if access exists, None otherwise
    """
    try:
        user = User.objects.get(id=user_id)
        project = Projects.objects.get(id=project_id)
        
        # Check if user is the project owner
        if project.user.id == int(user_id):
            return True, None, 'admin'
        
        # Check for shared access
        share_filter = {
            'project_id': project_id,
            'shared_with': user,
            'is_active': True
        }
        
        if file_type and file_name and sheet_name:
            # File-specific access with sheet name
            share_filter.update({
                'share_type': 'file',
                'file_type': file_type,
                'file_name': file_name,
                'sheet_name': sheet_name
            })
        elif file_type and file_name:
            # File-specific access without sheet name
            share_filter.update({
                'share_type': 'file',
                'file_type': file_type,
                'file_name': file_name
            })
        else:
            # Project-level access
            share_filter['share_type'] = 'project'
        
        try:
            share = ProjectShare.objects.get(**share_filter)
            return True, share, share.permission_level
        except ProjectShare.DoesNotExist:
            return False, None, None
            
    except (User.DoesNotExist, Projects.DoesNotExist):
        return False, None, None

@method_decorator(csrf_exempt, name='dispatch')

class MergeFile(APIView):
    def post(self, request):
        # Extract data from request  
        sheet_name = request.data.get('sheet_name')
        project_id = request.data.get('project_id')
        file_name = request.data.get('file_name')
        merge_file = request.FILES.get('merge_file')
        common_columns = request.data.get('common_columns')
        file_type = request.data.get('file_type')
        
        # file_type = file_name.split("\\")[0]
        file_name=file_name.split("\\")[1]

        print('Starting merge process...')
        print(f'Project ID: {project_id}')
        print(f'File Name: {file_name}')
        print(f'Sheet Name: {sheet_name}')
        
        # Clean up file_name if it contains media path
        if file_name and (file_name.startswith("media\\") or file_name.startswith("media/")):
            file_name = file_name.replace("media\\", "").replace("media/", "")

        # Ensure project exists
        try:
            project = Projects.objects.get(id=project_id)
            print(f'User ID: {project.user.id}')
        except Projects.DoesNotExist:
            return Response({'error': 'Project not found'}, status=404)

        try:
            # Correct KPI directory
            base_folder = os.path.join(settings.MEDIA_ROOT, f'user_{project.user.id}/project_{project.id}/{file_type}')
            print(f'Base folder: {base_folder}')
            
            # Determine the correct file path
            if os.path.isdir(os.path.join(base_folder, file_name)):
                # If file_name is a directory and sheet_name is a file
                original_file_path = os.path.join(base_folder, file_name, sheet_name)
            else:
                # If file_name is part of the file path
                original_file_path = os.path.join(base_folder, file_name, sheet_name)
            
            original_file_path = os.path.normpath(original_file_path)
            print(f'Original file path -> {original_file_path}')
            
            # Check if the file exists
            if not os.path.exists(original_file_path):
                return Response({'error': f'File not found at: {original_file_path}'}, status=404)
            
            print("2");

            # Load original and merge files
            original_df = pd.read_csv(original_file_path) if original_file_path.endswith('.csv') else pd.read_excel(original_file_path)
            merge_df = pd.read_excel(merge_file)

            print('Original File Columns:', original_df.columns.tolist())
            print('Merge File Columns:', merge_df.columns.tolist())
            print("3");

            # Standardizing column names (strip + lowercase)
            original_df.columns = original_df.columns.str.strip()
            merge_df.columns = merge_df.columns.str.strip()

            # Ensure common_columns is a proper list
            if isinstance(common_columns, str):
                try:
                    common_columns = eval(common_columns)
                except Exception as e:
                    return Response({'error': f'Invalid common_columns format: {str(e)}'}, status=400)

            if not isinstance(common_columns, list):
                return Response({'error': 'common_columns should be a list'}, status=400)

            common_columns = [col.strip() for col in common_columns]

            print('Common Columns for Merging:', common_columns)

            # Validate all common columns exist in both
            for col in common_columns:
                if col not in original_df.columns:
                    return Response({'error': f"Column '{col}' not found in original file"}, status=400)
                if col not in merge_df.columns:
                    return Response({'error': f"Column '{col}' not found in merge file"}, status=400)

            # Convert common columns to string to avoid merge issues
            for col in common_columns:
                original_df[col] = original_df[col].fillna('').astype(str)
                merge_df[col] = merge_df[col].fillna('').astype(str)
                print(f"Converted column '{col}' to string")

            # ---- MODIFIED MERGE LOGIC START ----

            # Get unique columns from each dataframe (excluding common columns)
            original_only_columns = [col for col in original_df.columns if col not in common_columns]
            merge_only_columns = [col for col in merge_df.columns if col not in common_columns]
            
            # Filter out any duplicates between original_only and merge_only (keeping original file's columns)
            merge_only_columns = [col for col in merge_only_columns if col not in original_only_columns]
            
            # Perform the merge on common columns
            merged_df = pd.merge(
                original_df[common_columns + original_only_columns], 
                merge_df[common_columns + merge_only_columns],
                how='outer', 
                on=common_columns
            )
            # ---- MODIFIED MERGE LOGIC END ----
            
            merged_df.drop_duplicates(inplace=True)
            merged_df = merged_df.replace({np.nan: None})

            # Get file directory for git commit
            file_dir = os.path.dirname(original_file_path)
            relative_file_dir = os.path.relpath(file_dir, start=os.path.join(settings.MEDIA_ROOT, f'user_{project.user.id}/project_{project.id}'))
            file_basename = os.path.basename(original_file_path)

            # Save back to original file (overwrite)
            if original_file_path.endswith('.csv'):
                merged_df.to_csv(original_file_path, index=False)
            else:
                merged_df.to_excel(original_file_path, index=False)

            # Git Commit
            self.commit_to_git(
                os.path.join(settings.MEDIA_ROOT, f'user_{project.user.id}/project_{project.id}'),
                project.user,
                project.id,
                relative_file_dir,
                file_basename
            )

            # Log the merge event
            user = get_logging_user(request, getattr(project, 'user', None))
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "merge_file", details=f"Merged file: {file_name}, sheet: {sheet_name}, type: {file_type}", ip_address=ip)

            return Response({
                "message": "Merged and updated the original file successfully",
                "merged_file": {
                    "columns": merged_df.columns.tolist(),
                    "data": merged_df.values.tolist()
                }
            }, status=200)

        except Exception as e:
            print(f'Error: {str(e)}')
            return Response({'error': str(e)}, status=500)

    def commit_to_git(self, project_folder, user, project_id, file_subfolder, sheet_name):
        try:
            if not os.path.exists(os.path.join(project_folder, ".git")):
                subprocess.run(["git", "init"], cwd=project_folder)
                subprocess.run(["git", "config", "user.name", user.name], cwd=project_folder)
                subprocess.run(["git", "config", "user.email", user.email], cwd=project_folder)

            file_path_relative = os.path.join(file_subfolder, sheet_name).replace("\\", "/")
            subprocess.run(["git", "add", file_path_relative], cwd=project_folder)
            commit_message = f"merge - {user.id}/{project_id}/{file_subfolder}/{sheet_name}"
            subprocess.run(["git", "commit", "-m", commit_message], cwd=project_folder)
            print(f"Git commit done for {file_path_relative}")
        except Exception as e:
            print(f"Git commit failed: {str(e)}")


class Save(APIView):
    def post(self, request):
        try:
            # Extract payload data
            sheet_name = request.data.get("sheet_name")  # e.g., "Jan - Dec (2).csv"
            project_id = request.data.get("project_id")
            file_name = request.data.get("file_name")  # e.g., "sample _sg_fronx"
            file_type = request.data.get("file_type")  # Expected: "media" or "kpi"
            data_stage = request.data.get("data_stage")  # Expected: "cleaned" or other stages

            print("\nReceived Save Request with Payload:")
            print(f"  - sheet_name: {sheet_name}")
            print(f"  - project_id: {project_id}")
            print(f"  - file_name: {file_name}")
            print(f"  - file_type: {file_type}")
            print(f"  - data_stage: {data_stage}")

            # Log the file download event
            user = get_logging_user(request, None)
            ip = request.META.get('REMOTE_ADDR')
            # Validate project existence
            try:
                project = Projects.objects.get(id=project_id)
                print(f"Project found: user_{project.user.id}/project_{project.id}")
            except Projects.DoesNotExist:
                return Response({"error": "Project not found"}, status=404)

            # If user is still None, use project.user as fallback
            if user is None and hasattr(project, 'user'):
                user = project.user

            log_user_action(user, "download_file", details=f"Download file: {file_name}, sheet: {sheet_name}, type: {file_type}", ip_address=ip)

            # Clean up file_name if it contains path information
            file_name = os.path.basename(file_name)
            if file_name and (file_name.startswith("media\\") or file_name.startswith("media/")):
                file_name = file_name.replace("media\\", "").replace("media/", "")

            # Validate file_type
            if file_type not in ["media", "kpi"]:
                return Response({"error": "Invalid file_type. Must be 'media' or 'kpi'"}, status=400)

            # Construct correct file path inside project folder using file_type directly
            file_path = os.path.join(
                settings.MEDIA_ROOT,
                f"user_{project.user.id}/project_{project.id}/{file_type}/{file_name}/{sheet_name}"
            )
            file_path = os.path.normpath(file_path)  # Normalize path
            print(f"Checking file at: {file_path}")

            # Check if the file exists
            if not os.path.exists(file_path):
                print(f"File not found: {file_path}")

                # Debugging - List actual files in the directory
                directory = os.path.dirname(file_path)
                if os.path.exists(directory):
                    print(f"Listing files in {directory}:")
                    print(os.listdir(directory))  # Print all files in the directory
                else:
                    print(f"Directory does not exist: {directory}")

                return Response({"error": "File not found", "expected_path": file_path}, status=404)

            # Extract file extension and validate
            file_extension = os.path.splitext(sheet_name)[-1].lower()
            print(f"Extracted file extension: {file_extension}")

            if file_extension not in [".xlsx", ".csv"]:
                return Response({"error": "Only XLSX and CSV files are supported"}, status=400)

            output = BytesIO()  # In-memory buffer for streaming response

            # Load the file based on its format
            if file_extension == ".xlsx":
                # Load all sheet names
                sheet_names = pd.ExcelFile(file_path, engine="openpyxl").sheet_names
                print(f"Available sheets: {sheet_names}")

                # Ensure the requested sheet exists
                if sheet_name not in sheet_names:
                    return Response({
                        "error": f"Sheet '{sheet_name}' not found in {actual_file_name}.",
                        "available_sheets": sheet_names
                    }, status=400)

                # Read the requested sheet
                df = pd.read_excel(file_path, sheet_name=sheet_name, engine="openpyxl")
            elif file_extension == ".csv":
                df = pd.read_csv(file_path)  # Let pandas infer types

            # Convert numeric-like columns back to numeric format
            for col in df.select_dtypes(include=["object"]).columns:
                try:
                    df[col] = pd.to_numeric(df[col], errors="ignore")
                except Exception:
                    pass

            # Save processed file in memory
            if file_extension == ".xlsx":
                df.to_excel(output, index=False, engine="openpyxl")
            else:
                df.to_csv(output, index=False)

            output.seek(0)  # Move buffer position to the beginning

            # Return file as a streaming response
            response = StreamingHttpResponse(output, content_type="application/octet-stream")
            response['Content-Disposition'] = f'attachment; filename="{sheet_name}"'
            response['Cache-Control'] = 'no-cache, no-store, must-revalidate'
            response['Pragma'] = 'no-cache'
            response['Expires'] = '0'
            return response

        except Exception as e:
            print(f"Error in Save API: {str(e)}")
            return Response({"error": f"Internal server error: {str(e)}"}, status=500)


class Mapping(APIView):
    def post(self, request):
        sheet_name = request.data.get("sheet_name")  
        project_id = request.data.get("project_id")
        file_name = request.data.get("file_name")  
        column_mappings = request.data.get("column_mappings")  
        mapping_file = request.FILES.get("mapping_file")
        file_type=request.data.get("file_type")
        file_name = os.path.basename(file_name)
        if not project_id or not file_name or not mapping_file or not column_mappings or not sheet_name:
            return Response({"error": "Missing required fields"}, status=400)
        try:
            project = Projects.objects.get(id=project_id)
        except Projects.DoesNotExist:
            return Response({"error": "Project not found"}, status=404)
        file_folder = os.path.join(settings.MEDIA_ROOT, f"user_{project.user.id}/project_{project.id}/{file_type}")
        original_file_path = os.path.join(file_folder, file_name, sheet_name)
        original_file_path = os.path.normpath(original_file_path)
        if not os.path.exists(original_file_path):
            return Response({"error": "File not found in KPI folder"}, status=404)
        file_ext = os.path.splitext(sheet_name)[1].lower()
        file_size_mb = os.path.getsize(original_file_path) / (1024 * 1024)
        try:
            if file_size_mb > 50:
                # Use the new robust Spark session management
                from .spark_utils import spark_session_context, validate_spark_session
                
                with spark_session_context() as spark:
                    # Validate session before use
                    if not validate_spark_session(spark):
                        raise Exception("Invalid Spark session. Please try again.")
                    
                    # Read file with optimized settings
                    if file_ext in ['.xlsx', '.xls']:
                        df = spark.read \
                            .format("com.crealytics.spark.excel") \
                            .option("header", True) \
                            .option("inferSchema", False) \
                            .option("maxRowsInMemory", 50000) \
                            .option("maxColumns", 20000) \
                            .load(original_file_path)
                    elif file_ext == '.csv':
                        df = spark.read \
                            .option("header", True) \
                            .option("inferSchema", False) \
                            .option("maxRowsInMemory", 50000) \
                            .option("maxColumns", 20000) \
                            .csv(original_file_path)
                    else:
                        return Response({"error": "Unsupported file type. Only .csv, .xls, .xlsx allowed."}, status=400)
                    
                    # Cache the DataFrame for better performance
                    df = df.cache()
                    
                    # Read mapping file
                    mapping_df = pd.read_excel(mapping_file)
                    
                    # Process column mappings with proper column escaping
                    column_mappings_dict = json.loads(column_mappings) if isinstance(column_mappings, str) else column_mappings
                    
                    for original_col, mapping_col in column_mappings_dict.items():
                        if original_col in df.columns and mapping_col in mapping_df.columns:
                            # Create mapping DataFrame
                            mapping_col_index = mapping_df.columns.get_loc(mapping_col)
                            next_col_index = mapping_col_index + 1
                            
                            if next_col_index < len(mapping_df.columns):
                                next_col_name = mapping_df.columns[next_col_index]
                                mapping_spark = spark.createDataFrame(mapping_df[[mapping_col, next_col_name]])
                                mapping_spark = mapping_spark.withColumnRenamed(mapping_col, "key").withColumnRenamed(next_col_name, "value")
                                
                                # Escape column names that contain spaces or special characters
                                escaped_original_col = f"`{original_col}`" if ' ' in original_col or any(char in original_col for char in ['-', '.', '/', '\\']) else original_col
                                escaped_new_col = f"`{original_col}_New`" if ' ' in f"{original_col}_New" or any(char in f"{original_col}_New" for char in ['-', '.', '/', '\\']) else f"{original_col}_New"
                                
                                # Perform the join with proper column escaping
                                df = df.join(mapping_spark, df[original_col] == mapping_spark["key"], "left") \
                                       .withColumn(escaped_new_col, expr(f"coalesce(value, {escaped_original_col})")) \
                                       .drop("key", "value")
                            else:
                                print(f"Warning: No next column exists after '{mapping_col}' in mapping_df")
                    
                    # Fill NA values
                    df = df.na.fill("NA")
                    
                    # Save the result
                    df.toPandas().to_csv(original_file_path, index=False)
                    result_df = df.limit(50).toPandas()
                    
                    # Uncache to free memory
                    df.unpersist()
            else:
                if file_ext in ['.xlsx', '.xls']:
                    df = pd.read_excel(original_file_path)
                elif file_ext == '.csv':
                    df = pd.read_csv(original_file_path)
                else:
                    return Response({"error": "Unsupported file type. Only .csv, .xls, .xlsx allowed."}, status=400)
                mapping_df = pd.read_excel(mapping_file)
                if isinstance(column_mappings, str):
                    try:
                        column_mappings = json.loads(column_mappings)
                    except json.JSONDecodeError:
                        return Response({"error": "Invalid JSON format for column_mappings"}, status=400)
                for original_col, mapping_col in column_mappings.items():
                    try:
                        new_col_name = f"{original_col}_New"
                        if original_col in df.columns and mapping_col in mapping_df.columns:
                            mapping_col_index = mapping_df.columns.get_loc(mapping_col)
                            next_col_index = mapping_col_index + 1
                            if next_col_index < len(mapping_df.columns):
                                next_col_name = mapping_df.columns[next_col_index]
                                mapping_dict = dict(zip(mapping_df[mapping_col], mapping_df[next_col_name]))
                                df[new_col_name] = df[original_col].map(mapping_dict).fillna(df[original_col])
                            else:
                                print(f"Warning: No next column exists after '{mapping_col}' in mapping_df")
                    except Exception as e:
                        return Response({"error": f"Error processing column mapping: {str(e)}"}, status=500)
                df.replace([np.inf, -np.inf], np.nan, inplace=True)
                df.fillna("NA", inplace=True)
                if file_ext in ['.xlsx', '.xls']:
                    df.to_excel(original_file_path, index=False)
                elif file_ext == '.csv':
                    df.to_csv(original_file_path, index=False)
                result_df = df
            self.commit_to_git(
                os.path.join(settings.MEDIA_ROOT, f"user_{project.user.id}/project_{project.id}"),
                project.user,
                project.id,
                file_type,
                file_name,
                sheet_name or ''
            )
            user = None
            if hasattr(request, 'user') and request.user.is_authenticated:
                user = request.user
            elif project and hasattr(project, 'user'):
                user = project.user
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "mapping_file", details=f"Mapped file: {file_name}, sheet: {sheet_name}, type: {file_type}", ip_address=ip)

            return Response({
                "message": "Mapped and overwritten the original file successfully",
                "mapped_file": {
                    "columns": result_df.columns.tolist(),
                    "data": result_df.values.tolist()
                }
            })
        except Exception as e:
            return Response({"error": f"Unexpected error: {str(e)}"}, status=500)

    def commit_to_git(self, project_folder, user, project_id, file_type, file_name, sheet_name):
        try:
            # Initialize git repo if it doesn't exist
            if not os.path.exists(os.path.join(project_folder, ".git")):
                subprocess.run(["git", "init"], cwd=project_folder)
                subprocess.run(["git", "config", "user.name", user.name], cwd=project_folder)
                subprocess.run(["git", "config", "user.email", user.email], cwd=project_folder)

            # Make relative path to the file/folder
            file_path_relative = os.path.join(file_type, file_name)

            # Stage the changes
            subprocess.run(["git", "add", file_path_relative], cwd=project_folder)

            # Commit message for mapping action
            commit_message = f"mapping - {user.id}/{project_id}/{file_type}/{file_name}/{sheet_name}"
            subprocess.run(["git", "commit", "-m", commit_message], cwd=project_folder)

        except Exception as e:
            print(f"Git commit failed: {str(e)}")



class Melting(APIView):
    def post(self, request):
        project_id = request.data.get("project_id")
        user_id = request.data.get("user_id")
        file_name = request.data.get("file_name")
        id_vars = request.data.get("id_vars")
        value_vars = request.data.get("value_vars")
        var_name = request.data.get("var_name", "Variable")
        value_name = request.data.get("value_name", "Value")
        sheet_name = request.data.get("sheet_name")
        file_type = request.data.get("file_type")
        file_name = os.path.basename(file_name)
        if not project_id or not file_name:
            return Response({"error": "Missing required fields"}, status=400)
        try:
            project = Projects.objects.get(id=project_id)
        except Projects.DoesNotExist:
            return Response({"error": "Project not found"}, status=404)
        user = project.user
        project_base_folder = os.path.join(settings.MEDIA_ROOT, f"user_{user.id}/project_{project_id}")
        file_folder = os.path.join(project_base_folder, file_type)
        original_file_path = os.path.join(file_folder, file_name, sheet_name)
        if not os.path.exists(original_file_path):
            return Response({"error": f"Original file not found in {file_type} folder"}, status=404)
        file_size_mb = os.path.getsize(original_file_path) / (1024 * 1024)
        try:
            if file_size_mb > 50:
                # Use the new robust Spark session management
                from .spark_utils import spark_session_context, validate_spark_session
                
                with spark_session_context() as spark:
                    # Validate session before use
                    if not validate_spark_session(spark):
                        raise Exception("Invalid Spark session. Please try again.")
                    
                    # Read CSV with optimized settings
                    df = spark.read \
                        .option("header", True) \
                        .option("inferSchema", False) \
                        .option("maxRowsInMemory", 50000) \
                        .option("maxColumns", 20000) \
                        .csv(original_file_path)
                    
                    # Cache the DataFrame for better performance
                    df = df.cache()
                    
                    # Handle comma-separated input
                    if isinstance(id_vars, str):
                        id_vars = [col.strip() for col in id_vars.split(",")]
                    if isinstance(value_vars, str):
                        value_vars = [col.strip() for col in value_vars.split(",")]
                    
                    # Validate that all columns exist
                    all_columns = df.columns
                    missing_id_vars = [col for col in id_vars if col not in all_columns]
                    missing_value_vars = [col for col in value_vars if col not in all_columns]
                    
                    if missing_id_vars or missing_value_vars:
                        missing_cols = missing_id_vars + missing_value_vars
                        raise Exception(f"Columns not found in dataset: {', '.join(missing_cols)}")
                    
                    # Create the melt operation using proper column escaping
                    n = len(value_vars)
                    
                    # Build the stack expression with proper column escaping
                    stack_parts = []
                    for v in value_vars:
                        # Escape column names that contain spaces or special characters
                        escaped_col = f"`{v}`" if ' ' in v or any(char in v for char in ['-', '.', '/', '\\']) else v
                        stack_parts.append(f"'{v}', {escaped_col}")
                    
                    exprs = ", ".join(stack_parts)
                    
                    # Build the select expression with proper column escaping
                    select_cols = []
                    for col in id_vars:
                        # Escape column names that contain spaces or special characters
                        escaped_col = f"`{col}`" if ' ' in col or any(char in col for char in ['-', '.', '/', '\\']) else col
                        select_cols.append(escaped_col)
                    
                    # Add the stack expression
                    select_cols.append(f"stack({n}, {exprs}) as ({var_name}, {value_name})")
                    
                    # Perform the melt operation
                    melted = df.selectExpr(*select_cols)
                    melted = melted.na.fill("NA")
                    
                    # Save the result
                    melted.toPandas().to_csv(original_file_path, index=False)
                    melted_df = melted.limit(50).toPandas()
                    
                    # Uncache to free memory
                    df.unpersist()
            else:
                try:
                    df = pd.read_csv(original_file_path)
                except UnicodeDecodeError:
                    df = pd.read_csv(original_file_path, encoding='latin1')
                if isinstance(id_vars, str):
                    id_vars = [col.strip() for col in id_vars.split(",")]
                if isinstance(value_vars, str):
                    value_vars = [col.strip() for col in value_vars.split(",")]
                melted_df = df.melt(id_vars=id_vars, value_vars=value_vars, var_name=var_name, value_name=value_name)
                melted_df.replace([np.inf, -np.inf], np.nan, inplace=True)
                melted_df.fillna("NA", inplace=True)
                melted_df.to_csv(original_file_path, index=False)
            self.commit_to_git(project_base_folder, user, project_id, file_type, file_name, sheet_name)
            user = None
            if hasattr(request, 'user') and request.user.is_authenticated:
                user = request.user
            elif project and hasattr(project, 'user'):
                user = project.user
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "melt_file", details=f"Melted file: {file_name}, sheet: {sheet_name}, type: {file_type}", ip_address=ip)

            return Response({
                "message": "Melted and overwritten the original file successfully",
                "melted_file": {
                    "columns": melted_df.columns.tolist(),
                    "data": melted_df.values.tolist()
                }
            }, status=200)
        except Exception as e:
            return Response({"error": f"Error during melting: {str(e)}"}, status=500)

    def commit_to_git(self, project_folder, user, project_id, file_type, file_name, sheet_name):
        try:
            if not os.path.exists(os.path.join(project_folder, ".git")):
                subprocess.run(["git", "init"], cwd=project_folder)
                subprocess.run(["git", "config", "user.name", user.name], cwd=project_folder)
                subprocess.run(["git", "config", "user.email", user.email], cwd=project_folder)

            file_path_relative = os.path.join(file_type, file_name)
            subprocess.run(["git", "add", file_path_relative], cwd=project_folder)
            commit_message = f"melted - {user.id}/{project_id}/{file_type}/{file_name}/{sheet_name}"
            subprocess.run(["git", "commit", "-m", commit_message], cwd=project_folder)
        except Exception as e:
            print(f"Git commit failed: {str(e)}")


class SheetInfo(APIView):
    def post(self, request):
        try:
            # Extract payload data
            file_type = request.data.get('file_type')
            file_name = request.data.get('file_name')
            project_id = request.data.get('project_id')
            sheet_name = request.data.get('sheet_name')

            file_name = os.path.basename(file_name)

            if not all([file_type, file_name, project_id, sheet_name]):
                return Response({"error": "Missing required fields"}, status=400)

            # Get project
            try:
                project = Projects.objects.get(id=project_id)
            except Projects.DoesNotExist:
                return Response({"error": "Project not found"}, status=404)

            # Construct file path
            if file_type == 'concatenated':
                # Search all subfolders for the sheet_name
                concatenated_base = os.path.join(settings.MEDIA_ROOT, f"user_{project.user.id}/project_{project.id}/concatenated")
                found = False
                file_path = None
                if os.path.exists(concatenated_base):
                    for folder in os.listdir(concatenated_base):
                        folder_path = os.path.join(concatenated_base, folder)
                        candidate = os.path.join(folder_path, sheet_name)
                        if os.path.isfile(candidate):
                            file_path = candidate
                            found = True
                            break
                if not found:
                    return Response({"error": "File not found"}, status=404)
            else:
                file_path = os.path.join(
                    settings.MEDIA_ROOT,
                    f"user_{project.user.id}/project_{project.id}/{file_type}/{file_name}/{sheet_name}"
                )
                file_path = os.path.normpath(file_path)
                if not os.path.exists(file_path):
                    return Response({"error": "File not found"}, status=404)

            # Read the file
            file_extension = os.path.splitext(sheet_name)[1].lower()
            if file_extension == '.csv':
                try:
                    df = pd.read_csv(file_path, encoding='utf-8')
                except UnicodeDecodeError:
                    df = pd.read_csv(file_path, encoding='latin1')
            elif file_extension in ['.xlsx', '.xls']:
                df = pd.read_excel(file_path)
            else:
                return Response({"error": "Unsupported file format"}, status=400)

            # Replace infinite values with None
            df = df.replace([np.inf, -np.inf], np.nan)
            
            column_info = {}

            for column in df.columns:
                column_data = df[column]
                null_count = column_data.isnull().sum()
                blank_count = (column_data == '').sum() if column_data.dtype == 'object' else 0
                total_empty = null_count + blank_count

                # Get unique values safely
                unique_values = column_data.dropna().unique()
                unique_values_list = make_json_safe(unique_values.tolist()[:100])  # limit to 100 unique elements

                column_info[column] = {
                    "null_count": int(total_empty),
                    "data_type": str(column_data.dtype),
                    "unique_values": int(column_data.nunique()),
                    "unique_elements": unique_values_list
                }

                # Check if numeric
                try:
                    numeric_data = pd.to_numeric(column_data, errors='coerce')
                    if not numeric_data.isnull().all():
                        # Handle infinite values in numeric calculations
                        numeric_data = numeric_data.replace([np.inf, -np.inf], np.nan)
                        min_val = numeric_data.min()
                        max_val = numeric_data.max()
                        mean_val = numeric_data.mean()
                        
                        column_info[column].update({
                            "type": "numeric",
                            "min": make_json_safe(min_val),
                            "max": make_json_safe(max_val),
                            "average": make_json_safe(mean_val),
                            "unique_values": int(numeric_data.nunique())
                        })
                        continue
                except:
                    pass

                # Check if datetime
                try:
                    date_data = pd.to_datetime(column_data, errors='coerce', dayfirst=False)
                    if not date_data.isnull().all():
                        granularity = 'daily'
                        date_data_clean = date_data.dropna()
                        date_diff = date_data_clean.diff().median()

                        if date_diff.days >= 365:
                            granularity = 'yearly'
                            freq = 'Y'
                            existing_dates = set(date_data_clean.dt.year)
                            expected_dates = set(pd.date_range(start=date_data_clean.min(), end=date_data_clean.max(), freq=freq).year)
                            missing_dates = sorted(expected_dates - existing_dates)
                            missing_formatted = [str(y) for y in missing_dates]
                        elif date_diff.days >= 28:
                            granularity = 'monthly'
                            freq = 'M'
                            existing_dates = set(date_data_clean.dt.to_period('M'))
                            expected_dates = set(pd.date_range(start=date_data_clean.min(), end=date_data_clean.max(), freq=freq).to_period('M'))
                            missing_dates = sorted(expected_dates - existing_dates)
                            missing_formatted = [str(m) for m in missing_dates]
                        else:
                            granularity = 'daily'
                            freq = 'D'
                            existing_dates = set(date_data_clean.dt.normalize())
                            expected_dates = set(pd.date_range(start=date_data_clean.min(), end=date_data_clean.max(), freq=freq))
                            missing_dates = sorted(expected_dates - existing_dates)
                            missing_formatted = [d.strftime('%Y-%m-%d') for d in missing_dates]

                        column_info[column].update({
                            "type": "datetime",
                            "granularity": granularity,
                            "min_date": date_data_clean.min().strftime('%Y-%m-%d %H:%M:%S'),
                            "max_date": date_data_clean.max().strftime('%Y-%m-%d %H:%M:%S'),
                            "missing_dates": missing_formatted
                        })
                        continue
                except:
                    pass

                # Fallback to string
                column_info[column].update({
                    "type": "string",
                    "unique_values": int(column_data.nunique())
                })
            user = None
            if hasattr(request, 'user') and request.user.is_authenticated:
                user = request.user
            elif project and hasattr(project, 'user'):
                user = project.user
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "column_info", details=f"Column information retrieved successfully", ip_address=ip)

            return Response({
                "message": "Column information retrieved successfully",
                "columns": column_info
            }, status=200)

        except Exception as e:
            return Response({"error": f"Error analyzing sheet: {str(e)}"}, status=500)

class CleaningColumns(APIView):
    def post(self, request):
        file_type = request.data.get('file_type')
        file_name = request.data.get('file_name')
        project_id = request.data.get('project_id')
        sheet_name = request.data.get('sheet_name')  # will act as CSV file name here

        print(sheet_name)
        analysis_only = request.data.get('analysis_only', False)
        selected_columns = request.data.get('selected_columns', [])
        options = request.data.get('options', {})

        file_name = os.path.basename(file_name)
        sheet_name = os.path.basename(sheet_name)  # sheet_name = CSV filename now

        if not all([file_type, file_name, project_id, sheet_name]):
            return Response({'error': 'Missing required fields'}, status=400)

        try:
            project = Projects.objects.get(id=project_id)
        except Projects.DoesNotExist:
            return Response({'error': 'Project not found'}, status=404)

        user = project.user
        project_folder = os.path.join(settings.MEDIA_ROOT, f"user_{user.id}/project_{project.id}")
        csv_folder_path = os.path.join(project_folder, file_type, file_name)
        csv_file_path = os.path.join(csv_folder_path, sheet_name)

        if not default_storage.exists(csv_file_path):
            return Response({'error': 'CSV file not found'}, status=404)

        local_path = os.path.join(settings.MEDIA_ROOT, csv_file_path)
        file_size_mb = os.path.getsize(local_path) / (1024 * 1024)

        try:
            if file_size_mb > 50:
                # Use the new robust Spark session management
                from .spark_utils import spark_session_context, validate_spark_session
                
                with spark_session_context() as spark:
                    # Validate session before use
                    if not validate_spark_session(spark):
                        raise Exception("Invalid Spark session. Please try again.")
                    
                    # Read CSV with optimized settings
                    df = spark.read \
                        .option("header", True) \
                        .option("inferSchema", False) \
                        .option("maxRowsInMemory", 50000) \
                        .option("maxColumns", 20000) \
                        .csv(local_path)
                    
                    # Cache the DataFrame for better performance
                    df = df.cache()
                    
                    # Apply cleaning operations
                    if not analysis_only:
                        # Lowercase columns
                        for col_name in options.get("lowercase_columns", []):
                            if col_name in df.columns:
                                df = df.withColumn(col_name, lower(col(col_name)))
                        
                        # Remove spaces
                        for col_name in options.get("remove_spaces_columns", []):
                            if col_name in df.columns:
                                df = df.withColumn(col_name, regexp_replace(col(col_name), " ", ""))
                        
                        # Remove special chars
                        for col_name, char_map in options.get("remove_special_chars", {}).items():
                            if col_name in df.columns:
                                for char, remove in char_map.items():
                                    if remove:
                                        df = df.withColumn(col_name, regexp_replace(col(col_name), re.escape(char), ""))
                        
                        # Fill NA
                        df = df.na.fill("NA")
                        
                        # Save back to CSV with optimized settings
                        df.toPandas().to_csv(local_path, index=False)
                    
                    # Get preview data for response
                    cleaned_data = df.limit(50).toPandas()
                    
                    # Uncache to free memory
                    df.unpersist()
                    
                    # For analysis_only mode, we don't need special chars analysis with Spark
                    special_chars_analysis = {} if analysis_only else {}
                    
                    if analysis_only:
                        return Response({
                            'status': 'success',
                            'special_chars_analysis': special_chars_analysis,
                            'columns': cleaned_data.columns.tolist(),
                            'total_rows': df.count() if not analysis_only else None
                        }, status=200)
            else:
                with default_storage.open(csv_file_path, 'rb') as f:
                    df = pd.read_csv(f)
                special_chars_analysis = self.analyze_special_characters(df)
                if analysis_only:
                    return Response({
                        'status': 'success',
                        'special_chars_analysis': special_chars_analysis,
                        'columns': df.columns.tolist()
                    }, status=200)
                if selected_columns:
                    invalid_columns = [col for col in selected_columns if col not in df.columns]
                    if invalid_columns:
                        return Response({
                            'error': f'Invalid columns: {", ".join(invalid_columns)}',
                            'special_chars_analysis': special_chars_analysis
                        }, status=400)
                cleaned_data = df.copy()
                # Apply cleaning options
                for column in options.get("lowercase_columns", []):
                    if column in cleaned_data.columns:
                        cleaned_data[column] = cleaned_data[column].astype(str).str.lower()
                cleaned_data = self.remove_spaces(cleaned_data, options.get("remove_spaces_columns", []))
                cleaned_data = self.remove_special_characters(cleaned_data, options.get("remove_special_chars", {}))
                cleaned_data = cleaned_data.replace([np.inf, -np.inf], np.nan)
                cleaned_data = cleaned_data.fillna("NA")
                # Save cleaned CSV back
                with default_storage.open(csv_file_path, 'w') as f:
                    cleaned_data.to_csv(f, index=False, lineterminator='\n')
            self.commit_to_git(project_folder, user, project_id, file_type, file_name, sheet_name)
            user = None
            if hasattr(request, 'user') and request.user.is_authenticated:
                user = request.user
            elif project and hasattr(project, 'user'):
                user = project.user
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "clean_file", details=f"Cleaned file: {file_name}, sheet: {sheet_name}, type: {file_type}", ip_address=ip)
            return Response({
                'message': 'CSV file cleaned successfully',
                'cleaned_sheet': {
                    sheet_name: {
                        'columns': cleaned_data.columns.tolist(),
                        'data': cleaned_data.values.tolist()
                    }
                },
                'special_chars_analysis': special_chars_analysis
            }, status=200)
        except Exception as e:
            return Response({'error': str(e)}, status=500)

    def commit_to_git(self, project_folder, user, project_id, file_type, file_name, sheet_name):
        try:
            if not os.path.exists(os.path.join(project_folder, ".git")):
                subprocess.run(["git", "init"], cwd=project_folder)
                subprocess.run(["git", "config", "user.name", user.name], cwd=project_folder)
                subprocess.run(["git", "config", "user.email", user.email], cwd=project_folder)

            subprocess.run(["git", "add", "."], cwd=project_folder)
            commit_message = f"cleaned - {user.id}/{project_id}/{file_type}/{file_name}/{sheet_name}"
            subprocess.run(["git", "commit", "-m", commit_message], cwd=project_folder)
        except Exception as e:
            print(f"Git commit failed: {str(e)}")

    def analyze_special_characters(self, df):
        special_chars_analysis = {}
        for column in df.columns:
            values = df[column].astype(str).fillna('')
            special_chars = set(re.findall(r'[^a-zA-Z0-9]', ' '.join(values)))
            char_stats = {}
            total_char_count = 0
            for char in special_chars:
                if char == ' ':
                    continue
                row_indices = [i for i, value in enumerate(values) if char in value]
                count_per_row = [value.count(char) for value in values if char in value]
                total_char_count += sum(count_per_row)
                char_stats[char] = {
                    'total_occurrences': sum(count_per_row),
                    'rows_containing': len(row_indices),
                    'min_row_index': min(row_indices) if row_indices else None,
                    'max_row_index': max(row_indices) if row_indices else None,
                    'avg_occurrences_per_row': round(sum(count_per_row) / len(row_indices), 2) if row_indices else 0
                }
            blank_cells_count = df[column].isnull().sum()
            special_chars_analysis[column] = {
                'special_characters': char_stats,
                'blank_cells_count': blank_cells_count,
                'total_special_char_occurrences': total_char_count
            }
        return special_chars_analysis

    def remove_spaces(self, df, columns):
        for column in columns:
            if column in df.columns:
                df[column] = df[column].astype(str).str.replace(" ", "", regex=True)
        return df

    def remove_special_characters(self, df, options):
        for column, char_map in options.items():
            if column in df.columns:
                for char, remove in char_map.items():
                    if remove:
                        df[column] = df[column].astype(str).str.replace(re.escape(char), "", regex=True)
        return df


class FileUploadView(APIView):
    def post(self, request):
        file_type = request.data.get('file_type')
        file_name = request.data.get('file_name')
        project_id = request.data.get('project_id')
        user_id = request.data.get('user_id')  # For direct project access
        share_id = request.data.get('share_id')  # For shared project access
        permission_level = request.data.get('permission_level')
        is_shared = request.data.get('is_shared', False)  # Boolean flag for shared access

        # Validate required fields
        if not file_name or not project_id or not file_type:
            return Response({'error': 'Missing required fields'}, status=400)

        # Handle shared project access
        if is_shared:
            if not share_id:
                return Response({'error': 'share_id is required for shared projects'}, status=400)
            
            # Get the share
            try:
                share = ProjectShare.objects.get(id=share_id, is_active=True)
            except ProjectShare.DoesNotExist:
                return Response({'error': 'Share not found or revoked'}, status=404)
            
            # Validate project matches
            if share.project_id != int(project_id):
                return Response({'error': 'Project ID mismatch'}, status=400)
            
            user_id = share.shared_with.id
            permission_level = share.permission_level
            
            # Additional validation for file-specific shares
            if share.share_type == 'file':
                # For file-specific shares, validate the file matches
                if share.file_type != file_type or share.file_name != file_name:
                    return Response({
                        'error': f'Access denied. This share only allows access to {share.file_type}/{share.file_name}'
                    }, status=403)
            
            # For shared access, use the file_type from the request (not from share)
            # This allows accessing any file within the shared scope
        else:
            # For direct project access, user_id is required
            if not user_id:
                return Response({'error': 'user_id is required for direct project access'}, status=400)
            
            # Check if user has access to this project/file
            has_access, share_object, permission_level = check_project_access(
                user_id, project_id, file_type, file_name
            )
            
            if not has_access:
                return Response({
                    'error': 'Access denied. You don\'t have permission to access this project/file.'
                }, status=403)

        # Ensure project exists
        try:
            project = Projects.objects.get(id=project_id)
        except Projects.DoesNotExist:
            return Response({'error': 'Project not found'}, status=404)

        # Define folder structure
        project_folder = f"user_{project.user.id}/project_{project.id}"
        
        # Extract just the filenames from the database paths (handle both full paths and just filenames)
        def extract_filename(file_path):
            """Extract filename from either full path or just filename"""
            if '\\' in file_path or '/' in file_path:
                return os.path.basename(file_path)
            return file_path
        
        last_name_kpi = [extract_filename(file) for file in project.kpi_file] if project.kpi_file else []
        last_name_media = [extract_filename(file) for file in project.media_file] if project.media_file else []
        last_name_concatenated = project.concatenated_file if hasattr(project, 'concatenated_file') and isinstance(project.concatenated_file, list) else []

        # Validate that the file exists in the project
        file_exists = False
        if file_name in last_name_kpi and file_type == 'kpi':
            file_exists = True
        elif file_name in last_name_media and file_type == 'media':
            file_exists = True
        elif file_type == 'concatenated':
            # For concatenated files, check if any folder contains this CSV file
            concatenated_folder = os.path.join(settings.MEDIA_ROOT, project_folder, "concatenated")
            if os.path.exists(concatenated_folder):
                for folder_name in os.listdir(concatenated_folder):
                    folder_path = os.path.join(concatenated_folder, folder_name)
                    if os.path.isdir(folder_path):
                        # Check if this folder contains the CSV file
                        for csv_file in os.listdir(folder_path):
                            if csv_file == file_name or csv_file.endswith('.csv'):
                                file_exists = True
                                # Update file_name to the folder name for further processing
                                file_name = folder_name
                                break
                        if file_exists:
                            break
        elif file_name in last_name_concatenated and file_type == 'concatenated':
            file_exists = True
        
        if not file_exists:
            return Response({
                'error': f'File "{file_name}" of type "{file_type}" not found in project. Available files: KPI={last_name_kpi}, Media={last_name_media}, Concatenated={last_name_concatenated}'
            }, status=404)

        # Handle file name based on file type
        if file_type == 'concatenated':
            # For concatenated files, use the folder name as the file name
            if '/' in file_name:
                # If file_name includes the CSV name, extract just the folder name
                file_name = file_name.split('/')[0]
            else:
                file_name = os.path.basename(file_name)
        else:
            # For kpi and media files, use basename as before
            file_name = os.path.basename(file_name)

        # Build path based on file type
        # If file_type is None (project-level share), determine it from file_name
        if file_type is None:
            if file_name in last_name_kpi:
                file_type = 'kpi'
            elif file_name in last_name_media:
                file_type = 'media'
            elif file_name in last_name_concatenated:
                file_type = 'concatenated'
            else:
                return Response({'error': 'File not associated with the project'}, status=400)
        
        if file_name in last_name_kpi and file_type == 'kpi':
            base_folder = os.path.join(settings.MEDIA_ROOT, project_folder, "kpi", file_name)
        elif file_name in last_name_media and file_type == 'media':
            base_folder = os.path.join(settings.MEDIA_ROOT, project_folder, "media", file_name)
        elif file_type == 'concatenated':
            # For concatenated files, use the folder name directly (which was updated during validation)
            base_folder = os.path.join(settings.MEDIA_ROOT, project_folder, "concatenated", file_name)
        elif file_name in last_name_concatenated and file_type == 'concatenated':
            base_folder = os.path.join(settings.MEDIA_ROOT, project_folder, "concatenated", file_name)
        else:
            return Response({'error': 'File not associated with the project'}, status=400)

        base_folder = os.path.normpath(base_folder)

        if not os.path.exists(base_folder):
            return Response({'error': 'Folder not found'}, status=404)

        # Collect all CSV files inside the folder
        csv_files = []
        for root, dirs, files in os.walk(base_folder):
            for file in files:
                if file.endswith('.csv'):
                    csv_files.append(os.path.join(root, file))

        if not csv_files:
            return Response({'error': 'No CSV files found inside the folder'}, status=404)

        sheets_data = {}

        try:
            for csv_file in csv_files:
                # Relativize the path for Django's default_storage
                rel_csv_file = os.path.relpath(csv_file, settings.MEDIA_ROOT)

                # Determine file size
                file_size_mb = os.path.getsize(csv_file) / (1024 * 1024)

                # Open CSV file using default_storage in binary mode, wrap in TextIOWrapper for utf-8
                try:
                    with default_storage.open(rel_csv_file, 'rb') as f:
                        csvfile = io.TextIOWrapper(f, encoding='utf-8')
                        reader = csv.reader(csvfile)
                        rows = list(reader)
                except UnicodeDecodeError:
                    with default_storage.open(rel_csv_file, 'rb') as f:
                        csvfile = io.TextIOWrapper(f, encoding='latin1')
                        reader = csv.reader(csvfile)
                        rows = list(reader)

                if not rows:
                    continue  # Skip empty CSV files

                columns = rows[0]  # Header
                if file_size_mb < 50:
                    data = rows[1:]  # All data rows
                else:
                    data = rows[1:1001]  # Top 1000 rows

                column_types = {}

                # Common error values to ignore
                error_values = {"#VALUE!", "#N/A", "#DIV/0!", "#REF!", "#NAME?", "#NULL!", "#NUM!"}

                for col_index, col_name in enumerate(columns):
                    column_types[col_name] = "Unknown"
                    for row in data:
                        try:
                            value = row[col_index]
                            if value and value not in error_values:
                                if value.replace('.', '', 1).isdigit():
                                    if '.' in value:
                                        column_types[col_name] = "float"
                                    else:
                                        column_types[col_name] = "int"
                                else:
                                    column_types[col_name] = "str"
                                break
                        except (IndexError, ValueError, TypeError):
                            continue

                # Use CSV file name as "sheet name"
                sheet_name = os.path.basename(csv_file)
                sheets_data[sheet_name] = {
                    'columns': columns,
                    'data': data,
                    'column_types': column_types,
                    'hidden': False  # CSV can't have hidden sheets
                }
            user = None
            if hasattr(request, 'user') and request.user.is_authenticated:
                user = request.user
            elif project and hasattr(project, 'user'):
                user = project.user
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "file_data", details=f"File data retrieved successfully", ip_address=ip)

            return Response({
                'message': 'CSV file(s) data retrieved successfully',
                'sheets_data': sheets_data,
                'permission_level': permission_level
            }, status=200)

        except Exception as e:
            return Response({'error': f'Error reading the CSV files: {str(e)}'}, status=500)
        


class SignupView(APIView):
    def post(self, request):
        username = request.data.get('username')
        password = request.data.get('password')
        email = request.data.get('email')

        if not username or not password or not email:
            return Response({'error': 'Missing required fields'}, status=400)

        # Check if email already exists
        if User.objects.filter(email=email).exists():
            return Response({'error': 'E-Mail ID already exist'}, status=400)

        # Check if username already exists
        if User.objects.filter(username=username).exists():
            return Response({'error': 'Username already exist'}, status=400)

        user = User.objects.create(username=username, email=email, password=password)
        refresh = RefreshToken.for_user(user)
        ip = request.META.get('REMOTE_ADDR')
        log_user_action(user, "signup", details=f"User signed up successfully", ip_address=ip)

        return Response({
            'user': {
                'id': user.id,
                'username': user.username,
                'email': user.email
            },
            'refresh': str(refresh),
            'access': str(refresh.access_token),
        }, status=201)


class SigninView(APIView):
    def post(self, request):
        username = request.data.get('username')
        password = request.data.get('password')
        
        # Authenticate user with username and password
        user = CustomAuthBackend.authenticate(username=username, password=password)
        print(f"DEBUG: Authentication result - User: {user}, Username: {getattr(user, 'username', 'N/A') if user else 'None'}")
        
        if user:
            # Bypass OTP for username 'pc1'
            print(f"DEBUG: Username received: '{username}', Type: {type(username)}")
            print(f"DEBUG: User object username: '{user.username}', Type: {type(user.username)}")
            if username and username.strip() == "pc1":
                refresh = RefreshToken.for_user(user)
                ip = request.META.get('REMOTE_ADDR')
                log_user_action(user, "login", details="User logged in (OTP bypassed for pc1)", ip_address=ip)
                return Response({
                    'user': {
                        'id': user.id,
                        'username': user.username,
                        'email': user.email
                    },
                    'refresh': str(refresh),
                    'access': str(refresh.access_token),
                    'message': 'Login successful (OTP bypassed)'
                })
            # Generate OTP and send to user's email (for all other users)
            from .models import OTPToken
            from django.core.mail import send_mail
            from django.conf import settings
            
            try:
                # Create OTP token
                otp_token = OTPToken.create_for_user(user)
                
                # Send OTP email
                subject = 'Login Verification Code'
                message = f'''Hello {user.username},
                Your login verification code is: {otp_token.otp}
                This code will expire in 10 minutes.
                If you didn't request this code, please ignore this email.
                
                Best regards,
                Sanitify, Skewb AI
                '''
                
                send_mail(
                    subject,
                    message,
                    settings.DEFAULT_FROM_EMAIL,
                    [user.email],
                    fail_silently=False,
                )
                
                # Log the OTP request
                ip = request.META.get('REMOTE_ADDR')
                log_user_action(user, "otp_requested", details="OTP requested for login", ip_address=ip)
                
                return Response({
                    'message': 'OTP sent to your email',
                    'user_id': user.id,
                    'email': user.email,
                    'requires_otp': True
                })
                
            except Exception as e:
                # Log the error
                ip = request.META.get('REMOTE_ADDR')
                log_user_action(user, "otp_error", details=f"Failed to send OTP: {str(e)}", ip_address=ip)
                
                return Response({
                    'error': 'Failed to send OTP. Please try again.'
                }, status=500)
        else:
            return Response({'error': 'Invalid credentials'}, status=401)
    

class VerifyOTPView(APIView):
    def post(self, request):
        user_id = request.data.get('user_id')
        otp = request.data.get('otp')
        
        if not user_id or not otp:
            return Response({'error': 'User ID and OTP are required'}, status=400)
        
        try:
            user = User.objects.get(id=user_id)
        except User.DoesNotExist:
            return Response({'error': 'User not found'}, status=404)
        
        # Get the most recent valid OTP for this user
        from .models import OTPToken
        try:
            otp_token = OTPToken.objects.filter(
                user=user,
                otp=otp,
                is_used=False
            ).latest('created_at')
        except OTPToken.DoesNotExist:
            return Response({'error': 'Invalid or expired OTP'}, status=400)
        
        # Check if OTP is valid and not expired
        if not otp_token.is_valid():
            return Response({'error': 'OTP has expired'}, status=400)
        
        # Mark OTP as used
        otp_token.is_used = True
        otp_token.save()
        
        # Generate JWT tokens
        refresh = RefreshToken.for_user(user)
        
        # Log the successful login
        ip = request.META.get('REMOTE_ADDR')
        log_user_action(user, "login", details="User logged in with OTP verification", ip_address=ip)
        
        return Response({
            'user': {
                'id': user.id,
                'username': user.username,
                'email': user.email
            },
            'refresh': str(refresh),
            'access': str(refresh.access_token),
            'message': 'Login successful'
        })
class UploadProject(APIView):
    def post(self, request):
        user_id = request.data.get('user_id')
        name = request.data.get('project_name')
        kpi_files = request.FILES.getlist('kpi_file')
        media_files = request.FILES.getlist('media_file')

        if not user_id or not name:
            return Response({'error': 'Missing required fields'}, status=400)

        if not kpi_files and not media_files:
            return Response({'error': 'No files provided for KPI or Media'}, status=400)

        try:
            user = User.objects.get(id=user_id)
        except User.DoesNotExist:
            return Response({'error': 'User not found'}, status=404)

        # Log the project upload event
        ip = request.META.get('REMOTE_ADDR')
        log_user_action(user, "upload_project", details=f"Project name: {name}", ip_address=ip)

        # Create new project
        project = Projects.objects.create(user=user, name=name, kpi_file=[], media_file=[])
        project_folder = os.path.join(settings.MEDIA_ROOT, f"user_{user.id}/project_{project.id}")
        os.makedirs(project_folder, exist_ok=True)

        if not os.path.exists(os.path.join(project_folder, ".git")):
            subprocess.run(["git", "init"], cwd=project_folder)
            subprocess.run(["git", "config", "user.name", user.username], cwd=project_folder)
            subprocess.run(["git", "config", "user.email", user.email], cwd=project_folder)

        spark = None  # Spark session, only created if needed
        spark_needed = False
        all_files = [(file, 'kpi') for file in kpi_files] + [(file, 'media') for file in media_files]
        
        # Optimized file size checking - avoid writing files twice
        for file, _ in all_files:
            # Get file size from Django's file object without writing to disk
            file.seek(0, 2)  # Seek to end
            file_size_bytes = file.tell()
            file.seek(0)  # Reset to beginning
            file_size_mb = file_size_bytes / (1024 * 1024)
            
            if file_size_mb >= 50:
                spark_needed = True
                print(f"Large file detected: {file.name} ({file_size_mb:.2f} MB) - will use Spark")
            else:
                print(f"Small file detected: {file.name} ({file_size_mb:.2f} MB) - will use pandas")
                
        if spark_needed:
            spark = get_spark_session()

        def process_file(file, file_type):
            print(f"Processing {file.name} ({file_type})...")
            if file_type == 'kpi':
                last_obj = Projects.objects.order_by('-kpi_id').first()
                id_field = 'kpi_id'
                updated_list = project.kpi_file
                base_subdir = 'kpi'
            else:
                last_obj = Projects.objects.order_by('-media_id').first()
                id_field = 'media_id'
                updated_list = project.media_file
                base_subdir = 'media'

            file_id = (getattr(last_obj, id_field, 0) or 0) + 1
            file_basename = os.path.splitext(os.path.basename(file.name))[0]
            file_folder = os.path.join(project_folder, base_subdir, file_basename)
            os.makedirs(file_folder, exist_ok=True)

            temp_path = os.path.join(file_folder, file.name)
            with open(temp_path, 'wb') as f:
                for chunk in file.chunks():
                    f.write(chunk)

            file_size_mb = os.path.getsize(temp_path) / (1024 * 1024)
            file_extension = os.path.splitext(file.name)[1].lower()

            # Handle CSV files directly without conversion
            if file_extension == '.csv':
                # For CSV files, just copy them directly without conversion
                csv_path = os.path.join(file_folder, file.name)
                if temp_path != csv_path:
                    import shutil
                    shutil.copy2(temp_path, csv_path)
                    os.remove(temp_path)
                
                commit_msg = f"updated - {project.user.id}/{project.id}/{base_subdir}/{file_basename}/{file.name}"
                subprocess.run(["git", "add", csv_path], cwd=project_folder)
                subprocess.run(["git", "commit", "-m", commit_msg], cwd=project_folder)
                
            # Handle Excel files with conversion to CSV
            elif file_extension in ['.xlsx', '.xls']:
                if file_size_mb < 50:
                    with pd.ExcelFile(temp_path) as xls:
                        for sheet_name in xls.sheet_names:
                            df = xls.parse(sheet_name, dtype=str)
                            sheet_path = os.path.join(file_folder, f"{sheet_name}.csv")
                            df.to_csv(sheet_path, index=False)
                            commit_msg = f"updated - {project.user.id}/{project.id}/{base_subdir}/{file_basename}/{sheet_name}"
                            subprocess.run(["git", "add", sheet_path], cwd=project_folder)
                            subprocess.run(["git", "commit", "-m", commit_msg], cwd=project_folder)
                else:
                    # Use the new robust Spark session management for large files
                    from .spark_utils import spark_session_context, validate_spark_session
                    
                    try:
                        xls = pd.ExcelFile(temp_path, engine='openpyxl')
                        sheet_names = xls.sheet_names
                    except Exception as e:
                        print(f"❌ Failed to extract sheet names: {e}")
                        raise
                    
                    def convert_sheet(sheet_name):
                        output_path = os.path.join(file_folder, f"{sheet_name}.csv")
                        commit_msg = f"updated - {project.user.id}/{project.id}/{base_subdir}/{file_basename}/{sheet_name}"
                        
                        with spark_session_context() as spark:
                            # Validate session before use
                            if not validate_spark_session(spark):
                                raise Exception("Invalid Spark session. Please try again.")
                            
                            try:
                                df = spark.read \
                                    .format("com.crealytics.spark.excel") \
                                    .option("dataAddress", f"'{sheet_name}'!A1") \
                                    .option("header", "true") \
                                    .option("inferSchema", "false") \
                                    .option("maxRowsInMemory", 50000) \
                                    .option("maxColumns", 20000) \
                                    .option("treatEmptyValuesAsNulls", "true") \
                                    .option("workbookPassword", None) \
                                    .load(temp_path)
                                
                                # Cache for better performance
                                df = df.cache()
                                df.toPandas().to_csv(output_path, index=False)
                                df.unpersist()  # Free memory
                                
                                subprocess.run(["git", "add", output_path], cwd=project_folder)
                                subprocess.run(["git", "commit", "-m", commit_msg], cwd=project_folder)
                                
                            except Exception as e:
                                if "RecordFormatException" in str(e):
                                    import warnings
                                    warnings.filterwarnings("ignore", category=UserWarning)
                                    try:
                                        xls = pd.ExcelFile(temp_path, engine='openpyxl')
                                        df = xls.parse(sheet_name, dtype=str)
                                        df.to_csv(output_path, index=False)
                                        subprocess.run(["git", "add", output_path], cwd=project_folder)
                                        subprocess.run(["git", "commit", "-m", commit_msg], cwd=project_folder)
                                        print(f"✅ Fallback to pandas successful for sheet: {sheet_name}")
                                    except Exception as pe:
                                        print(f"❌ Pandas fallback failed for {sheet_name}: {pe}")
                                        raise
                                else:
                                    print(f"❌ Unexpected Spark error for {sheet_name}: {e}")
                                    raise
                    
                    # Use optimized thread pool for parallel processing
                    with ThreadPoolExecutor(max_workers=min(8, len(sheet_names))) as executor:
                        futures = [executor.submit(convert_sheet, sheet) for sheet in sheet_names]
                        for f in futures:
                            f.result()
            else:
                # Unsupported file type
                os.remove(temp_path)
                raise Exception(f"Unsupported file type: {file_extension}. Only .csv, .xlsx, and .xls files are supported.")
            
            updated_list.append(file_basename)
            setattr(project, id_field, file_id)

        updated_kpi_files = []
        for file in kpi_files:
            process_file(file, 'kpi')
            updated_kpi_files.append(os.path.splitext(file.name)[0])

        updated_media_files = []
        for file in media_files:
            process_file(file, 'media')
            updated_media_files.append(os.path.splitext(file.name)[0])

        project.save()

        if spark:
            spark.stop()

        return Response({
            'message': 'Project uploaded successfully',
            'project_id': project.id,
            'name': project.name,
            'uploaded_kpi_files': updated_kpi_files,
            'uploaded_media_files': updated_media_files,
        }, status=201)



class UserProjectsView(APIView):
    def post(self, request):
        user_id = request.data.get('user_id')

        if not user_id:
            return Response({'error': 'User ID is required'}, status=400)

        user = get_object_or_404(User, id=user_id)
        projects = Projects.objects.filter(user=user)

        project_data = []
        for project in projects:
            # Ensure kpi_file and media_file are lists
            kpi_files = project.kpi_file if isinstance(project.kpi_file, list) else []
            media_files = project.media_file if isinstance(project.media_file, list) else []
            
            # Get concatenated files
            concatenated_files = project.concatenated_file if hasattr(project, 'concatenated_file') and isinstance(project.concatenated_file, list) else []

            # Process concatenated files to get actual CSV filenames
            concatenated_file_data = []
            for timestamp_folder in concatenated_files:
                concatenated_folder_path = os.path.join(settings.MEDIA_ROOT, f"user_{user.id}/project_{project.id}/concatenated/{timestamp_folder}")
                if os.path.exists(concatenated_folder_path):
                    # Look for CSV files in the timestamp folder
                    for file in os.listdir(concatenated_folder_path):
                        if file.endswith('.csv'):
                            concatenated_file_data.append({
                                'id': len(concatenated_file_data) + 1,
                                'path': f"user_{user.id}/project_{project.id}/concatenated/{timestamp_folder}/{file}",
                                'name': f"{timestamp_folder}/{file}"
                            })
            
            project_data.append({
                'id': project.id,
                'name': project.name,
                'files': {
                    'kpis': [{'id': idx + 1, 'path': f"user_{user.id}/project_{project.id}/kpi/{file_name}", 'name': file_name} 
                             for idx, file_name in enumerate(kpi_files)],
                    'media': [{'id': idx + 1, 'path': f"user_{user.id}/project_{project.id}/media/{file_name}", 'name': file_name} 
                              for idx, file_name in enumerate(media_files)],
                    'concatenated': concatenated_file_data
                }
            })
        # Robust user logging
        logging_user = get_logging_user(request, user)
        if logging_user:
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(logging_user, "view_projects", details="User viewed projects", ip_address=ip)
        return Response({'projects': project_data}, status=200)


class DeleteProject(APIView):
    def delete(self, request):
        project_id = request.data.get("project_id")
        

        if not project_id :
            return Response({"error": "Project ID and User ID are required"}, status=400)

        try:
            project = Projects.objects.get(id=project_id)
        except Projects.DoesNotExist:
            return Response({"error": "Project not found"}, status=404)

        # Define project folder path
        project_path = os.path.join(settings.MEDIA_ROOT, f"user_{project.user.id}/project_{project.id}")
        project_path = os.path.abspath(os.path.normpath(project_path))

        print("Deleting project at:", project_path)

        try:
            # Check if it's a Git repository
            git_folder = os.path.join(project_path, ".git")
            if os.path.exists(git_folder):
                print("Removing Git repository...")

                # Run git commands to clean up the repository
                subprocess.run(["git", "-C", project_path, "gc", "--prune=now"], check=False)
                subprocess.run(["git", "-C", project_path, "rm", "-rf", "--cached", "."], check=False)

                # Force remove .git folder
                shutil.rmtree(git_folder, ignore_errors=True)

            # Now delete the project folder
            if os.path.exists(project_path):
                def remove_readonly(func, path, _):
                    """Force remove read-only files like .git/index.lock"""
                    os.chmod(path, stat.S_IWRITE)
                    func(path)

                shutil.rmtree(project_path, onerror=remove_readonly)

            # Delete the project from the database
            project.delete()
            user = None
            if hasattr(request, 'user') and request.user.is_authenticated:
                user = request.user
            elif project and hasattr(project, 'user'):
                user = project.user
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "delete_project", details=f"Project deleted successfully", ip_address=ip)
            return Response({"message": "Project, files, and Git repository deleted successfully"}, status=200)

        except Exception as e:
            print("Error while deleting project:", traceback.format_exc())  # Print full error details
            return Response({"error": str(e)}, status=500)
        


class UpdateProject(APIView):
    def post(self, request):
        project_id = request.data.get('project_id')
        user_id = request.data.get('user_id')
        kpi_files = request.FILES.getlist('kpi_file')
        media_files = request.FILES.getlist('media_file')

        if not project_id or not user_id:
            return Response({'error': 'Missing required fields'}, status=400)

        try:
            project = Projects.objects.get(id=project_id, user_id=user_id)
        except Projects.DoesNotExist:
            return Response({'error': 'Project not found or access denied'}, status=404)

        project_folder = os.path.join(settings.MEDIA_ROOT, f"user_{project.user.id}/project_{project.id}")
        os.makedirs(project_folder, exist_ok=True)

        try:
            user = User.objects.get(id=project.user.id)
        except User.DoesNotExist:
            return Response({'error': 'User not found'}, status=404)

        if not os.path.exists(os.path.join(project_folder, ".git")):
            subprocess.run(["git", "init"], cwd=project_folder)
            subprocess.run(["git", "config", "user.name", user.username], cwd=project_folder)
            subprocess.run(["git", "config", "user.email", user.email], cwd=project_folder)

        spark = None  # Spark session, only created if needed
        spark_needed = False
        all_files = [(file, 'kpi') for file in kpi_files] + [(file, 'media') for file in media_files]
        
        # Optimized file size checking - avoid writing files twice
        for file, _ in all_files:
            # Get file size from Django's file object without writing to disk
            file.seek(0, 2)  # Seek to end
            file_size_bytes = file.tell()
            file.seek(0)  # Reset to beginning
            file_size_mb = file_size_bytes / (1024 * 1024)
            
            if file_size_mb >= 50:
                spark_needed = True
                print(f"Large file detected: {file.name} ({file_size_mb:.2f} MB) - will use Spark")
            else:
                print(f"Small file detected: {file.name} ({file_size_mb:.2f} MB) - will use pandas")
                
        if spark_needed:
            spark = get_spark_session()

        def process_file(file, file_type):
            print(f"Processing {file.name} ({file_type})...")
            if file_type == 'kpi':
                last_obj = Projects.objects.order_by('-kpi_id').first()
                id_field = 'kpi_id'
                updated_list = project.kpi_file
                base_subdir = 'kpi'
            else:
                last_obj = Projects.objects.order_by('-media_id').first()
                id_field = 'media_id'
                updated_list = project.media_file
                base_subdir = 'media'

            file_id = (getattr(last_obj, id_field, 0) or 0) + 1
            file_basename = os.path.splitext(os.path.basename(file.name))[0]
            file_folder = os.path.join(project_folder, base_subdir, file_basename)
            os.makedirs(file_folder, exist_ok=True)

            temp_path = os.path.join(file_folder, file.name)
            with open(temp_path, 'wb') as f:
                for chunk in file.chunks():
                    f.write(chunk)

            file_size_mb = os.path.getsize(temp_path) / (1024 * 1024)
            file_extension = os.path.splitext(file.name)[1].lower()

            # Handle CSV files directly without conversion
            if file_extension == '.csv':
                # For CSV files, just copy them directly without conversion
                csv_path = os.path.join(file_folder, file.name)
                if temp_path != csv_path:
                    import shutil
                    shutil.copy2(temp_path, csv_path)
                    os.remove(temp_path)
                
                commit_msg = f"updated - {project.user.id}/{project.id}/{base_subdir}/{file_basename}/{file.name}"
                subprocess.run(["git", "add", csv_path], cwd=project_folder)
                subprocess.run(["git", "commit", "-m", commit_msg], cwd=project_folder)
                
            # Handle Excel files with conversion to CSV
            elif file_extension in ['.xlsx', '.xls']:
                if file_size_mb < 50:
                    with pd.ExcelFile(temp_path) as xls:
                        for sheet_name in xls.sheet_names:
                            df = xls.parse(sheet_name, dtype=str)
                            sheet_path = os.path.join(file_folder, f"{sheet_name}.csv")
                            df.to_csv(sheet_path, index=False)
                            commit_msg = f"updated - {project.user.id}/{project.id}/{base_subdir}/{file_basename}/{sheet_name}"
                            subprocess.run(["git", "add", sheet_path], cwd=project_folder)
                            subprocess.run(["git", "commit", "-m", commit_msg], cwd=project_folder)
                else:
                    # Use the new robust Spark session management for large files
                    from .spark_utils import spark_session_context, validate_spark_session
                    
                    try:
                        xls = pd.ExcelFile(temp_path, engine='openpyxl')
                        sheet_names = xls.sheet_names
                    except Exception as e:
                        print(f"❌ Failed to extract sheet names: {e}")
                        raise
                    
                    def convert_sheet(sheet_name):
                        output_path = os.path.join(file_folder, f"{sheet_name}.csv")
                        commit_msg = f"updated - {project.user.id}/{project.id}/{base_subdir}/{file_basename}/{sheet_name}"
                        
                        with spark_session_context() as spark:
                            # Validate session before use
                            if not validate_spark_session(spark):
                                raise Exception("Invalid Spark session. Please try again.")
                            
                            try:
                                df = spark.read \
                                    .format("com.crealytics.spark.excel") \
                                    .option("dataAddress", f"'{sheet_name}'!A1") \
                                    .option("header", "true") \
                                    .option("inferSchema", "false") \
                                    .option("maxRowsInMemory", 50000) \
                                    .option("maxColumns", 20000) \
                                    .option("treatEmptyValuesAsNulls", "true") \
                                    .option("workbookPassword", None) \
                                    .load(temp_path)
                                
                                # Cache for better performance
                                df = df.cache()
                                df.toPandas().to_csv(output_path, index=False)
                                df.unpersist()  # Free memory
                                
                                subprocess.run(["git", "add", output_path], cwd=project_folder)
                                subprocess.run(["git", "commit", "-m", commit_msg], cwd=project_folder)
                                
                            except Exception as e:
                                if "RecordFormatException" in str(e):
                                    import warnings
                                    warnings.filterwarnings("ignore", category=UserWarning)
                                    try:
                                        xls = pd.ExcelFile(temp_path, engine='openpyxl')
                                        df = xls.parse(sheet_name, dtype=str)
                                        df.to_csv(output_path, index=False)
                                        subprocess.run(["git", "add", output_path], cwd=project_folder)
                                        subprocess.run(["git", "commit", "-m", commit_msg], cwd=project_folder)
                                        print(f"✅ Fallback to pandas successful for sheet: {sheet_name}")
                                    except Exception as pe:
                                        print(f"❌ Pandas fallback failed for {sheet_name}: {pe}")
                                        raise
                                else:
                                    print(f"❌ Unexpected Spark error for {sheet_name}: {e}")
                                    raise
                    
                    # Use optimized thread pool for parallel processing
                    with ThreadPoolExecutor(max_workers=min(8, len(sheet_names))) as executor:
                        futures = [executor.submit(convert_sheet, sheet) for sheet in sheet_names]
                        for f in futures:
                            f.result()
            else:
                # Unsupported file type
                os.remove(temp_path)
                raise Exception(f"Unsupported file type: {file_extension}. Only .csv, .xlsx, and .xls files are supported.")
            
            updated_list.append(file_basename)
            setattr(project, id_field, file_id)

        updated_kpi_files = []
        for file in kpi_files:
            process_file(file, 'kpi')
            updated_kpi_files.append(os.path.splitext(file.name)[0])

        updated_media_files = []
        for file in media_files:
            process_file(file, 'media')
            updated_media_files.append(os.path.splitext(file.name)[0])

        project.save()

        if spark:
            spark.stop()

        logging_user = get_logging_user(request, user)
        if logging_user:
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(logging_user, "update_project", details=f"Project updated successfully", ip_address=ip)

        return Response({
            'message': 'Project files updated successfully',
            'project_id': project.id,
            'updated_kpi_files': updated_kpi_files,
            'updated_media_files': updated_media_files,
        }, status=200)


class GetSpecificSheetCommits(APIView):
    def post(self, request):
        project_id = request.data.get("project_id")
        user_id = request.data.get("user_id")
        file_name = request.data.get("file_name")
        sheet_name = request.data.get("sheet_name")
        file_type = request.data.get("file_type")
        send_only_commits = request.data.get("send_only_commits")

        original_file_name = file_name
        file_name = file_name.replace("\\", "/").split("/")[-1]
        git_relative_path = f"{file_type}/{original_file_name}/{sheet_name}"

        if not project_id or not user_id or not file_name or not sheet_name:
            return Response({"error": "Missing required parameters"}, status=400)

        # Check access (existing code remains the same)
        has_access = False
        share_object = None
        permission_level = None
        has_access, share_object, permission_level = check_project_access(
            user_id, project_id, file_type, original_file_name, sheet_name
        )
        if not has_access:
            has_access, share_object, permission_level = check_project_access(
                user_id, project_id, file_type, file_name, sheet_name
            )
        if not has_access:
            has_access, share_object, permission_level = check_project_access(
                user_id, project_id
            )
        if not has_access:
            return Response({"error": "Access denied to this project or file"}, status=403)

        try:
            project = Projects.objects.get(id=project_id)
            project_owner_id = project.user.id
        except Projects.DoesNotExist:
            return Response({"error": "Project not found"}, status=404)

        project_folder = os.path.join(settings.MEDIA_ROOT, f"user_{project_owner_id}/project_{project_id}")
        git_relative_path = f"{file_type}/{file_name}/{sheet_name}"

        try:
            

            # Get commit history
            result = subprocess.run(
                ["git", "log", "--pretty=%H %P", "--reverse", "--", git_relative_path],
                cwd=project_folder,
                capture_output=True,
                text=True
            )

            if result.returncode != 0:
                return Response({"error": "Failed to fetch commits"}, status=500)

            raw_commits = result.stdout.strip().split("\n")
            commit_map = {}
            commit_order = []

            # First pass: build commit metadata
            for raw_commit in raw_commits:
                parts = raw_commit.strip().split()
                if not parts:
                    continue

                commit_hash = parts[0]
                parents = parts[1:] if len(parts) > 1 else []

                msg_result = subprocess.run(
                    ["git", "show", "-s", "--format=%s", commit_hash],
                    cwd=project_folder,
                    capture_output=True,
                    text=True
                )
                message = msg_result.stdout.strip() if msg_result.returncode == 0 else ""
                message_lower = message.lower()

                if "undo" in message_lower:
                    operation_type = "undo"
                elif "redo" in message_lower:
                    operation_type = "redo"
                elif "cleaned" in message_lower:
                    operation_type = "cleaning"
                elif "melted" in message_lower:
                    operation_type = "melting"
                elif "mapping" in message_lower:
                    operation_type = "mapping"
                elif "merging" in message_lower:
                    operation_type = "merging"
                elif "uploaded" in message_lower:
                    operation_type = "uploaded"
                elif "script" in message_lower:
                    operation_type = "script"
                elif "updated" in message_lower:
                    operation_type = "updated"
                elif "google" in message_lower:
                    operation_type = "google sheet"
                else:
                    operation_type = "other"

                if send_only_commits:
                    sheet_data = None
                else:
                    try:
                        sheet_data_result = subprocess.run(
                            ['git', 'show', f'{commit_hash}:{git_relative_path}'],
                            cwd=project_folder,
                            capture_output=True,
                            text=True,
                            timeout=10
                        )
                        if sheet_data_result.returncode == 0 and sheet_data_result.stdout.strip():
                            from io import StringIO
                            sheet_data_df = pd.read_csv(StringIO(sheet_data_result.stdout.strip()), dtype=str)
                            sheet_data_df = sheet_data_df.replace([np.nan, np.inf, -np.inf], "value error")
                            sheet_data = {
                                "columns": sheet_data_df.columns.tolist(),
                                "data": sheet_data_df.values.tolist()
                            }
                        else:
                            sheet_data = None
                    except Exception as e:
                        print(f"Error retrieving sheet data: {e}")
                        sheet_data = None

                commit_map[commit_hash] = {
                    "hash": commit_hash,
                    "message": message,
                    "operation_type": operation_type,
                    "parents": parents,
                    "children": [],
                    "sheet_data": sheet_data
                }
                commit_order.append(commit_hash)

            # Second pass: build initial children relationships
            for commit_hash in commit_order:
                for parent_hash in commit_map[commit_hash]["parents"]:
                    if parent_hash in commit_map:
                        commit_map[parent_hash]["children"].append(commit_hash)

            # Third pass: handle special operations (undo, script, etc.)
            for commit_hash in commit_order:
                commit = commit_map[commit_hash]
                operation_type = commit["operation_type"]

                if operation_type == "undo":
                    # For undo operations, we need to skip the immediate parent
                    # and connect to the grandparent instead
                    if len(commit["parents"]) == 1:
                        parent_hash = commit["parents"][0]
                        if parent_hash in commit_map:
                            parent_commit = commit_map[parent_hash]
                            # Find the grandparent (the parent of the operation being undone)
                            if len(parent_commit["parents"]) > 0:
                                grandparent_hash = parent_commit["parents"][0]
                                # Update the undo's parent to be the grandparent
                                commit["parents"] = [grandparent_hash]
                                # Update the grandparent's children to include this undo
                                if grandparent_hash in commit_map:
                                    commit_map[grandparent_hash]["children"].append(commit_hash)
                                # Remove this undo from the immediate parent's children
                                if commit_hash in parent_commit["children"]:
                                    parent_commit["children"].remove(commit_hash)

                elif operation_type == "script":
                    # Script operations should be children of the previous "updated" operation
                    # and parents of subsequent operations (cleaning, melting, etc.)
                    # We'll handle this by ensuring the script operation maintains its relationships
                    pass  # The basic parent-child relationships are already correct

            # Fourth pass: ensure script operations are properly connected to their children
            # This handles cases where cleaning/melting operations should be children of script operations
            script_commits = [c for c in commit_map.values() if c["operation_type"] == "script"]
            for script_commit in script_commits:
                script_hash = script_commit["hash"]
                # Find all commits that happened after this script commit
                script_index = commit_order.index(script_hash)
                subsequent_commits = commit_order[script_index + 1:]
                for subsequent_hash in subsequent_commits:
                    subsequent_commit = commit_map[subsequent_hash]
                    # If this is a cleaning/melting operation and it's not already connected to the script
                    if subsequent_commit["operation_type"] in ["cleaning", "melting", "mapping"]:
                        if script_hash not in subsequent_commit["parents"]:
                            # Add script as a parent
                            subsequent_commit["parents"].append(script_hash)
                            # Add this commit as a child of the script
                            script_commit["children"].append(subsequent_hash)
                            # Remove any duplicate parent relationships
                            subsequent_commit["parents"] = list(set(subsequent_commit["parents"]))

            # Log the action
            user = get_logging_user(request, User.objects.get(id=user_id))
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "get_commits",
                          details=f"Retrieved commits for project {project_id}, file: {file_name}, sheet: {sheet_name}",
                          ip_address=ip)

            return Response({
                "commits": [commit_map[hash] for hash in commit_order],
                "permission_level": permission_level
            }, status=200)

        except Exception as e:
            return Response({"error": str(e)}, status=500)

class GetSpecificSheetCommitsArray(APIView):
    def post(self, request):
        
        project_id = request.data.get("project_id")
        user_id = request.data.get("user_id")
        file_name = request.data.get("file_name")
        sheet_name = request.data.get("sheet_name")
        file_type = request.data.get("file_type")
        send_only_commits = request.data.get("send_only_commits")
        # Sanitize file_name to remove invalid characters and spaces
        file_name = file_name.replace("\\", "/").split("/")[-1]
        # Remove leading/trailing spaces and replace with underscores 
        file_name = file_name.strip().replace(" ", "_")

        if not project_id or not user_id or not file_name or not sheet_name:
            return Response({"error": "Missing required parameters"}, status=400)

        # Check project access (including shared projects)
        has_access, share_object, permission_level = check_project_access(
            user_id, project_id, file_type, file_name, sheet_name
        )
        
        if not has_access:
            return Response({"error": "Access denied to this project or file"}, status=403)

        # Get the project owner's user ID for the correct folder path
        try:
            project = Projects.objects.get(id=project_id)
            project_owner_id = project.user.id
        except Projects.DoesNotExist:
            return Response({"error": "Project not found"}, status=404)

        # Use project owner's folder for git operations
        project_folder = os.path.join(settings.MEDIA_ROOT, f"user_{project_owner_id}/project_{project_id}")
        git_relative_path = f"{file_type}/{file_name}/{sheet_name}"

        try:
            result = subprocess.run(
                ["git", "log", "--oneline", "--reverse", "--", git_relative_path],
                cwd=project_folder,
                capture_output=True,
                text=True
            )

            if result.returncode != 0:
                return Response({"error": "Failed to fetch commits"}, status=500)

            raw_commits = result.stdout.strip().split("\n")
            commit_map = {}
            parent_stack = []
            current_branch_parents = []
            last_valid_commit = None

            for i, raw_commit in enumerate(raw_commits):
                parts = raw_commit.split(" ", 1)
                if len(parts) < 2:
                    continue

                commit_hash, message = parts
                message_lower = message.lower()

                # Identify operation
                if "undo" in message_lower:
                    operation_type = "undo"
                elif "redo" in message_lower:
                    operation_type = "redo"
                elif "cleaned" in message_lower:
                    operation_type = "cleaning"
                elif "melted" in message_lower:
                    operation_type = "melting"
                elif "mapped" in message_lower:
                    operation_type = "mapping"
                elif "merged" in message_lower:
                    operation_type = "merging"
                elif "uploaded" in message_lower:
                    operation_type = "uploaded"
                else:
                    operation_type = "other"

                parents = []
                if operation_type == "uploaded":
                    parent_stack = [commit_hash]
                elif operation_type == "undo":
                    if parent_stack:
                        last_valid_commit = parent_stack.pop()
                else:
                    if operation_type == "redo" and last_valid_commit:
                        parents = [last_valid_commit]
                    elif parent_stack:
                        parents = [parent_stack[-1]]
                    else:
                        parents = []

                    parent_stack.append(commit_hash)

                # Retrieve sheet data
                if send_only_commits:
                    sheet_data = None
                else:
                    try:
                        sheet_data_result = subprocess.run(
                            ['git', 'show', f'{commit_hash}:{git_relative_path}'],
                            cwd=project_folder,
                            capture_output=True,
                            text=True,
                            timeout=10
                        )

                        if sheet_data_result.returncode == 0 and sheet_data_result.stdout.strip():
                            from io import StringIO
                            sheet_data_df = pd.read_csv(StringIO(sheet_data_result.stdout.strip()), dtype=str)
                            sheet_data_df = sheet_data_df.replace([np.nan, np.inf, -np.inf], "value error")
                            sheet_data = {
                                "columns": sheet_data_df.columns.tolist(),
                                "data": sheet_data_df.values.tolist()
                            }
                        else:
                            sheet_data = None
                    except Exception as e:
                        print(f"Error retrieving sheet data: {e}")
                        sheet_data = None

                commit_map[commit_hash] = {
                    "hash": commit_hash,
                    "message": message,
                    "operation_type": operation_type,
                    "parents": parents,
                    "children": [],
                    "sheet_data": sheet_data
                }

            # Set up children for each commit
            
            for commit in commit_map.values():
                for parent_hash in commit["parents"]:
                    if parent_hash in commit_map:
                        commit_map[parent_hash]["children"].append(commit["hash"])

            # Exclude undo and redo commits from the response
            filtered_commits = [commit for commit in commit_map.values() if commit["operation_type"] not in ["undo", "redo"]]

            # Log the event
            user = get_logging_user(request, User.objects.get(id=user_id))
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "get_specific_sheet_commits_array", details=f"Retrieved commits array for project {project_id}, file: {file_name}, sheet: {sheet_name}", ip_address=ip)

            return Response({
                "commits": filtered_commits,
                "permission_level": permission_level
            }, status=200)

        except Exception as e:
            return Response({"error": str(e)}, status=500)


class UndoRedoSheet(APIView):
    def post(self, request):
        project_id = request.data.get("project_id")
        user_id = request.data.get("user_id")
        file_name = request.data.get("file_name")
        sheet_name = request.data.get("sheet_name")
        file_type = request.data.get("file_type")
        action = request.data.get("action")  # 'undo' or 'redo'
        hash = request.data.get("hash")

        # Sanitize file_name to remove invalid characters and spaces
        file_name = file_name.split("\\")[1] if "\\" in file_name else file_name
        # Remove leading/trailing spaces and replace spaces with underscores
        # file_name = file_name.strip().replace(" ", "_")

        if not all([project_id, user_id, file_name, sheet_name, file_type, action]):
            return Response({"error": "Missing required parameters"}, status=400)

        # Check project access (including shared projects)
        has_access, share_object, permission_level = check_project_access(
            user_id, project_id, file_type, file_name, sheet_name
        )
        
        if not has_access:
            return Response({"error": "Access denied to this project or file"}, status=403)

        # Get the project owner's user ID for the correct folder path
        try:
            project = Projects.objects.get(id=project_id)
            project_owner_id = project.user.id
        except Projects.DoesNotExist:
            return Response({"error": "Project not found"}, status=404)

        # Use project owner's folder for git operations
        project_folder = os.path.join(settings.MEDIA_ROOT, f"user_{project_owner_id}/project_{project_id}")
        sheet_path = os.path.join(project_folder, file_type, file_name, f"{sheet_name}")

        sheet_path=os.path.normpath(sheet_path)
        print(f"Debug - sheet_path: {sheet_path}")
        if not os.path.exists(sheet_path):
            return Response({"error": "Sheet file not found"}, status=404)

        try:
            print("1")

            if isinstance(hash, str):
                hash = json.loads(hash)
            commit_hash = hash.get("hash")
            # Get commit history for the specific sheet
            result = subprocess.run(
                ["git", "log", "--oneline", "--", sheet_path],
                cwd=project_folder,
                capture_output=True,
                text=True
            )

            print("2")

            if result.returncode != 0:
                return Response({"error": "Failed to fetch commit history"}, status=500)

            commits = result.stdout.strip().split("\n")
            if not commits or len(commits) < 2:
                return Response({"error": "No previous state available for undo/redo"}, status=400)
            
            print("3")

            commit_hashes = [commit.split(" ")[0] for commit in commits]
            
            target_index = commit_hashes.index(commit_hash)
            target_commit = commit_hashes[target_index]
            
            print("4")

            if not target_commit:
                return Response({"error": "No commit found for undo/redo"}, status=400)

            # Restore the sheet file from the selected commit
            subprocess.run(["git", "checkout", target_commit, "--", sheet_path], cwd=project_folder)

            df_restored = pd.read_csv(sheet_path)

            df_restored = df_restored.replace([float("inf"), -float("inf")], None)
            df_restored = df_restored.fillna("NA")
            try:
                self.commit_to_git(
                    project_folder,
                    project.user,
                    project_id,
                    file_type,
                    file_name,
                    sheet_name,
                    action
                )
            except Exception as e:
                print(f"Failed to create commit: {str(e)}")

            # Log the event
            user = get_logging_user(request, User.objects.get(id=user_id))
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "undo_redo_sheet", details=f"{action.capitalize()} sheet successful for project {project_id}, file: {file_name}, sheet: {sheet_name}", ip_address=ip)

            return Response({
                "message": f"{action.capitalize()} successful",
                "sheet_data": {
                    "columns": df_restored.columns.tolist(),
                    "data": df_restored.values.tolist()
                }
            }, status=200)

        except Exception as e:
            return Response({"error": str(e)}, status=500)
        
    def commit_to_git(self, project_folder, user, project_id, file_type, file_name, sheet_name, action):
        # Dummy commit logic — update based on your implementation
        commit_msg = f"{action.capitalize()} action on {sheet_name} by {user.username}"
        subprocess.run(["git", "add", "."], cwd=project_folder)
        subprocess.run(["git", "commit", "-m", commit_msg], cwd=project_folder)


class PivotEDAAnalysis(APIView):
    def post(self, request):
        try:
            data = request.data.get('data')
            if not data:
                return Response({"error": "No data provided"}, status=400)
            df = pd.DataFrame(data)
            eda_results = {
                "basic_info": {
                    "total_rows": len(df),
                    "total_columns": len(df.columns),
                    "memory_usage": df.memory_usage(deep=True).sum(),
                    "duplicate_rows": df.duplicated().sum(),
                    "column_types": {
                        col: (
                            "numerical" if pd.api.types.is_numeric_dtype(df[col])
                            else "datetime" if pd.api.types.is_datetime64_any_dtype(df[col])
                            else "categorical"
                        )
                        for col in df.columns
                    }
                },
                "missing_values": self.analyze_missing_values(df),
                "numerical_stats": self.analyze_numerical_data(df),
                "categorical_stats": self.analyze_categorical_data(df),
                "correlation": self.analyze_correlations(df),
            }
           
            return Response({"eda_results": eda_results}, status=200)
        except Exception as e:
            return Response({"error": str(e)}, status=500)

    def safe_float(self, val):
        try:
            if val is None:
                return None
            f = float(val)
            if math.isinf(f) or math.isnan(f):
                return None
            return f
        except Exception:
            return None

    def analyze_missing_values(self, df):
        missing_info = df.isnull().sum()
        return {
            column: {
                "missing_count": int(count),
                "missing_percentage": float(count / len(df) * 100) if len(df) > 0 else 0
            }
            for column, count in missing_info.items() if count > 0
        }

    def analyze_numerical_data(self, df):
        numerical_columns = df.select_dtypes(include=['int64', 'float64']).columns
        results = {}
        for col in numerical_columns:
            if col in df.columns:
                data = df[col].dropna()
                try:
                    stats = {
                        "mean": self.safe_float(data.mean()) if not data.empty else None,
                        "median": self.safe_float(data.median()) if not data.empty else None,
                        "std": self.safe_float(data.std()) if not data.empty else None,
                        "min": self.safe_float(data.min()) if not data.empty else None,
                        "max": self.safe_float(data.max()) if not data.empty else None,
                    }
                    if not data.empty:
                        hist_counts, hist_bins = np.histogram(data, bins='auto')
                        stats["histogram"] = {
                            "bins": [self.safe_float(b) for b in hist_bins.tolist()],
                            "counts": [int(c) for c in hist_counts.tolist()]
                        }
                    results[col] = stats
                except Exception:
                    results[col] = {"error": "Could not compute statistics"}
        return results

    def analyze_categorical_data(self, df):
        categorical_columns = df.select_dtypes(include=['object']).columns
        results = {}
        for col in categorical_columns:
            if col in df.columns:
                value_counts = df[col].value_counts()
                results[col] = {
                    "unique_values": int(df[col].nunique()),
                    "top_categories": {
                        str(k): int(v) for k, v in value_counts.head(10).items()
                    }
                }
        return results

    def analyze_correlations(self, df):
        numerical_columns = df.select_dtypes(include=['int64', 'float64']).columns
        if len(numerical_columns) > 1:
            df_clean = df[numerical_columns].replace([float('inf'), float('-inf')], None)
            correlation_matrix = df_clean.corr().fillna(0)
            return {
                col1: {
                    col2: self.safe_float(correlation_matrix.loc[col1, col2])
                    for col2 in correlation_matrix.columns if col1 != col2
                }
                for col1 in correlation_matrix.columns
            }
        return {}


class DownloadEDAExcel(APIView):
    parser_classes = [JSONParser]

    def post(self, request):
        eda_results = request.data.get('eda_results')
        personalized_plots = request.data.get('personalized_plots')
        if not eda_results:
            return Response({'error': 'No EDA results provided'}, status=400)

        wb = Workbook()
        ws_summary = wb.active
        ws_summary.title = 'Summary'

        # Write basic info
        basic_info = eda_results.get('basic_info') or eda_results.get('summary', {}).get('dataset_info', {})
        if basic_info:
            ws_summary.append(['Metric', 'Value'])
            for k, v in basic_info.items():
                if isinstance(v, dict):
                    continue
                ws_summary.append([k, v])

        # Write missing values
        missing_values = eda_results.get('missing_values') or eda_results.get('data_quality')
        if missing_values:
            ws_missing = wb.create_sheet('Missing Values')
            ws_missing.append(['Column', 'Missing Count', 'Missing %'])
            for col, stats in missing_values.items():
                ws_missing.append([col, stats.get('missing_count') or stats.get('count'), stats.get('missing_percentage') or stats.get('percentage')])
            # Add bar chart for missing values
            chart = BarChart()
            chart.title = 'Missing Values by Column'
            chart.x_axis.title = 'Column'
            chart.y_axis.title = 'Missing Count'
            data = Reference(ws_missing, min_col=2, min_row=1, max_row=ws_missing.max_row, max_col=2)
            cats = Reference(ws_missing, min_col=1, min_row=2, max_row=ws_missing.max_row)
            chart.add_data(data, titles_from_data=True)
            chart.set_categories(cats)
            ws_missing.add_chart(chart, 'E2')

        # Write numerical stats
        numerical_stats = eda_results.get('numerical_stats') or eda_results.get('detailed_analysis', {}).get('numerical')
        if numerical_stats:
            ws_num = wb.create_sheet('Numerical Stats')
            ws_num.append(['Column', 'Mean', 'Median', 'Std', 'Variance', 'Min', 'Max'])
            for col, stats in numerical_stats.items():
                ws_num.append([
                    col,
                    stats.get('mean'),
                    stats.get('median'),
                    stats.get('std'),
                    stats.get('variance'),
                    stats.get('min'),
                    stats.get('max'),
                ])
            # Add histogram/bar chart for each numerical column with histogram data
            for col, stats in numerical_stats.items():
                if 'histogram' in stats:
                    hist = stats['histogram']
                    ws_hist = wb.create_sheet(f'{col} Histogram')
                    ws_hist.append(['Bin', 'Count'])
                    for bin_val, count in zip(hist['bins'], hist['counts']):
                        ws_hist.append([bin_val, count])
                    chart = BarChart()
                    chart.title = f'Histogram of {col}'
                    chart.x_axis.title = 'Bin'
                    chart.y_axis.title = 'Count'
                    data = Reference(ws_hist, min_col=2, min_row=1, max_row=ws_hist.max_row, max_col=2)
                    cats = Reference(ws_hist, min_col=1, min_row=2, max_row=ws_hist.max_row)
                    chart.add_data(data, titles_from_data=True)
                    chart.set_categories(cats)
                    ws_hist.add_chart(chart, 'E2')

        # Write categorical stats
        categorical_stats = eda_results.get('categorical_stats') or eda_results.get('detailed_analysis', {}).get('categorical')
        if categorical_stats:
            ws_cat = wb.create_sheet('Categorical Stats')
            ws_cat.append(['Column', 'Unique Values', 'Top Category', 'Top Count'])
            for col, stats in categorical_stats.items():
                top_cat = None
                top_count = None
                if 'top_categories' in stats:
                    top_items = list(stats['top_categories'].items())
                    if top_items:
                        top_cat, top_count = top_items[0]
                ws_cat.append([
                    col,
                    stats.get('unique_values'),
                    top_cat,
                    top_count
                ])
            # Pie chart for each categorical column with top_categories
            for col, stats in categorical_stats.items():
                if 'top_categories' in stats:
                    ws_pie = wb.create_sheet(f'{col} Pie')
                    ws_pie.append(['Category', 'Count'])
                    for cat, count in stats['top_categories'].items():
                        ws_pie.append([cat, count])
                    chart = PieChart()
                    chart.title = f'Pie Chart of {col}'
                    data = Reference(ws_pie, min_col=2, min_row=1, max_row=ws_pie.max_row, max_col=2)
                    labels = Reference(ws_pie, min_col=1, min_row=2, max_row=ws_pie.max_row)
                    chart.add_data(data, titles_from_data=True)
                    chart.set_categories(labels)
                    ws_pie.add_chart(chart, 'E2')

        # Write correlation matrix
        correlation = eda_results.get('correlation') or eda_results.get('detailed_analysis', {}).get('correlations')
        if correlation:
            ws_corr = wb.create_sheet('Correlations')
            cols = list(correlation.keys())
            ws_corr.append([''] + cols)
            for row in cols:
                ws_corr.append([row] + [correlation[row].get(col, '') for col in cols])
            # Add heatmap as a bar chart for the first row (Excel doesn't support heatmaps natively)
            chart = BarChart()
            chart.title = 'Correlation (First Variable)'
            chart.x_axis.title = 'Variable'
            chart.y_axis.title = 'Correlation'
            data = Reference(ws_corr, min_col=2, min_row=1, max_row=2, max_col=len(cols)+1)
            cats = Reference(ws_corr, min_col=2, min_row=1, max_row=1, max_col=len(cols)+1)
            chart.add_data(data, titles_from_data=True)
            chart.set_categories(cats)
            ws_corr.add_chart(chart, 'E10')

        # --- Add personalized plots if present ---
        if personalized_plots:
            for idx, plot in enumerate(personalized_plots):
                chart_type = plot.get('chartType', 'bar').lower()
                chart_data = plot.get('chartData', {})
                labels = chart_data.get('labels', [])
                datasets = chart_data.get('datasets', [])
                sheet_title = f'Personalized Plot {idx+1}'
                ws = wb.create_sheet(title=sheet_title[:31])
                # Write header
                ws.append(['Label'] + [ds.get('label', f'Data{i+1}') for i, ds in enumerate(datasets)])
                # Write data
                for i, label in enumerate(labels):
                    row = [label]
                    for ds in datasets:
                        data = ds.get('data', [])
                        row.append(data[i] if i < len(data) else None)
                    ws.append(row)
                # Add chart for all datasets
                if labels and datasets:
                    data_len = len(labels)
                    data_ref = Reference(ws, min_col=2, min_row=1, max_row=data_len+1, max_col=1+len(datasets))
                    cats_ref = Reference(ws, min_col=1, min_row=2, max_row=data_len+1)
                    if chart_type == 'line':
                        chart = LineChart()
                    else:
                        chart = BarChart()
                        chart.type = "col"
                        if chart_type == 'stacked':
                            chart.grouping = 'stacked'
                        elif chart_type in ['stacked100', 'stacked_100', '100stacked', '100%stacked']:
                            chart.grouping = 'percentStacked'
                    chart.title = sheet_title
                    chart.x_axis.title = 'Label'
                    chart.y_axis.title = 'Value'
                    chart.add_data(data_ref, titles_from_data=True)
                    chart.set_categories(cats_ref)
                    ws.add_chart(chart, f"E2")

        # Save to a temporary file and return as response
        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
            wb.save(tmp.name)
            tmp.seek(0)
            response = HttpResponse(tmp.read(), content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
            response['Content-Disposition'] = 'attachment; filename="eda_report.xlsx"'
            return response


class DownloadEDAPPTX(APIView):
    parser_classes = [JSONParser]

    def post(self, request):
        eda_results = request.data.get('eda_results')
        personalized_plots = request.data.get('personalized_plots')
        if not eda_results:
            return Response({'error': 'No EDA results provided'}, status=400)

        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = "EDA Report"
        slide.placeholders[1].text = "Generated by Sanitifi"

        # Numerical plots (histograms)
        numerical_stats = eda_results.get('numerical_stats') or eda_results.get('detailed_analysis', {}).get('numerical')
        if numerical_stats:
            for col, stats in numerical_stats.items():
                hist = stats.get('histogram')
                if hist and hist.get('bins') and hist.get('counts'):
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    slide.shapes.title.text = f"Histogram: {col}"
                    chart_data = CategoryChartData()
                    bins = [str(b) for b in hist['bins'][:-1]]
                    chart_data.categories = bins
                    chart_data.add_series('Count', hist['counts'])
                    x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4)
                    chart = slide.shapes.add_chart(
                        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
                    ).chart
                    chart.has_legend = False
                    chart.value_axis.has_major_gridlines = True
                    chart.category_axis.tick_labels.font.size = Pt(10)
                    chart.value_axis.tick_labels.font.size = Pt(10)

        # Categorical plots (pie charts)
        categorical_stats = eda_results.get('categorical_stats') or eda_results.get('detailed_analysis', {}).get('categorical')
        if categorical_stats:
            for col, stats in categorical_stats.items():
                top_cats = stats.get('top_categories') or stats.get('value_counts')
                if top_cats:
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    slide.shapes.title.text = f"Pie Chart: {col}"
                    chart_data = CategoryChartData()
                    chart_data.categories = list(top_cats.keys())
                    chart_data.add_series('Count', list(top_cats.values()))
                    x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4)
                    chart = slide.shapes.add_chart(
                        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
                    ).chart
                    chart.has_legend = True
                    chart.legend.font.size = Pt(10)

        # --- Add personalized plots if present ---
        if personalized_plots:
            for idx, plot in enumerate(personalized_plots):
                chart_type = plot.get('plot_config', {}).get('chartType', 'bar').lower()
                chart_data = plot.get('chart_data', {})
                labels = chart_data.get('labels', [])
                datasets = chart_data.get('datasets', [])
                
                # Get chart title from plot config or use default
                chart_title = plot.get('plot_config', {}).get('selectedY', f'Personalized Plot {idx+1}')
                
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.shapes.title.text = chart_title
                
                chart_data_obj = CategoryChartData()
                chart_data_obj.categories = labels
                
                # Add each dataset as a series
                for ds in datasets:
                    series_name = ds.get('label', 'Data')
                    series_data = ds.get('data', [])
                    chart_data_obj.add_series(series_name, series_data)
                
                # Choose chart type based on plot_config
                if chart_type == 'line':
                    chart_type_obj = XL_CHART_TYPE.LINE
                elif chart_type == 'stacked':
                    chart_type_obj = XL_CHART_TYPE.COLUMN_STACKED
                elif chart_type in ['stacked100', 'stacked_100', '100stacked', '100%stacked']:
                    chart_type_obj = XL_CHART_TYPE.COLUMN_STACKED_100
                else:
                    chart_type_obj = XL_CHART_TYPE.COLUMN_CLUSTERED
                
                x_pos, y_pos, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4)
                chart = slide.shapes.add_chart(
                    chart_type_obj, x_pos, y_pos, cx, cy, chart_data_obj
                ).chart
                
                # Configure chart appearance
                chart.has_legend = True
                chart.value_axis.has_major_gridlines = True
                chart.category_axis.tick_labels.font.size = Pt(10)
                chart.value_axis.tick_labels.font.size = Pt(10)
                
                # Add axis titles
                chart.category_axis.has_title = True
                chart.value_axis.has_title = True
                chart.category_axis.axis_title.text_frame.text = plot.get('plot_config', {}).get('selectedX', 'Label')
                chart.value_axis.axis_title.text_frame.text = plot.get('plot_config', {}).get('selectedY', 'Value')
                
                # Set chart title
                chart.chart_title.text_frame.text = chart_title

        # Save to a temporary file and return as response
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            prs.save(tmp.name)
            tmp.seek(0)
            response = HttpResponse(tmp.read(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
            response['Content-Disposition'] = 'attachment; filename="eda_report.pptx"'
            return response


class SheetCommitGraph(APIView):
    def post(self, request):
        project_id = request.data.get("project_id")
        user_id = request.data.get("user_id")
        file_name = request.data.get("file_name")
        sheet_name = request.data.get("sheet_name")
        file_type = request.data.get("file_type")

        # Sanitize file_name to remove invalid characters and spaces
        file_name = file_name.split("\\")[-1] if "\\" in file_name else file_name
        # Remove leading/trailing spaces and replace spaces with underscores
        file_name = file_name.strip().replace(" ", "_")

        if not all([project_id, user_id, file_name, sheet_name, file_type]):
            return Response({"error": "Missing required parameters"}, status=400)

        # Check project access (including shared projects)
        has_access, share_object, permission_level = check_project_access(
            user_id, project_id, file_type, file_name, sheet_name
        )
        
        if not has_access:
            return Response({"error": "Access denied to this project or file"}, status=403)

        # Get the project owner's user ID for the correct folder path
        try:
            project = Projects.objects.get(id=project_id)
            project_owner_id = project.user.id
        except Projects.DoesNotExist:
            return Response({"error": "Project not found"}, status=404)

        # Use project owner's folder for git operations
        project_folder = os.path.join(settings.MEDIA_ROOT, f"user_{project_owner_id}/project_{project_id}")
        sheet_path = os.path.join(project_folder, file_type, file_name, f"{sheet_name}")
        sheet_path = os.path.normpath(sheet_path)

        # Get commit log with parents and hashes
        result = subprocess.run(
            ["git", "log", "--pretty=format:%H|%P|%s", "--", sheet_path],
            cwd=project_folder,
            capture_output=True,
            text=True
        )
        if result.returncode != 0:
            return Response({"error": "Failed to fetch commit history"}, status=500)

        commits = []
        for line in result.stdout.strip().split("\n"):
            if not line:
                continue
            parts = line.split("|")
            commit_hash = parts[0]
            parents = parts[1].split() if parts[1] else []
            message = parts[2] if len(parts) > 2 else ""
            commits.append({
                "hash": commit_hash,
                "parents": parents,
                "message": message,
            })

        # Build children map
        hash_to_commit = {c["hash"]: c for c in commits}
        for c in commits:
            c["children"] = []
        for c in commits:
            for parent in c["parents"]:
                if parent in hash_to_commit:
                    hash_to_commit[parent]["children"].append(c["hash"])

        # Log the event
        user = get_logging_user(request, User.objects.get(id=user_id))
        ip = request.META.get('REMOTE_ADDR')
        log_user_action(user, "sheet_commit_graph", details=f"Retrieved commit graph for project {project_id}, file: {file_name}, sheet: {sheet_name}", ip_address=ip)

        return Response({"commits": commits})

class AddDateColumns(APIView):
    def post(self, request):
        try:
            file_type = request.data.get('file_type')
            file_name = request.data.get('file_name')
            project_id = request.data.get('project_id')
            sheet_name = request.data.get('sheet_name')
            datetime_columns = request.data.get('datetime_columns', [])

            if not all([file_type, file_name, project_id, sheet_name, datetime_columns]):
                return Response({'error': 'Missing required fields'}, status=400)

            file_name = os.path.basename(file_name)

            try:
                project = Projects.objects.get(id=project_id)
            except Projects.DoesNotExist:
                return Response({'error': 'Project not found'}, status=404)

            file_path = os.path.join(
                settings.MEDIA_ROOT,
                f"user_{project.user.id}/project_{project.id}/{file_type}/{file_name}/{sheet_name}"
            )
            file_path = os.path.normpath(file_path)

            if not os.path.exists(file_path):
                return Response({'error': 'File not found'}, status=404)

            file_extension = os.path.splitext(sheet_name)[1].lower()
            if file_extension == '.csv':
                df = pd.read_csv(file_path)
            elif file_extension in ['.xlsx', '.xls']:
                df = pd.read_excel(file_path)
            else:
                return Response({'error': 'Unsupported file format'}, status=400)

            added_columns = []
            for col in datetime_columns:
                if col in df.columns:
                    df[col] = df[col].astype(str).str.replace(r'\s+', ' ', regex=True).str.strip()
                    date_series = pd.to_datetime(df[col], errors='coerce')

                    year_col = f'{col}_year'
                    month_col = f'{col}_month'
                    month_year_col = f'{col}_month_year'

                    if year_col not in df.columns:
                        df[year_col] = date_series.dt.year
                        added_columns.append(year_col)
                    if month_col not in df.columns:
                        df[month_col] = date_series.dt.month
                        added_columns.append(month_col)
                    if month_year_col not in df.columns:
                        df[month_year_col] = date_series.dt.strftime('%Y-%m')
                        added_columns.append(month_year_col)
                else:
                    return Response({'error': f'Column {col} not found in sheet.'}, status=400)

            # Save the updated file
            if file_extension == '.csv':
                df.to_csv(file_path, index=False)
            else:
                df.to_excel(file_path, index=False)

            # Sanitize data for JSON serialization
            import numpy as np
            def json_safe(val):
                if val is None:
                    return None
                if isinstance(val, (float, np.floating)):
                    if np.isnan(val) or np.isinf(val):
                        return None
                    return float(val)
                if isinstance(val, (int, np.integer)):
                    return int(val)
                if isinstance(val, (np.generic, np.ndarray)):
                    return val.item() if hasattr(val, "item") else str(val)
                return val

            df = df.replace([np.inf, -np.inf], np.nan)
            df = df.astype(object).where(pd.notnull(df), None)

            # Use a robust list comprehension for sample_data
            sample_data = [[json_safe(cell) for cell in row] for row in df.head(10).values.tolist()]

            # ... after saving the updated file ...
            df = df.replace([np.inf, -np.inf], np.nan)
            df = df.astype(object).where(pd.notnull(df), None)
            def json_safe(val):
                if val is None:
                    return None
                if isinstance(val, (float, np.floating)):
                    if np.isnan(val) or np.isinf(val):
                        return None
                    return float(val)
                if isinstance(val, (int, np.integer)):
                    return int(val)
                if isinstance(val, (np.generic, np.ndarray)):
                    return val.item() if hasattr(val, "item") else str(val)
                return val
            sheet_data = {
                sheet_name: {
                    'columns': df.columns.tolist(),
                    'data': [[json_safe(cell) for cell in row] for row in df.values.tolist()]
                }
            }
            # Log the event
            user = None
            if hasattr(request, 'user') and request.user.is_authenticated:
                user = request.user
            elif project and hasattr(project, 'user'):
                user = project.user
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "add_date_columns", details=f"Date columns added successfully", ip_address=ip)


            return Response({
                'message': 'Date columns added successfully',
                'added_columns': added_columns,
                'sheet_data': sheet_data
            }, status=200)
        except Exception as e:
            return Response({'error': str(e)}, status=500)


class CreateGoogleSheet(APIView):
    def post(self, request):
        file_type = request.data.get('file_type')
        file_name = request.data.get('file_name')
        project_id = request.data.get('project_id')
        sheet_name = request.data.get('sheet_name')
        if not all([file_type, file_name, project_id, sheet_name]):
            return Response({'error': 'Missing required fields'}, status=400)
        file_name = os.path.basename(file_name)
        try:
            project = Projects.objects.get(id=project_id)
        except Projects.DoesNotExist:
            return Response({'error': 'Project not found'}, status=404)
        file_path = os.path.join(
            settings.MEDIA_ROOT,
            f"user_{project.user.id}/project_{project.id}/{file_type}/{file_name}/{sheet_name}"
        )
        file_path = os.path.normpath(file_path)
        if not os.path.exists(file_path):
            return Response({'error': 'File not found'}, status=404)
        file_extension = os.path.splitext(sheet_name)[1].lower()
        try:
            if file_extension == '.csv':
                try:
                    df = pd.read_csv(file_path, encoding='utf-8')
                except UnicodeDecodeError:
                    df = pd.read_csv(file_path, encoding='latin1')
            elif file_extension in ['.xlsx', '.xls']:
                df = pd.read_excel(file_path)
            else:
                return Response({'error': 'Unsupported file format'}, status=400)
            data = [df.columns.tolist()] + df.astype(str).values.tolist()
            sheets_service, drive_service = get_gsheet_service()
            # Use a unique key for this file/sheet
            sheet_key = f"{file_type}|{file_name}|{sheet_name}"
            google_sheet_ids = project.google_sheet_ids or {}
            spreadsheet_id = google_sheet_ids.get(sheet_key)
            if spreadsheet_id:
                # Sheet already exists, update its data
                try:
                    sheets_service.spreadsheets().values().update(
                        spreadsheetId=spreadsheet_id,
                        range="Sheet1!A1",
                        valueInputOption="RAW",
                        body={"values": data}
                    ).execute()
                    # Make sure permissions are still correct (in case sheet was created before this logic)
                    try:
                        drive_service.permissions().create(
                            fileId=spreadsheet_id,
                            body={"role": "writer", "type": "anyone"},
                        ).execute()
                    except Exception:
                        pass
                except Exception as e:
                    # If update fails (e.g., sheet deleted), create a new one
                    spreadsheet_id = None
            if not spreadsheet_id:
                # Create new Google Sheet, upload data, and store mapping
                spreadsheet = sheets_service.spreadsheets().create(
                    body={"properties": {"title": f"EditData_{sheet_name}"}},
                    fields="spreadsheetId"
                ).execute()
                spreadsheet_id = spreadsheet.get('spreadsheetId')
                sheets_service.spreadsheets().values().update(
                    spreadsheetId=spreadsheet_id,
                    range="Sheet1!A1",
                    valueInputOption="RAW",
                    body={"values": data}
                ).execute()
                drive_service.permissions().create(
                    fileId=spreadsheet_id,
                    body={"role": "writer", "type": "anyone"},
                ).execute()
                # Save mapping
                google_sheet_ids[sheet_key] = spreadsheet_id
                project.google_sheet_ids = google_sheet_ids
                project.save()
            sheet_url = f"https://docs.google.com/spreadsheets/d/{spreadsheet_id}/edit"
            return Response({"sheet_url": sheet_url, "sheet_id": spreadsheet_id}, status=200)
        except Exception as e:
            return Response({'error': f'Google Sheets error: {str(e)}'}, status=500)

class UpdateFromGoogleSheet(APIView):
    def post(self, request):
        file_type = request.data.get('file_type')
        file_name = request.data.get('file_name')
        project_id = request.data.get('project_id')
        sheet_name = request.data.get('sheet_name')
        google_sheet_id = request.data.get('google_sheet_id')
        if not all([file_type, file_name, project_id, sheet_name, google_sheet_id]):
            return Response({'error': 'Missing required fields'}, status=400)
        file_name = os.path.basename(file_name)
        try:
            project = Projects.objects.get(id=project_id)
        except Projects.DoesNotExist:
            return Response({'error': 'Project not found'}, status=404)
        file_path = os.path.join(
            settings.MEDIA_ROOT,
            f"user_{project.user.id}/project_{project.id}/{file_type}/{file_name}/{sheet_name}"
        )
        file_path = os.path.normpath(file_path)
        file_extension = os.path.splitext(sheet_name)[1].lower()
        try:
            sheets_service, _ = get_gsheet_service()
            result = sheets_service.spreadsheets().values().get(
                spreadsheetId=google_sheet_id,
                range="Sheet1"
            ).execute()
            values = result.get('values', [])
            if not values or len(values) < 2:
                return Response({'error': 'No data found in Google Sheet'}, status=400)
            df = pd.DataFrame(values[1:], columns=values[0])
            # Save to backend file (csv or xlsx)
            if file_extension == '.csv':
                try:
                    df.to_csv(file_path, index=False, encoding='utf-8')
                except Exception as e:
                    return Response({'error': f'Failed to write CSV: {str(e)}'}, status=500)
            elif file_extension in ['.xlsx', '.xls']:
                try:
                    df.to_excel(file_path, index=False, engine='openpyxl')
                except Exception as e:
                    return Response({'error': f'Failed to write Excel: {str(e)}'}, status=500)
            else:
                return Response({'error': 'Unsupported file format'}, status=400)

            # --- GIT COMMIT LOGIC (MATCH Save.commit_to_git) ---
            try:
                user = None
                if hasattr(request, 'user') and request.user.is_authenticated:
                    user = request.user
                elif project and hasattr(project, 'user'):
                    user = project.user
                project_folder = os.path.join(settings.MEDIA_ROOT, f"user_{project.user.id}/project_{project.id}")
                # Initialize git repo if it doesn't exist
                if not os.path.exists(os.path.join(project_folder, ".git")):
                    subprocess.run(["git", "init"], cwd=project_folder)
                    subprocess.run(["git", "config", "user.name", user.name], cwd=project_folder)
                    subprocess.run(["git", "config", "user.email", user.email], cwd=project_folder)
                file_path_relative = os.path.join(file_type, file_name)
                subprocess.run(["git", "add", file_path_relative], cwd=project_folder)
                commit_message = f"google sheet - {user.id}/{project_id}/{file_type}/{file_name}/{sheet_name}"
                subprocess.run(["git", "commit", "-m", commit_message], cwd=project_folder)
            except Exception as e:
                print(f"Git commit failed: {str(e)}")
            # --- END GIT COMMIT LOGIC ---

            # Read the file again to ensure it's saved and to return the latest data
            try:
                if file_extension == '.csv':
                    updated_df = pd.read_csv(file_path, encoding='utf-8')
                else:
                    updated_df = pd.read_excel(file_path, engine='openpyxl')
            except Exception as e:
                return Response({'error': f'Failed to read updated file: {str(e)}'}, status=500)
            # Sanitize data for JSON
            import numpy as np
            def json_safe(val):
                if val is None:
                    return None
                if isinstance(val, (float, np.floating)):
                    if np.isnan(val) or np.isinf(val):
                        return None
                    return float(val)
                if isinstance(val, (int, np.integer)):
                    return int(val)
                if isinstance(val, (np.generic, np.ndarray)):
                    return val.item() if hasattr(val, "item") else str(val)
                return val
            updated_df = updated_df.replace([np.inf, -np.inf], np.nan)
            updated_df = updated_df.astype(object).where(pd.notnull(updated_df), None)
            safe_data = [[json_safe(cell) for cell in row] for row in updated_df.values.tolist()]
            user = None
            if hasattr(request, 'user') and request.user.is_authenticated:
                user = request.user
            elif project and hasattr(project, 'user'):
                user = project.user
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "update_from_google_sheet", details=f"Data updated from Google Sheet successfully", ip_address=ip)

            return Response({
                'message': 'Data updated from Google Sheet successfully',
                'columns': updated_df.columns.tolist(),
                'data': safe_data
            }, status=200)
        except Exception as e:
            return Response({'error': f'Google Sheets error: {str(e)}'}, status=500)


class UpdateSheetData(APIView):
    """
    API endpoint to update sheet data directly.
    """
    parser_classes = [JSONParser]
    
    def post(self, request):
        try:
            # Extract required parameters
            file_type = request.data.get('file_type')
            file_name = request.data.get('file_name')
            project_id = request.data.get('project_id')
            sheet_name = request.data.get('sheet_name')
            user_id = request.data.get('user_id')
            operation_type = request.data.get('operation_type', "Sheet data updated")

            # Extract update data
            update_data = request.data.get('update_data', {})
            
            # Validate required fields
            if not all([file_type, file_name, project_id, sheet_name, user_id]):
                return Response({
                    'error': 'Missing required fields: file_type, file_name, project_id, sheet_name, user_id'
                }, status=400)
            
            file_name = os.path.basename(file_name)
            
            # Check project access
            has_access, share_object, permission_level = check_project_access(
                user_id, project_id, file_type, file_name, sheet_name
            )
            
            if not has_access:
                return Response({
                    'error': 'Access denied. You do not have permission to update this sheet.'
                }, status=403)
            
            # Check if user has edit permissions
            if permission_level == 'view':
                return Response({
                    'error': 'You only have view permissions for this sheet.'
                }, status=403)
            
            # Get project
            try:
                project = Projects.objects.get(id=project_id)
            except Projects.DoesNotExist:
                return Response({'error': 'Project not found'}, status=404)
            
            # Construct file path
            if file_type == 'concatenated':
                concatenated_base = os.path.join(settings.MEDIA_ROOT, f"user_{project.user.id}/project_{project.id}/concatenated")
                found = False
                file_path = None
                if os.path.exists(concatenated_base):
                    for folder in os.listdir(concatenated_base):
                        folder_path = os.path.join(concatenated_base, folder)
                        candidate = os.path.join(folder_path, sheet_name)
                        if os.path.isfile(candidate):
                            file_path = candidate
                            found = True
                            break
                if not found:
                    return Response({"error": "File not found"}, status=404)
            else:
                file_path = os.path.join(
                    settings.MEDIA_ROOT,
                    f"user_{project.user.id}/project_{project.id}/{file_type}/{file_name}/{sheet_name}"
                )
                file_path = os.path.normpath(file_path)
                if not os.path.exists(file_path):
                    return Response({"error": "File not found"}, status=404)
            
            # Only support full update
            if 'columns' in update_data and 'data' in update_data:
                new_columns = update_data['columns']
                new_data = update_data['data']
                
                # Validate data structure
                if not isinstance(new_columns, list) or not isinstance(new_data, list):
                    return Response({
                        'error': 'Invalid data format. columns and data must be lists.'
                    }, status=400)
                
                # Create new DataFrame
                df = pd.DataFrame(new_data, columns=new_columns)
            else:
                return Response({
                    'error': 'For full update, both columns and data are required.'
                }, status=400)
            
            # Save the updated file
            file_extension = os.path.splitext(sheet_name)[1].lower()
            try:
                if file_extension == '.csv':
                    df.to_csv(file_path, index=False, encoding='utf-8')
                elif file_extension in ['.xlsx', '.xls']:
                    df.to_excel(file_path, index=False, engine='openpyxl')
                else:
                    return Response({"error": "Unsupported file format"}, status=400)
            except Exception as e:
                return Response({
                    'error': f'Failed to save updated file: {str(e)}'
                }, status=500)
            
            # Commit to git if project has git tracking
            project_folder = os.path.join(settings.MEDIA_ROOT, f"user_{project.user.id}/project_{project.id}")
            if os.path.exists(os.path.join(project_folder, '.git')):
                try:
                    self.commit_to_git(
                        project_folder,
                        project.user,
                        project_id,
                        file_type,
                        file_name,
                        sheet_name,
                        operation_type
                    )
                except Exception as e:
                    print(f"Warning: Git commit failed: {str(e)}")
            
            # Prepare response data
            def json_safe(val):
                if val is None:
                    return None
                if isinstance(val, (float, np.floating)):
                    if np.isnan(val) or np.isinf(val):
                        return None
                    return float(val)
                if isinstance(val, (int, np.integer)):
                    return int(val)
                if isinstance(val, (np.generic, np.ndarray)):
                    return val.item() if hasattr(val, "item") else str(val)
                return val
            
            # Clean data for JSON response
            df_clean = df.replace([np.inf, -np.inf], np.nan)
            df_clean = df_clean.astype(object).where(pd.notnull(df_clean), None)
            safe_data = [[json_safe(cell) for cell in row] for row in df_clean.values.tolist()]
            
            # Log the action
            user = User.objects.get(id=user_id)
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "update_sheet_data", 
                          details=f"Sheet data updated successfully - File: {file_name}, Sheet: {sheet_name}", 
                          ip_address=ip)
            
            return Response({
                'message': 'Sheet data updated successfully',
                'columns': df.columns.tolist(),
                'data': safe_data,
                'rows_count': len(df),
                'columns_count': len(df.columns),
                'operation_type': operation_type
            }, status=200)
            
        except Exception as e:
            return Response({
                'error': f'An error occurred while updating sheet data: {str(e)}'
            }, status=500)
    
    def commit_to_git(self, project_folder, user, project_id, file_type, file_name, sheet_name, operation_type):
        commit_msg = f"{operation_type} on {sheet_name} by {user.username}"
        subprocess.run(["git", "add", "."], cwd=project_folder)
        subprocess.run(["git", "commit", "-m", commit_msg], cwd=project_folder)





class DeleteFile(APIView):
    def delete(self, request):
        project_id = request.data.get('project_id')
        file_type = request.data.get('file_type')
        file_name = request.data.get('file_name')

        if not all([project_id, file_type, file_name]):
            return Response({'error': 'Missing required fields'}, status=400)

        try:
            project = Projects.objects.get(id=project_id)
        except Projects.DoesNotExist:
            return Response({'error': 'Project not found'}, status=404)

        if file_type == 'concatenated':
            # Search all subfolders for the file_name
            concatenated_base = os.path.join(settings.MEDIA_ROOT, f"user_{project.user.id}/project_{project.id}/concatenated")
            found = False
            file_path = None
            folder_path = None
            if os.path.exists(concatenated_base):
                for folder in os.listdir(concatenated_base):
                    candidate_folder = os.path.join(concatenated_base, folder)
                    candidate = os.path.join(candidate_folder, file_name)
                    if os.path.isfile(candidate):
                        file_path = candidate
                        folder_path = candidate_folder
                        found = True
                        break
            if not found:
                return Response({'error': 'File or folder not found'}, status=404)
        else:
            file_path = os.path.join(
                settings.MEDIA_ROOT,
                f"user_{project.user.id}/project_{project.id}/{file_type}/{file_name}"
            )
            file_path = os.path.normpath(file_path)
            folder_path = file_path
            if not os.path.exists(file_path):
                return Response({'error': 'File or folder not found'}, status=404)

        try:
            if os.path.isfile(file_path):
                os.remove(file_path)
                # If concatenated, remove the folder if empty
                if file_type == 'concatenated' and folder_path and os.path.isdir(folder_path):
                    if not os.listdir(folder_path):
                        import shutil
                        shutil.rmtree(folder_path)
            else:
                import shutil
                shutil.rmtree(file_path)

            return Response({'message': 'File or folder deleted successfully'}, status=200)
        except Exception as e:
            return Response({'error': str(e)}, status=500)

class TimeoutException(Exception):
    pass

def run_with_timeout(func, args=(), kwargs={}, timeout_duration=10):
    """Run a function with a timeout using threading"""
    result = []
    error = []
    
    def target():
        try:
            result.append(func(*args, **kwargs))
        except Exception as e:
            error.append(e)
    
    thread = threading.Thread(target=target)
    thread.daemon = True
    thread.start()
    thread.join(timeout_duration)
    
    if thread.is_alive():
        # Thread is still running, timeout occurred
        raise TimeoutException("Script execution timed out")
    
    if error:
        raise error[0]
    
    return result[0]

class CustomScriptRun(APIView):
    def post(self, request):
        try:
            # Validate required fields
            required_fields = ['file_type', 'file_name', 'project_id', 'sheet_name', 'script_content']
            if not all(request.data.get(field) for field in required_fields):
                return Response({'error': 'Missing required fields'}, status=400)

            file_type = request.data['file_type']
            file_name = request.data['file_name']
            project_id = request.data['project_id']
            sheet_name = request.data['sheet_name']
            script_content = request.data['script_content']
            action = request.data.get('action', 'preview')

            # Validate action parameter
            if action not in ['preview', 'save']:
                return Response({'error': 'Invalid action parameter'}, status=400)

            # Validate script content length
            if len(script_content) > 10000:  # 10KB limit
                return Response({'error': 'Script content too large (max 10KB)'}, status=400)

            # Check for dangerous patterns in script
            dangerous_patterns = [
                r'import\s+(os|sys|subprocess|shutil|glob|socket|requests|urllib|pickle|ctypes)',
                r'__import__\s*\(',
                r'eval\s*\(',
                r'exec\s*\(',
                r'open\s*\(',
                r'file\s*\(',
                r'\.to_csv\s*\(',
                r'\.to_excel\s*\(',
                r'\.to_json\s*\(',
                r'\.to_pickle\s*\(',
                r'while\s+True:',
                r'for\s+.+\s+in\s+.+:\s*while\s+True:',
            ]

            for pattern in dangerous_patterns:
                if re.search(pattern, script_content, re.IGNORECASE):
                    return Response({'error': f'Script contains potentially dangerous pattern: {pattern}'}, status=400)

            # Get project
            try:
                project = Projects.objects.get(id=project_id)
            except Projects.DoesNotExist:
                return Response({'error': 'Project not found'}, status=404)

            # Validate file path
            file_path = os.path.join(
                settings.MEDIA_ROOT,
                f"user_{project.user.id}/project_{project.id}/{file_type}/{file_name}/{sheet_name}"
            )
            file_path = os.path.normpath(file_path)

            # Security check to prevent directory traversal
            if not file_path.startswith(os.path.normpath(settings.MEDIA_ROOT)):
                return Response({'error': 'Invalid file path'}, status=400)

            if not os.path.exists(file_path):
                return Response({'error': 'File not found'}, status=404)

            # Read the file
            file_extension = os.path.splitext(sheet_name)[1].lower()
            try:
                # Get file size in MB
                file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
                
                if file_extension == '.csv':
                    try:
                        df = pd.read_csv(file_path, encoding='utf-8')
                    except UnicodeDecodeError:
                        df = pd.read_csv(file_path, encoding='latin1')
                elif file_extension in ['.xlsx', '.xls']:
                    df = pd.read_excel(file_path)
                else:
                    return Response({'error': 'Unsupported file format'}, status=400)
            except Exception as e:
                return Response({'error': f'Error reading file: {str(e)}'}, status=400)

            # Create a secure execution environment
            safe_globals = {
                'pd': pd,
                'np': np,
                'plt': plt,
                'sns': sns,
                'plotly': plotly,
                'px': plotly.express,
                'go': plotly.graph_objects,
                'scipy': scipy,
                'sklearn': sklearn,
                'datetime': datetime,
                'statsmodels': statsmodels,
                'df': df.copy(),  # Work with a copy of the dataframe
                '__builtins__': {
                    'None': None,
                    'False': False,
                    'True': True,
                    'bool': bool,
                    'int': int,
                    'float': float,
                    'str': str,
                    'list': list,
                    'dict': dict,
                    'tuple': tuple,
                    'set': set,
                    'len': len,
                    'range': range,
                    'sum': sum,
                    'min': min,
                    'max': max,
                    'sorted': sorted,
                    'enumerate': enumerate,
                    'zip': zip,
                    'isinstance': isinstance,
                    'KeyError': KeyError,
                    'ValueError': ValueError,
                    'TypeError': TypeError,
                    'IndexError': IndexError,
                    'Exception': Exception,
                },
                'print': lambda *args: None,  # Disable print
            }

            def execute_script():
                exec(script_content, safe_globals)
                return safe_globals.get('df')

            try:
                # Execute the script with timeout protection
                modified_df = run_with_timeout(execute_script, timeout_duration=10)
                
                if modified_df is None or not isinstance(modified_df, pd.DataFrame):
                    return Response({'error': 'Script must modify the dataframe variable "df"'}, status=400)

                # Validate dataframe size (prevent memory exhaustion)
                if len(modified_df) > 1000000:  # 1 million rows max
                    return Response({'error': 'Resulting dataframe too large (max 1 million rows)'}, status=400)

                # Handle preview or save action
                if action == 'preview':
                    # For preview, return complete data if file is small enough
                    if file_size_mb < 50:
                        preview_df = modified_df  # Use complete dataframe
                    else:
                        preview_df = modified_df.head(100)  # Use first 100 rows for large files
                    
                    preview_data = {
                        'columns': preview_df.columns.tolist(),
                        'data': make_json_safe(preview_df.replace([np.inf, -np.inf], np.nan)
                                     .fillna('NA')
                                     .values.tolist()),
                        'total_rows': len(modified_df),
                        'is_complete_data': file_size_mb < 50,
                        'file_size_mb': round(file_size_mb, 2)
                    }
                    user = None
                    if hasattr(request, 'user') and request.user.is_authenticated:
                        user = request.user
                    elif project and hasattr(project, 'user'):
                        user = project.user
                    ip = request.META.get('REMOTE_ADDR')
                    log_user_action(user, "custom_script_run", details=f"Script executed successfully (preview)", ip_address=ip)
                    return Response(make_json_safe({
                        'message': 'Script executed successfully (preview)',
                        'preview_data': preview_data
                    }), status=200)
                else:  # save action
                    # Save the modified dataframe
                    try:
                        if file_extension == '.csv':
                            modified_df.to_csv(file_path, index=False, encoding='utf-8')
                        else:
                            modified_df.to_excel(file_path, index=False, engine='openpyxl')

                        # Commit changes to git
                        project_folder = os.path.join(settings.MEDIA_ROOT, f"user_{project.user.id}/project_{project_id}")
                        
                        try:
                            # Initialize git if it doesn't exist
                            if not os.path.exists(os.path.join(project_folder, ".git")):
                                subprocess.run(["git", "init"], cwd=project_folder, check=True)
                                subprocess.run(["git", "config", "user.name", project.user.name], cwd=project_folder, check=True)
                                subprocess.run(["git", "config", "user.email", project.user.email], cwd=project_folder, check=True)
                        except Exception as e:
                            print(f"Error initializing git: {str(e)}")

                        # Commit the changes
                        try:
                            subprocess.run(["git", "add", "."], cwd=project_folder, check=True)
                            subprocess.run(["git", "commit", "-m", f"Script execution: {sheet_name}"], cwd=project_folder, check=True)
                        except Exception as e:
                            print(f"Error committing to git: {str(e)}")

                        # Prepare sheet data for response (like preview)
                        # --- Add column_types for frontend ---
                        def get_column_types(df):
                            types = {}
                            for col in df.columns:
                                dtype = str(df[col].dtype)
                                if dtype.startswith('int'):
                                    types[col] = 'int'
                                elif dtype.startswith('float'):
                                    types[col] = 'float'
                                elif dtype.startswith('bool'):
                                    types[col] = 'bool'
                                elif dtype.startswith('datetime'):
                                    types[col] = 'datetime'
                                else:
                                    types[col] = 'str'
                            return types
                        # --- End column_types ---
                        sheet_data = {
                            'columns': modified_df.columns.tolist(),
                            'data': make_json_safe(modified_df.replace([np.inf, -np.inf], np.nan).fillna('NA').values.tolist()),
                            'column_types': get_column_types(modified_df),
                            'total_rows': len(modified_df)
                        }

                        # Wrap sheet_data in a dict keyed by sheet_name for frontend compatibility
                        user = None
                        if hasattr(request, 'user') and request.user.is_authenticated:
                            user = request.user
                        elif project and hasattr(project, 'user'):
                            user = project.user
                        ip = request.META.get('REMOTE_ADDR')
                        log_user_action(user, "custom_script_run", details=f"Script executed and saved successfully", ip_address=ip)
                        return Response(make_json_safe({
                            'message': 'Script executed and saved successfully',
                            'file_path': file_path,
                            'total_rows': len(modified_df),
                            'sheet_data': {sheet_name: sheet_data}
                        }), status=200)

                    except Exception as e:
                        return Response({'error': f'Error saving file: {str(e)}'}, status=400)

            except TimeoutException:
                return Response({'error': 'Script execution timed out (max 10 seconds)'}, status=400)
            except Exception as e:
                return Response({'error': f'Error executing script: {str(e)}'}, status=400)

        except Exception as e:
            return Response({'error': f'Unexpected error: {str(e)}'}, status=400)


class SaveScript(APIView):
    def post(self, request):
        try:
            # Validate required fields
            required_fields = ['title', 'script_content', 'user_id']
            if not all(request.data.get(field) for field in required_fields):
                return Response({'error': 'Missing required fields'}, status=400)

            # Get user from user_id in payload
            try:
                user = User.objects.get(id=request.data['user_id'])
            except User.DoesNotExist:
                return Response({'error': 'User not found'}, status=404)

            # Create new script
            script = SavedScript(
                user=user,
                title=request.data['title'],
                script_content=request.data['script_content'],
                description=request.data.get('description', '')
            )
            script.save()

            user = None
            if hasattr(request, 'user') and request.user.is_authenticated:
                user = request.user
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "save_script", details=f"Script saved successfully", ip_address=ip)

            return Response({
                'message': 'Script saved successfully',
                'script_id': script.id,
                'title': script.title,
                'created_at': script.created_at,
                'updated_at': script.updated_at
            }, status=201)

        except Exception as e:
            return Response({'error': f'Error saving script: {str(e)}'}, status=400)



class FetchScripts(APIView):
    def post(self, request):  # Changed to POST to accept user_id in payload
        try:
            # Validate user_id in payload
            if not request.data.get('user_id'):
                return Response({'error': 'Missing user_id'}, status=400)

            try:
                user = User.objects.get(id=request.data['user_id'])
            except User.DoesNotExist:
                return Response({'error': 'User not found'}, status=404)

            # Get all scripts for the user
            scripts = SavedScript.objects.filter(user=user)
            
            scripts_list = [{
                'id': script.id,
                'title': script.title,
                'script_content': script.script_content,
                'description': script.description,
                'created_at': script.created_at,
                'updated_at': script.updated_at,
            } for script in scripts]

           

            return Response({
                'scripts': scripts_list
            }, status=200)

        except Exception as e:
            return Response({'error': f'Error fetching scripts: {str(e)}'}, status=400)

class ConcatenateSheets(APIView):
    def post(self, request):
        try:
            # Extract data from request
            project_id = request.data.get('project_id')
            file_name = request.data.get('file_name')
            file_type = request.data.get('file_type')
            sheet_names = request.data.get('sheet_names')  # List of sheet names to concatenate
            
            if not all([project_id, file_name, file_type, sheet_names]):
                return Response({'error': 'Missing required fields'}, status=400)

            file_name = os.path.basename(file_name)

            # Generate new sheet name from concatenated sheet names
            # Remove file extensions and join with '+'
            base_names = [os.path.splitext(name)[0] for name in sheet_names]
            new_sheet_name = '+'.join(base_names) + '.csv'

            # Get project
            try:
                project = Projects.objects.get(id=project_id)
            except Projects.DoesNotExist:
                return Response({'error': 'Project not found'}, status=404)

            # Construct base folder path
            base_folder = os.path.join(settings.MEDIA_ROOT, f'user_{project.user.id}/project_{project.id}/{file_type}')
            file_folder = os.path.join(base_folder, file_name)

            # Read and concatenate all sheets
            dfs = []
            for sheet_name in sheet_names:
                sheet_path = os.path.join(file_folder, sheet_name)
                sheet_path = os.path.normpath(sheet_path)
                
                if not os.path.exists(sheet_path):
                    return Response({'error': f'Sheet not found: {sheet_name}'}, status=404)

                # Read the sheet based on file extension
                file_extension = os.path.splitext(sheet_name)[1].lower()
                if file_extension == '.csv':
                    try:
                        df = pd.read_csv(sheet_path, encoding='utf-8')
                    except UnicodeDecodeError:
                        df = pd.read_csv(sheet_path, encoding='latin1')
                elif file_extension in ['.xlsx', '.xls']:
                    df = pd.read_excel(sheet_path)
                else:
                    return Response({'error': f'Unsupported file format for sheet: {sheet_name}'}, status=400)

                # Add source sheet name as a column
                df['source_sheet'] = sheet_name
                dfs.append(df)

            # Concatenate all dataframes
            if not dfs:
                return Response({'error': 'No valid sheets to concatenate'}, status=400)

            concatenated_df = pd.concat(dfs, ignore_index=True)

            # Handle infinite and NaN values before saving
            concatenated_df = concatenated_df.replace([np.inf, -np.inf], np.nan)
            concatenated_df = concatenated_df.fillna('NA')

            # Save the concatenated dataframe as CSV
            new_sheet_path = os.path.join(file_folder, new_sheet_name)
            new_sheet_path = os.path.normpath(new_sheet_path)

            # Save as CSV
            concatenated_df.to_csv(new_sheet_path, index=False, encoding='utf-8')

            # Git commit
            self.commit_to_git(
                os.path.join(settings.MEDIA_ROOT, f'user_{project.user.id}/project_{project.id}'),
                project.user,
                project.id,
                file_type,
                file_name,
                new_sheet_name
            )

            # Convert DataFrame to JSON-safe format
            safe_df = concatenated_df.copy()
            # Convert any remaining problematic values to strings
            for col in safe_df.columns:
                safe_df[col] = safe_df[col].astype(str)

            user = get_logging_user(request, getattr(project, 'user', None))
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "concatenate_sheets", details=f"Sheets concatenated successfully", ip_address=ip)

            return Response({
                'message': 'Sheets concatenated successfully',
                'new_sheet': {
                    'name': new_sheet_name,
                    'columns': safe_df.columns.tolist(),
                    'data': safe_df.values.tolist()
                }
            }, status=200)

        except Exception as e:
            return Response({'error': str(e)}, status=500)

    def commit_to_git(self, project_folder, user, project_id, file_type, file_name, sheet_name):
        try:
            if not os.path.exists(os.path.join(project_folder, ".git")):
                subprocess.run(["git", "init"], cwd=project_folder)
                subprocess.run(["git", "config", "user.name", user.name], cwd=project_folder)
                subprocess.run(["git", "config", "user.email", user.email], cwd=project_folder)

            file_path_relative = os.path.join(file_type, file_name, sheet_name).replace("\\", "/")
            subprocess.run(["git", "add", file_path_relative], cwd=project_folder)
            commit_message = f"concatenate - {user.id}/{project_id}/{file_type}/{file_name}/{sheet_name}"
            subprocess.run(["git", "commit", "-m", commit_message], cwd=project_folder)
            print(f"Git commit done for {file_path_relative}")
        except Exception as e:
            print(f"Git commit failed: {str(e)}")

class SavePlot(APIView):
    parser_classes = [JSONParser]

    def post(self, request):
        try:
            # Extract data from request
            user_id = request.data.get('user_id')
            project_id = request.data.get('project_id')
            file_type = request.data.get('file_type')
            file_name = request.data.get('file_name')
            sheet_name = request.data.get('sheet_name')
            plot_name = request.data.get('plot_name', f'Plot {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}')  # Default name if not provided
            plot_config = request.data.get('plot_config')
            chart_data = request.data.get('chart_data')
            chart_options = request.data.get('chart_options')

            # Validate required fields
            if not all([user_id, project_id, file_type, file_name, sheet_name, plot_config, chart_data, chart_options]):
                return Response({
                    'error': 'Missing required fields'
                }, status=400)

            # Validate file_type
            if file_type not in ['kpi', 'media']:
                return Response({
                    'error': 'Invalid file_type. Must be either "kpi" or "media"'
                }, status=400)

            # Check if user has access to this project/file
            has_access, share_object, permission_level = check_project_access(
                user_id, project_id, file_type, file_name, sheet_name
            )
            
            if not has_access:
                return Response({
                    'error': 'Access denied. You don\'t have permission to save plots for this project/file.'
                }, status=403)
            
            # Check if user has edit permissions
            if permission_level == 'view':
                return Response({
                    'error': 'Access denied. You only have view permissions. Edit permissions required to save plots.'
                }, status=403)

            # Get user and project
            try:
                user = User.objects.get(id=user_id)
                project = Projects.objects.get(id=project_id)
            except User.DoesNotExist:
                return Response({'error': 'User not found'}, status=404)
            except Projects.DoesNotExist:
                return Response({'error': 'Project not found'}, status=404)

            # Create or update saved plot
            saved_plot, created = SavedPlot.objects.update_or_create(
                user=user,
                project=project,
                file_type=file_type,
                file_name=file_name,
                sheet_name=sheet_name,
                plot_name=plot_name,
                defaults={
                    'plot_config': plot_config,
                    'chart_data': chart_data,
                    'chart_options': chart_options
                }
            )

            user = None
            if hasattr(request, 'user') and request.user.is_authenticated:
                user = request.user
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "save_plot", details=f"Plot saved successfully", ip_address=ip)

            return Response({
                'message': 'Plot saved successfully',
                'plot_id': saved_plot.id,
                'created': created,
                'updated_at': saved_plot.updated_at,
                'access_info': {
                    'permission_level': permission_level,
                    'is_owner': share_object is None
                }
            }, status=200)

        except Exception as e:
            return Response({'error': str(e)}, status=500)

@method_decorator(csrf_exempt, name='dispatch')
class FetchPlots(APIView):
    parser_classes = [JSONParser]

    def post(self, request):
        print("FetchPlots API called with data:", request.data)  # Debug log
        try:
            # Extract data from request
            user_id = request.data.get('user_id')
            project_id = request.data.get('project_id')
            file_type = request.data.get('file_type')
            file_name = request.data.get('file_name')
            sheet_name = request.data.get('sheet_name')

            print(f"Extracted data: user_id={user_id}, project_id={project_id}, file_type={file_type}, file_name={file_name}, sheet_name={sheet_name}")  # Debug log

            # Validate required fields
            if not all([user_id, project_id]):
                missing_fields = []
                if not user_id: missing_fields.append('user_id')
                if not project_id: missing_fields.append('project_id')
                print(f"Missing required fields: {missing_fields}")  # Debug log
                return Response({
                    'error': f'Missing required fields: {", ".join(missing_fields)}'
                }, status=400)

            # Check if user has access to this project
            has_access, share_object, permission_level = check_project_access(
                user_id, project_id, file_type, file_name, sheet_name
            )
            
            if not has_access:
                return Response({
                    'error': 'Access denied. You don\'t have permission to view plots for this project/file.'
                }, status=403)

            # Get user and project
            try:
                user = User.objects.get(id=user_id)
                project = Projects.objects.get(id=project_id)
                print(f"Found user: {user.username}, project: {project.name}")  # Debug log
            except User.DoesNotExist:
                print(f"User not found with id: {user_id}")  # Debug log
                return Response({'error': 'User not found'}, status=404)
            except Projects.DoesNotExist:
                print(f"Project not found with id: {project_id}")  # Debug log
                return Response({'error': 'Project not found'}, status=404)

            # Build filter for plots based on access level
            plot_filter = {
                'project': project
            }
            
            # If this is a file-specific share, only show plots for that file
            if share_object and share_object.share_type == 'file':
                plot_filter.update({
                    'file_type': share_object.file_type,
                    'file_name': share_object.file_name
                })
                
                # If specific sheet is shared, only show plots for that sheet
                if share_object.sheet_name:
                    plot_filter['sheet_name'] = share_object.sheet_name
            else:
                # Project-level sharing or owner access - apply optional filters
                if file_type:
                    plot_filter['file_type'] = file_type
                if file_name:
                    plot_filter['file_name'] = file_name
                if sheet_name:
                    plot_filter['sheet_name'] = sheet_name

            # Debug: Print all plots for this user/project combination
            all_plots = SavedPlot.objects.filter(project=project)
            print(f"Total plots for project {project_id}: {all_plots.count()}")
            for plot in all_plots:
                print(f"Plot ID: {plot.id}, Name: {plot.plot_name}, File: {plot.file_name}, Sheet: {plot.sheet_name}, Updated: {plot.updated_at}")

            # Fetch saved plots with filters
            saved_plots = SavedPlot.objects.filter(**plot_filter).order_by('-updated_at')

            print(f"Found {saved_plots.count()} plots matching filters")
            print(f"Filter criteria: {plot_filter}")
            
            # Debug: Print details of each matching plot
            for plot in saved_plots:
                print(f"Matching Plot ID: {plot.id}")
                print(f"  Plot Name: {plot.plot_name}")
                print(f"  File Type: {plot.file_type}")
                print(f"  File Name: {plot.file_name}")
                print(f"  Sheet Name: {plot.sheet_name}")
                print(f"  Updated At: {plot.updated_at}")
                print("  Plot Config:", plot.plot_config)
                print("---")

            # Format response
            plots_data = [{
                'id': plot.id,
                'plot_name': plot.plot_name,
                'file_type': plot.file_type,
                'file_name': plot.file_name,
                'sheet_name': plot.sheet_name,
                'plot_config': plot.plot_config,
                'chart_data': plot.chart_data,
                'chart_options': plot.chart_options,
                'created_at': plot.created_at,
                'updated_at': plot.updated_at
            } for plot in saved_plots]

            # --- Fix: Make filter_criteria serializable ---
            serializable_filter_criteria = {}
            for k, v in plot_filter.items():
                if k == 'project' and hasattr(v, 'id'):
                    serializable_filter_criteria[k] = v.id
                else:
                    serializable_filter_criteria[k] = v

            return Response({
                'plots': plots_data,
                'access_info': {
                    'permission_level': permission_level,
                    'is_owner': share_object is None,
                    'share_type': share_object.share_type if share_object else 'owner'
                },
                'debug_info': {
                    'total_plots': all_plots.count(),
                    'matching_plots': saved_plots.count(),
                    'filter_criteria': serializable_filter_criteria
                }
            }, status=200)

        except Exception as e:
            print(f"Error in FetchPlots: {str(e)}")  # Debug log
            return Response({
                'error': f'Error fetching plots: {str(e)}'
            }, status=500)

class ProjectDetails(APIView):
    def post(self, request):
        try:
            # Extract project_id and user_id from request
            project_id = request.data.get('project_id')
            user_id = request.data.get('user_id')
            
            if not project_id:
                return Response({"error": "project_id is required"}, status=400)
            
            if not user_id:
                return Response({"error": "user_id is required"}, status=400)

            # Check if user has access to this project
            has_access, share_object, permission_level = check_project_access(user_id, project_id)
            
            if not has_access:
                return Response({"error": "Access denied. You don't have permission to view this project."}, status=403)

            # Get project
            try:
                project = Projects.objects.get(id=project_id)
            except Projects.DoesNotExist:
                return Response({"error": "Project not found"}, status=404)

            # Base project folder path
            project_folder = os.path.join(settings.MEDIA_ROOT, f"user_{project.user.id}/project_{project.id}")
            
            # Initialize response data structure
            project_details = {
                "project_id": project.id,
                "project_name": project.name,
                "access_info": {
                    "permission_level": permission_level,
                    "is_owner": share_object is None,  # True if user is project owner
                    "shared_by": None
                },
                "files": {
                    "media": [],
                    "kpi": [],
                    "concatenated": []
                }
            }
            
            # Add sharing info if this is a shared project
            if share_object:
                project_details["access_info"]["shared_by"] = {
                    "id": share_object.shared_by.id,
                    "username": share_object.shared_by.username,
                    "email": share_object.shared_by.email
                }
                
                # For file-specific sharing, add shared file info
                if share_object.share_type == 'file':
                    project_details["shared_file"] = {
                        "file_type": share_object.file_type,
                        "file_name": share_object.file_name,
                        "sheet_name": share_object.sheet_name
                    }

            # Process files based on access level
            if share_object and share_object.share_type == 'file':
                # File-specific sharing - only show the shared file
                file_type = share_object.file_type
                file_name = share_object.file_name
                
                file_folder = os.path.join(project_folder, file_type)
                file_path = os.path.join(file_folder, file_name)
                
                if os.path.exists(file_path) and os.path.isdir(file_path):
                    file_info = {
                        "name": file_name,
                        "sheets": []
                    }
                    
                    # Get all sheets (CSV files) in the file directory
                    for sheet_file in os.listdir(file_path):
                        if sheet_file.endswith('.csv'):
                            # For file-specific sharing, only show the shared sheet if specified
                            if share_object.sheet_name and sheet_file != share_object.sheet_name:
                                continue
                                
                            sheet_path = os.path.join(file_path, sheet_file)
                            sheet_info = {
                                "name": sheet_file,
                                "size": os.path.getsize(sheet_path),
                                "last_modified": os.path.getmtime(sheet_path)
                            }
                            
                            # Read complete sheet data
                            try:
                                df = pd.read_csv(sheet_path, dtype=str)
                                df = df.replace([np.nan, np.inf, -np.inf], None)
                                sheet_info["columns"] = df.columns.tolist()
                                sheet_info["data"] = make_json_safe(df.values.tolist())
                            except Exception as e:
                                print(f"Error reading sheet {sheet_file}: {str(e)}")
                                sheet_info["error"] = "Could not read sheet data"
                            
                            file_info["sheets"].append(sheet_info)
                    
                    project_details["files"][file_type].append(file_info)
            else:
                # Project-level sharing or owner access - show all files
                for file_type in ["media", "kpi"]:
                    file_folder = os.path.join(project_folder, file_type)
                    if not os.path.exists(file_folder):
                        continue

                    # Get all files in the folder
                    for file_name in os.listdir(file_folder):
                        file_path = os.path.join(file_folder, file_name)
                        if os.path.isdir(file_path):  # Each file is a directory containing sheets
                            file_info = {
                                "name": file_name,
                                "sheets": []
                            }
                            
                            # Get all sheets (CSV files) in the file directory
                            for sheet_file in os.listdir(file_path):
                                if sheet_file.endswith('.csv'):
                                    sheet_path = os.path.join(file_path, sheet_file)
                                    sheet_info = {
                                        "name": sheet_file,
                                        "size": os.path.getsize(sheet_path),
                                        "last_modified": os.path.getmtime(sheet_path)
                                    }
                                    
                                    # Read complete sheet data
                                    try:
                                        df = pd.read_csv(sheet_path, dtype=str)
                                        df = df.replace([np.nan, np.inf, -np.inf], None)
                                        sheet_info["columns"] = df.columns.tolist()
                                        sheet_info["data"] = make_json_safe(df.values.tolist())
                                    except Exception as e:
                                        print(f"Error reading sheet {sheet_file}: {str(e)}")
                                        sheet_info["error"] = "Could not read sheet data"
                                    
                                    file_info["sheets"].append(sheet_info)
                            
                            project_details["files"][file_type].append(file_info)

            return Response(project_details, status=200)

        except Exception as e:
            return Response({"error": str(e)}, status=500)


class SaveReportPivot(APIView):
    parser_classes = [JSONParser]

    def post(self, request):
        try:
            # Extract data from request
            user_id = request.data.get('user_id')
            project_id = request.data.get('project_id')
            pivot_name = request.data.get('pivot_name')
            file_type = request.data.get('file_type')
            file_name = request.data.get('file_name')
            sheet_name = request.data.get('sheet_name')
            pivot_config = request.data.get('pivot_config')
            pivot_data = request.data.get('pivot_data')

            # Validate required fields
            if not all([user_id, project_id, pivot_name, file_type, file_name, sheet_name, pivot_config, pivot_data]):
                missing_fields = []
                if not user_id: missing_fields.append('user_id')
                if not project_id: missing_fields.append('project_id')
                if not pivot_name: missing_fields.append('pivot_name')
                if not file_type: missing_fields.append('file_type')
                if not file_name: missing_fields.append('file_name')
                if not sheet_name: missing_fields.append('sheet_name')
                if not pivot_config: missing_fields.append('pivot_config')
                if not pivot_data: missing_fields.append('pivot_data')
                return Response({
                    'error': f'Missing required fields: {", ".join(missing_fields)}'
                }, status=400)

            # Validate file_type
            if file_type not in ['kpi', 'media']:
                return Response({
                    'error': 'Invalid file_type. Must be either "kpi" or "media"'
                }, status=400)

            # Check if user has access to this project/file
            has_access, share_object, permission_level = check_project_access(
                user_id, project_id, file_type, file_name, sheet_name
            )
            
            if not has_access:
                return Response({
                    'error': 'Access denied. You don\'t have permission to save pivots for this project/file.'
                }, status=403)
            
            # Check if user has edit permissions
            if permission_level == 'view':
                return Response({
                    'error': 'Access denied. You only have view permissions. Edit permissions required to save pivots.'
                }, status=403)

            # Get user and project
            try:
                user = User.objects.get(id=user_id)
                project = Projects.objects.get(id=project_id)
            except User.DoesNotExist:
                return Response({'error': 'User not found'}, status=404)
            except Projects.DoesNotExist:
                return Response({'error': 'Project not found'}, status=404)

            # Create or update saved pivot
            saved_pivot, created = SavedPivot.objects.update_or_create(
                user=user,
                project=project,
                pivot_name=pivot_name,
                defaults={
                    'file_type': file_type,
                    'file_name': file_name,
                    'sheet_name': sheet_name,
                    'pivot_config': pivot_config,
                    'pivot_data': pivot_data
                }
            )

            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "save_pivot", details=f"Pivot table saved successfully", ip_address=ip)

            return Response({
                'message': 'Pivot table saved successfully',
                'pivot_id': saved_pivot.id,
                'created': created,
                'updated_at': saved_pivot.updated_at,
                'access_info': {
                    'permission_level': permission_level,
                    'is_owner': share_object is None
                }
            }, status=200)

        except Exception as e:
            return Response({'error': str(e)}, status=500)


class FetchReportPivot(APIView):
    parser_classes = [JSONParser]

    def post(self, request):
        try:
            # Extract data from request
            user_id = request.data.get('user_id')
            project_id = request.data.get('project_id')
            file_type = request.data.get('file_type')  # Optional filter
            file_name = request.data.get('file_name')  # Optional filter
            sheet_name = request.data.get('sheet_name')  # Optional filter

            # Validate required fields
            if not all([user_id, project_id]):
                missing_fields = []
                if not user_id: missing_fields.append('user_id')
                if not project_id: missing_fields.append('project_id')
                return Response({
                    'error': f'Missing required fields: {", ".join(missing_fields)}'
                }, status=400)

            # Check if user has access to this project
            has_access, share_object, permission_level = check_project_access(user_id, project_id)
            
            if not has_access:
                return Response({
                    'error': 'Access denied. You don\'t have permission to view pivots for this project.'
                }, status=403)

            # Get user and project
            try:
                user = User.objects.get(id=user_id)
                project = Projects.objects.get(id=project_id)
            except User.DoesNotExist:
                return Response({'error': 'User not found'}, status=404)
            except Projects.DoesNotExist:
                return Response({'error': 'Project not found'}, status=404)

            # Build filter for pivots based on access level
            pivot_filter = {
                'project': project
            }
            
            # If this is a file-specific share, only show pivots for that file
            if share_object and share_object.share_type == 'file':
                pivot_filter.update({
                    'file_type': share_object.file_type,
                    'file_name': share_object.file_name
                })
                
                # If specific sheet is shared, only show pivots for that sheet
                if share_object.sheet_name:
                    pivot_filter['sheet_name'] = share_object.sheet_name
            else:
                # Project-level sharing or owner access - apply optional filters
                if file_type:
                    pivot_filter['file_type'] = file_type
                if file_name:
                    pivot_filter['file_name'] = file_name
                if sheet_name:
                    pivot_filter['sheet_name'] = sheet_name

            # Fetch all saved pivots for this project with filters
            saved_pivots = SavedPivot.objects.filter(**pivot_filter).order_by('-updated_at')

            # Format response
            pivots_data = [{
                'id': pivot.id,
                'pivot_name': pivot.pivot_name,
                'file_type': pivot.file_type,
                'file_name': pivot.file_name,
                'sheet_name': pivot.sheet_name,
                'pivot_config': pivot.pivot_config,
                'pivot_data': pivot.pivot_data,
                'created_at': pivot.created_at,
                'updated_at': pivot.updated_at
            } for pivot in saved_pivots]

           

            return Response({
                'pivots': pivots_data,
                'total_pivots': saved_pivots.count(),
                'access_info': {
                    'permission_level': permission_level,
                    'is_owner': share_object is None,
                    'share_type': share_object.share_type if share_object else 'owner'
                }
            }, status=200)

        except Exception as e:
            return Response({'error': str(e)}, status=500)

class DeleteReportPivot(APIView):
    parser_classes = [JSONParser]

    def post(self, request):
        try:
            # Extract data from request
            user_id = request.data.get('user_id')
            project_id = request.data.get('project_id')
            pivot_id = request.data.get('pivot_id')

            # Validate required fields
            if not all([user_id, project_id, pivot_id]):
                missing_fields = []
                if not user_id: missing_fields.append('user_id')
                if not project_id: missing_fields.append('project_id')
                if not pivot_id: missing_fields.append('pivot_id')
                return Response({
                    'error': f'Missing required fields: {", ".join(missing_fields)}'
                }, status=400)

            # Get user and project
            try:
                user = User.objects.get(id=user_id)
                project = Projects.objects.get(id=project_id, user=user)
            except User.DoesNotExist:
                return Response({'error': 'User not found'}, status=404)
            except Projects.DoesNotExist:
                return Response({'error': 'Project not found'}, status=404)

            # Try to get and delete the pivot
            try:
                pivot = SavedPivot.objects.get(
                    id=pivot_id,
                    user=user,
                    project=project
                )
                pivot.delete()
                ip = request.META.get('REMOTE_ADDR')
                log_user_action(user, "delete_pivot", details=f"Pivot table deleted successfully", ip_address=ip)
                return Response({
                    'message': 'Pivot table deleted successfully',
                    'pivot_id': pivot_id
                }, status=200)
            except SavedPivot.DoesNotExist:
                return Response({
                    'error': 'Pivot table not found or you do not have permission to delete it'
                }, status=404)

        except Exception as e:
            return Response({'error': str(e)}, status=500)

class SavePivotPlot(APIView):
    parser_classes = [JSONParser]

    def post(self, request):
        try:
            # Extract data from request
            user_id = request.data.get('user_id')
            project_id = request.data.get('project_id')
            pivot_id = request.data.get('pivot_id')
            plot_name = request.data.get('plot_name', 'Default Pivot Plot')
            plot_config = request.data.get('plot_config')
            chart_data = request.data.get('chart_data')
            chart_options = request.data.get('chart_options', {})

            # Validate required fields
            if not all([user_id, project_id, pivot_id, plot_config, chart_data]):
                missing_fields = []
                if not user_id: missing_fields.append('user_id')
                if not project_id: missing_fields.append('project_id')
                if not pivot_id: missing_fields.append('pivot_id')
                if not plot_config: missing_fields.append('plot_config')
                if not chart_data: missing_fields.append('chart_data')
                return Response({
                    'error': f'Missing required fields: {", ".join(missing_fields)}'
                }, status=400)

            # Get user, project, and pivot
            try:
                user = User.objects.get(id=user_id)
            except User.DoesNotExist:
                return Response({'error': 'User not found'}, status=404)

            try:
                project = Projects.objects.get(id=project_id, user=user)
            except Projects.DoesNotExist:
                return Response({'error': 'Project not found'}, status=404)

            try:
                pivot = SavedPivot.objects.get(id=pivot_id, project=project, user=user)
            except SavedPivot.DoesNotExist:
                return Response({'error': 'Pivot table not found'}, status=404)

            # Extract activeFilters from plot_config if it exists
            active_filters = plot_config.get('activeFilters', {}) if plot_config else {}

            # Create or update the plot
            saved_plot, created = SavedPivotPlot.objects.update_or_create(
                user=user,
                project=project,
                pivot=pivot,
                plot_name=plot_name,
                defaults={
                    'plot_config': plot_config,
                    'chart_data': chart_data,
                    'chart_options': chart_options,
                    'active_filters': active_filters
                }
            )

            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "save_pivot_plot", details=f"Plot saved successfully", ip_address=ip)

            return Response({
                'message': 'Plot saved successfully',
                'plot_id': saved_plot.id,
                'plot_name': saved_plot.plot_name,
                'created_at': saved_plot.created_at,
                'updated_at': saved_plot.updated_at,
                'is_new': created
            }, status=201 if created else 200)

        except Exception as e:
            return Response({
                'error': f'Error saving plot: {str(e)}'
            }, status=500)

class DeletePivotPlot(APIView):
    parser_classes = [JSONParser]

    def post(self, request):
        try:
            # Extract data from request
            user_id = request.data.get('user_id')
            project_id = request.data.get('project_id')
            pivot_id = request.data.get('pivot_id')
            plot_id = request.data.get('plot_id')  # Optional: if provided, delete specific plot

            # Validate required fields
            if not all([user_id, project_id, pivot_id]):
                missing_fields = []
                if not user_id: missing_fields.append('user_id')
                if not project_id: missing_fields.append('project_id')
                if not pivot_id: missing_fields.append('pivot_id')
                return Response({
                    'error': f'Missing required fields: {", ".join(missing_fields)}'
                }, status=400)

            # Get user and project
            try:
                user = User.objects.get(id=user_id)
            except User.DoesNotExist:
                return Response({'error': 'User not found'}, status=404)

            try:
                project = Projects.objects.get(id=project_id, user=user)
            except Projects.DoesNotExist:
                return Response({'error': 'Project not found'}, status=404)

            try:
                pivot = SavedPivot.objects.get(id=pivot_id, project=project, user=user)
            except SavedPivot.DoesNotExist:
                return Response({'error': 'Pivot table not found'}, status=404)

            # Delete specific plot if plot_id is provided, otherwise delete all plots for the pivot
            if plot_id:
                try:
                    plot = SavedPivotPlot.objects.get(
                        id=plot_id,
                        user=user,
                        project=project,
                        pivot=pivot
                    )
                    plot.delete()
                    return Response({
                        'message': f'Plot {plot_id} deleted successfully'
                    }, status=200)
                except SavedPivotPlot.DoesNotExist:
                    return Response({
                        'error': 'Plot not found or you don\'t have permission to delete it'
                    }, status=404)
            else:
                # Delete all plots for this pivot
                deleted_count, _ = SavedPivotPlot.objects.filter(
                    user=user,
                    project=project,
                    pivot=pivot
                ).delete()
                
                return Response({
                    'message': f'Successfully deleted {deleted_count} plot(s)'
                }, status=200)

        except Exception as e:
            return Response({
                'error': f'Error deleting plot(s): {str(e)}'
            }, status=500)

class FetchPivotPlots(APIView):
    parser_classes = [JSONParser]

    def post(self, request):
        try:
            # Extract data from request
            user_id = request.data.get('user_id')
            project_id = request.data.get('project_id')
            pivot_id = request.data.get('pivot_id')

            # Validate required fields
            if not all([user_id, project_id, pivot_id]):
                missing_fields = []
                if not user_id: missing_fields.append('user_id')
                if not project_id: missing_fields.append('project_id')
                if not pivot_id: missing_fields.append('pivot_id')
                return Response({
                    'error': f'Missing required fields: {", ".join(missing_fields)}'
                }, status=400)

            # Get user
            try:
                user = User.objects.get(id=user_id)
            except User.DoesNotExist:
                return Response({'error': 'User not found'}, status=404)

            # Get project and check access
            try:
                project = Projects.objects.get(id=project_id)
            except Projects.DoesNotExist:
                return Response({'error': 'Project not found'}, status=404)

            # Check if user has access to this project
            has_access, share_object, permission_level = check_project_access(user_id, project_id)
            
            if not has_access:
                return Response({
                    'error': 'Access denied. You don\'t have permission to access this project.'
                }, status=403)

            # Get the pivot table
            try:
                # For project owner, get pivot by user and project
                if project.user.id == int(user_id):
                    pivot = SavedPivot.objects.get(id=pivot_id, project=project, user=user)
                else:
                    # For shared access, get pivot by project only (pivots belong to project owner)
                    pivot = SavedPivot.objects.get(id=pivot_id, project=project)
            except SavedPivot.DoesNotExist:
                return Response({'error': 'Pivot table not found'}, status=404)

            # Fetch all plots for this pivot
            # For project owner, get plots by user and project
            if project.user.id == int(user_id):
                plots = SavedPivotPlot.objects.filter(
                    user=user,
                    project=project,
                    pivot=pivot
                ).order_by('-updated_at')
            else:
                # For shared access, get plots by project only (plots belong to project owner)
                plots = SavedPivotPlot.objects.filter(
                    project=project,
                    pivot=pivot
                ).order_by('-updated_at')

            # Prepare the response data
            plots_data = []
            for plot in plots:
                plot_data = {
                    'plot_id': plot.id,
                    'plot_name': plot.plot_name,
                    'plot_config': plot.plot_config,
                    'chart_data': plot.chart_data,
                    'chart_options': plot.chart_options,
                    'active_filters': plot.active_filters,
                    'created_at': plot.created_at,
                    'updated_at': plot.updated_at,
                    'pivot_info': {
                        'pivot_id': pivot.id,
                        'pivot_name': pivot.pivot_name,
                        'file_type': pivot.file_type,
                        'file_name': pivot.file_name,
                        'sheet_name': pivot.sheet_name
                    }
                }
                plots_data.append(plot_data)

            return Response({
                'message': 'Plots fetched successfully',
                'plots': plots_data,
                'total_plots': len(plots_data),
                'permission_level': permission_level
            }, status=200)

        except Exception as e:
            return Response({
                'error': f'Error fetching plots: {str(e)}'
            }, status=500)


class ConcatenateProjectSheets(APIView):
    """
    API endpoint to concatenate specific sheets from different files within a project.
    """
    def post(self, request):
        project_id = request.data.get('project_id')
        sheet_selections = request.data.get('sheet_selections', {})
        user_id = request.data.get('user_id')
        update_database = request.data.get('update_database', False)  # Optional flag to update DB

        if not all([project_id, sheet_selections, user_id]):
            return Response({
                'success': False,
                'message': 'Missing required parameters: project_id, sheet_selections, user_id'
            }, status=400)

        try:
            project = Projects.objects.get(id=project_id, user_id=user_id)
        except Projects.DoesNotExist:
            return Response({
                'success': False,
                'message': 'Project not found or access denied'
            }, status=404)

        project_folder = os.path.join(settings.MEDIA_ROOT, f"user_{user_id}/project_{project_id}")
        
        if not os.path.exists(project_folder):
            return Response({
                'success': False,
                'message': 'Project folder not found'
            }, status=404)

        try:
            # Get all available files in the project
            all_files = {}
            
            # Get KPI files
            for kpi_file in project.kpi_file:
                kpi_path = os.path.join(project_folder, 'kpi', kpi_file)
                if os.path.exists(kpi_path):
                    all_files[kpi_file] = ('kpi', kpi_path)
            
            # Get media files
            for media_file in project.media_file:
                media_path = os.path.join(project_folder, 'media', media_file)
                if os.path.exists(media_path):
                    all_files[media_file] = ('media', media_path)

            # Scan file system for files not in database
            kpi_folder = os.path.join(project_folder, 'kpi')
            media_folder = os.path.join(project_folder, 'media')
            
            # Check KPI folder
            if os.path.exists(kpi_folder):
                for file_name in os.listdir(kpi_folder):
                    if file_name not in all_files:
                        file_path = os.path.join(kpi_folder, file_name)
                        if os.path.isfile(file_path) or os.path.isdir(file_path):
                            all_files[file_name] = ('kpi', file_path)
            
            # Check media folder
            if os.path.exists(media_folder):
                for file_name in os.listdir(media_folder):
                    if file_name not in all_files:
                        file_path = os.path.join(media_folder, file_name)
                        if os.path.isfile(file_path) or os.path.isdir(file_path):
                            all_files[file_name] = ('media', file_path)

            # Update database if requested
            if update_database:
                kpi_files_to_add = []
                media_files_to_add = []
                
                for file_name, (file_type, file_path) in all_files.items():
                    if file_type == 'kpi' and file_name not in project.kpi_file:
                        kpi_files_to_add.append(file_name)
                    elif file_type == 'media' and file_name not in project.media_file:
                        media_files_to_add.append(file_name)
                
                if kpi_files_to_add:
                    project.kpi_file.extend(kpi_files_to_add)
                if media_files_to_add:
                    project.media_file.extend(media_files_to_add)
                
                if kpi_files_to_add or media_files_to_add:
                    project.save()

            # Validate sheet selections
            valid_sheets = {}
            debug_info = {
                'available_files': list(all_files.keys()),
                'requested_files': list(sheet_selections.keys()),
                'file_details': {}
            }
            
            for file_id, sheet_names in sheet_selections.items():
                if file_id not in all_files:
                    debug_info['file_details'][file_id] = 'File not found in project'
                    continue
                
                file_type, file_path = all_files[file_id]
                available_sheets = []
                
                # Check if it's a CSV file
                if os.path.isfile(file_path) and file_path.endswith('.csv'):
                    available_sheets = [f"{file_id}_data"]
                # Check if it's a directory with CSV files (Excel converted)
                elif os.path.isdir(file_path):
                    csv_files = [f for f in os.listdir(file_path) if f.endswith('.csv')]
                    available_sheets = [f"{file_id}_{os.path.splitext(csv_file)[0]}" for csv_file in csv_files]
                
                debug_info['file_details'][file_id] = {
                    'file_type': file_type,
                    'file_path': file_path,
                    'available_sheets': available_sheets,
                    'requested_sheets': sheet_names
                }
                
                # Filter requested sheets that actually exist
                valid_sheet_names = [sheet for sheet in sheet_names if sheet in available_sheets]
                if valid_sheet_names:
                    valid_sheets[file_id] = (file_type, file_path, valid_sheet_names)

            if not valid_sheets:
                return Response({
                    'success': False,
                    'message': 'No valid sheets found for concatenation',
                    'debug_info': debug_info
                }, status=400)

            # Read and concatenate sheets
            concatenated_dfs = []
            all_columns = set()
            sheet_names_list = []

            for file_id, (file_type, file_path, sheet_names) in valid_sheets.items():
                for sheet_name in sheet_names:
                    try:
                        # Extract actual sheet name from the combined name
                        if sheet_name.endswith('_data'):
                            # It's a CSV file
                            df = pd.read_csv(file_path, dtype=str)
                            actual_sheet_name = file_id  # Use file name as sheet name
                        else:
                            # It's a sheet from Excel file
                            actual_sheet_name = sheet_name.replace(f"{file_id}_", "")
                            csv_path = os.path.join(file_path, f"{actual_sheet_name}.csv")
                            if os.path.exists(csv_path):
                                df = pd.read_csv(csv_path, dtype=str)
                            else:
                                continue
                        
                        # Add sheet name to list for naming
                        sheet_names_list.append(actual_sheet_name)
                        
                        # Clean the dataframe to handle inf, -inf, and NaN values
                        df = df.replace([np.inf, -np.inf], 'infinity')
                        df = df.fillna('')
                        
                        # Add source information
                        df['source_file'] = file_id
                        df['source_sheet'] = sheet_name
                        
                        concatenated_dfs.append(df)
                        all_columns.update(df.columns.tolist())
                        
                    except Exception as e:
                        print(f"Error reading sheet {sheet_name} from {file_id}: {e}")
                        continue

            if not concatenated_dfs:
                return Response({
                    'success': False,
                    'message': 'No valid data found in selected sheets'
                }, status=400)

            # Align columns across all dataframes
            aligned_dfs = []
            for df in concatenated_dfs:
                # Add missing columns with NaN values
                for col in all_columns:
                    if col not in df.columns:
                        df[col] = ''
                # Reorder columns to match
                df = df[list(all_columns)]
                aligned_dfs.append(df)

            # Concatenate all dataframes
            final_df = pd.concat(aligned_dfs, ignore_index=True)

            # Create a shorter, safer concatenated file name
            import hashlib
            import time
            
            # Create a shorter name from the first few characters of each sheet
            short_names = []
            for sheet_name in sheet_names_list:
                # Take first 10 characters, remove special characters
                short_name = ''.join(c for c in sheet_name[:10] if c.isalnum() or c in ' -_')
                short_names.append(short_name)
            
            # Create concatenated name with timestamp to ensure uniqueness
            timestamp = int(time.time())
            concatenated_name = f"concatenated_{timestamp}"
            
            # Create folder for the concatenated file
            concatenated_folder = os.path.join(project_folder, 'concatenated', concatenated_name)
            os.makedirs(concatenated_folder, exist_ok=True)

            # Save the CSV file with a descriptive name
            csv_filename = f"{'+'.join(short_names)}.csv"
            # Ensure the filename is not too long
            if len(csv_filename) > 200:
                csv_filename = f"concatenated_{timestamp}.csv"
            
            csv_path = os.path.join(concatenated_folder, csv_filename)
            final_df.to_csv(csv_path, index=False, encoding='utf-8')

            # Prepare sheet data for response
            sheets_data = {
                csv_filename: {
                    "columns": final_df.columns.tolist(),
                    "data": make_json_safe(final_df.values.tolist()[:100])  # Limit to first 100 rows and make JSON safe
                }
            }

            # Commit to git
            commit_msg = f"concatenated - {user_id}/{project_id}/concatenated/{concatenated_name}/{csv_filename}"
            subprocess.run(["git", "add", csv_path], cwd=project_folder)
            subprocess.run(["git", "commit", "-m", commit_msg], cwd=project_folder)

            # Update project with concatenated file
            if not hasattr(project, 'concatenated_file'):
                project.concatenated_file = []
            if concatenated_name not in project.concatenated_file:
                project.concatenated_file.append(concatenated_name)
            project.save()

            user = None
            if hasattr(request, 'user') and request.user.is_authenticated:
                user = request.user
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "concatenate_sheets", details=f"Sheets concatenated successfully", ip_address=ip)

            return Response({
                'success': True,
                'message': 'Sheets concatenated successfully',
                'concatenated_file': {
                    'name': concatenated_name,
                    'type': 'concatenated',
                    'sheets_data': sheets_data
                },
                'debug_info': debug_info
            }, status=200)

        except Exception as e:
            return Response({
                'success': False,
                'message': f'Error during sheet concatenation: {str(e)}'
            }, status=500)


class GetProjectFiles(APIView):
    """
    API endpoint to get all available files in a project for debugging.
    """
    def post(self, request):
        project_id = request.data.get('project_id')
        user_id = request.data.get('user_id')
        update_database = request.data.get('update_database', False)  # Optional flag to update DB

        if not all([project_id, user_id]):
            return Response({
                'success': False,
                'message': 'Missing required parameters: project_id, user_id'
            }, status=400)

        try:
            project = Projects.objects.get(id=project_id, user_id=user_id)
        except Projects.DoesNotExist:
            return Response({
                'success': False,
                'message': 'Project not found or access denied'
            }, status=404)

        project_folder = os.path.join(settings.MEDIA_ROOT, f"user_{user_id}/project_{project_id}")
        
        if not os.path.exists(project_folder):
            return Response({
                'success': False,
                'message': 'Project folder not found'
            }, status=404)

        try:
            # Get all files from the project folder
            all_files = []
            database_files = []
            filesystem_files = []
            missing_files = []

            # Check KPI files
            kpi_folder = os.path.join(project_folder, 'kpi')
            if os.path.exists(kpi_folder):
                for item in os.listdir(kpi_folder):
                    item_path = os.path.join(kpi_folder, item)
                    if os.path.isdir(item_path):
                        # This is a file folder, look for Excel files inside
                        for file in os.listdir(item_path):
                            if file.endswith(('.xlsx', '.xls')):
                                file_path = os.path.join(item, file)
                                all_files.append({
                                    'type': 'kpi',
                                    'name': file_path,
                                    'path': f"user_{user_id}/project_{project_id}/kpi/{file_path}",
                                    'exists_in_db': file_path in (project.kpi_file or []),
                                    'exists_in_fs': True
                                })
                                filesystem_files.append(file_path)
                                if file_path not in (project.kpi_file or []):
                                    missing_files.append(('kpi', file_path))

            # Check Media files
            media_folder = os.path.join(project_folder, 'media')
            if os.path.exists(media_folder):
                for item in os.listdir(media_folder):
                    item_path = os.path.join(media_folder, item)
                    if os.path.isdir(item_path):
                        # This is a file folder, look for Excel files inside
                        for file in os.listdir(item_path):
                            if file.endswith(('.xlsx', '.xls')):
                                file_path = os.path.join(item, file)
                                all_files.append({
                                    'type': 'media',
                                    'name': file_path,
                                    'path': f"user_{user_id}/project_{project_id}/media/{file_path}",
                                    'exists_in_db': file_path in (project.media_file or []),
                                    'exists_in_fs': True
                                })
                                filesystem_files.append(file_path)
                                if file_path not in (project.media_file or []):
                                    missing_files.append(('media', file_path))

            # Add database files that might not exist in filesystem
            for file_path in (project.kpi_file or []):
                database_files.append(file_path)
                if not any(f['name'] == file_path for f in all_files):
                    all_files.append({
                        'type': 'kpi',
                        'name': file_path,
                        'path': f"user_{user_id}/project_{project_id}/kpi/{file_path}",
                        'exists_in_db': True,
                        'exists_in_fs': False
                    })

            for file_path in (project.media_file or []):
                database_files.append(file_path)
                if not any(f['name'] == file_path for f in all_files):
                    all_files.append({
                        'type': 'media',
                        'name': file_path,
                        'path': f"user_{user_id}/project_{project_id}/media/{file_path}",
                        'exists_in_db': True,
                        'exists_in_fs': False
                    })

            # Update database if requested and there are missing files
            if update_database and missing_files:
                for file_type, file_name in missing_files:
                    if file_type == 'kpi':
                        if not hasattr(project, 'kpi_file') or project.kpi_file is None:
                            project.kpi_file = []
                        project.kpi_file.append(file_name)
                    elif file_type == 'media':
                        if not hasattr(project, 'media_file') or project.media_file is None:
                            project.media_file = []
                        project.media_file.append(file_name)
                
                project.save()

            return Response({
                'success': True,
                'project_id': project_id,
                'user_id': user_id,
                'files': all_files,
                'total_files': len(all_files),
                'database_files': len(database_files),
                'filesystem_files': len(filesystem_files),
                'missing_files': [f[1] for f in missing_files],
                'database_updated': update_database and bool(missing_files)
            }, status=200)

        except Exception as e:
            return Response({
                'success': False,
                'message': f'Error getting project files: {str(e)}'
            }, status=500)


# Project Sharing Views
class ShareProject(APIView):
    """
    API endpoint to share a project or specific files with other users.
    """
    def post(self, request):
        try:
            # Extract data from request
            project_id = request.data.get('project_id')
            shared_by_user_id = request.data.get('shared_by_user_id')
            shared_with_email = request.data.get('shared_with_email')
            share_type = request.data.get('share_type', 'project')  # 'project' or 'file'
            permission_level = request.data.get('permission_level', 'view')  # 'view', 'edit', 'admin'
            
            # For file-specific sharing
            file_type = request.data.get('file_type')  # 'kpi' or 'media'
            file_name = request.data.get('file_name')
            sheet_name = request.data.get('sheet_name')
            
            # Validate required fields
            if not all([project_id, shared_by_user_id, shared_with_email]):
                return Response({
                    'success': False,
                    'message': 'project_id, shared_by_user_id, and shared_with_email are required'
                }, status=400)
            
            # Validate share_type
            if share_type not in ['project', 'file']:
                return Response({
                    'success': False,
                    'message': 'share_type must be either "project" or "file"'
                }, status=400)
            
            # Validate permission_level
            if permission_level not in ['view', 'edit', 'admin']:
                return Response({
                    'success': False,
                    'message': 'permission_level must be either "view", "edit", or "admin"'
                }, status=400)
            
            # For file-specific sharing, validate file fields
            if share_type == 'file':
                if not all([file_type, file_name]):
                    return Response({
                        'success': False,
                        'message': 'file_type and file_name are required for file-specific sharing'
                    }, status=400)
                
                if file_type not in ['kpi', 'media', 'concatenated']:
                    return Response({
                        'success': False,
                        'message': 'file_type must be either "kpi", "media", or "concatenated"'
                    }, status=400)
            
            # Get users
            try:
                shared_by_user = User.objects.get(id=shared_by_user_id)
            except User.DoesNotExist:
                return Response({
                    'success': False,
                    'message': 'Shared by user not found'
                }, status=404)
            
            try:
                shared_with_user = User.objects.get(email=shared_with_email)
            except User.DoesNotExist:
                return Response({
                    'success': False,
                    'message': 'User with provided email not found'
                }, status=404)
            
            # Check if user is trying to share with themselves
            if shared_by_user.id == shared_with_user.id:
                return Response({
                    'success': False,
                    'message': 'Cannot share project with yourself'
                }, status=400)
            
            # Get project
            try:
                project = Projects.objects.get(id=project_id, user=shared_by_user)
            except Projects.DoesNotExist:
                return Response({
                    'success': False,
                    'message': 'Project not found or you do not have permission to share it'
                }, status=404)
            
            # For file-specific sharing, validate file exists
            if share_type == 'file':
                project_folder = os.path.join(settings.MEDIA_ROOT, f"user_{project.user.id}/project_{project.id}")
                
                if file_type == 'concatenated':
                    # For concatenated files, check if any folder contains this CSV file
                    concatenated_folder = os.path.join(project_folder, "concatenated")
                    if os.path.exists(concatenated_folder):
                        file_found = False
                        for folder_name in os.listdir(concatenated_folder):
                            folder_path = os.path.join(concatenated_folder, folder_name)
                            if os.path.isdir(folder_path):
                                # Check if this folder contains the CSV file
                                for csv_file in os.listdir(folder_path):
                                    if csv_file == file_name or csv_file.endswith('.csv'):
                                        file_found = True
                                        # Update file_name to the folder name for storage
                                        file_name = folder_name
                                        break
                            if file_found:
                                break
                        
                        if not file_found:
                            return Response({
                                'success': False,
                                'message': f'Concatenated file {file_name} not found in project'
                            }, status=404)
                    else:
                        return Response({
                            'success': False,
                            'message': 'Concatenated folder not found in project'
                        }, status=404)
                else:
                    # For kpi and media files, check the regular path
                    file_path = os.path.join(project_folder, file_type, file_name)
                    
                    if not os.path.exists(file_path):
                        return Response({
                            'success': False,
                            'message': f'File {file_name} not found in {file_type} folder'
                        }, status=404)
                    
                    # If sheet_name is provided, validate it exists
                    if sheet_name:
                        sheet_path = os.path.join(file_path, sheet_name)
                        if not os.path.exists(sheet_path):
                            return Response({
                                'success': False,
                                'message': f'Sheet {sheet_name} not found in file {file_name}'
                            }, status=404)
            
            # Check if share already exists
            share_filter = {
                'project': project,
                'shared_with': shared_with_user,
                'is_active': True
            }
            user = get_logging_user(request, getattr(project, 'user', None))
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "share_project", details=f"Project shared successfully", ip_address=ip)
            
            if share_type == 'file':
                share_filter.update({
                    'share_type': 'file',
                    'file_type': file_type,
                    'file_name': file_name,
                    'sheet_name': sheet_name
                })
            else:
                share_filter['share_type'] = 'project'
            
            existing_share = ProjectShare.objects.filter(**share_filter).first()
            
            if existing_share:
                # Update existing share
                existing_share.permission_level = permission_level
                existing_share.updated_at = timezone.now()
                existing_share.save()
                # Send email notification for update
                try:
                    subject = f"Project Share Updated: {project.name}"
                    if share_type == 'file':
                        message = f"Hi {shared_with_user.username},\n\nYou have been granted updated {permission_level} access to the file '{file_name}' in the project '{project.name}' by {shared_by_user.username}.\n\nLogin to your account to access the shared file."
                    else:
                        message = f"Hi {shared_with_user.username},\n\nYou have been granted updated {permission_level} access to the project '{project.name}' by {shared_by_user.username}.\n\nLogin to your account to access the shared project."
                    send_mail(subject, message, settings.DEFAULT_FROM_EMAIL, [shared_with_user.email], fail_silently=True)
                except Exception as e:
                    print("email not sent", e)
                    pass  # Do not block API on email failure
                return Response({
                    'success': True,
                    'message': 'Project share updated successfully',
                    'share_id': existing_share.id
                }, status=200)
            else:
                # Create new share
                share = ProjectShare.objects.create(
                    project=project,
                    shared_by=shared_by_user,
                    shared_with=shared_with_user,
                    share_type=share_type,
                    permission_level=permission_level,
                    file_type=file_type if share_type == 'file' else None,
                    file_name=file_name if share_type == 'file' else None,
                    sheet_name=sheet_name if share_type == 'file' else None
                )
                # Send email notification for new share
                try:
                    subject = f"Project Shared: {project.name} on Sanitify"
                    if share_type == 'file':
                        message = f"Hi {shared_with_user.username},\n\nYou have been granted {permission_level} access to the file '{file_name}' in the project '{project.name}' by {shared_by_user.username}.\n\nLogin to your account to access the shared file."
                        print("message", message)
                    else:
                        message = f"Hi {shared_with_user.username},\n\nYou have been granted {permission_level} access to the project '{project.name}' by {shared_by_user.username}.\n\nLogin to your account to access the shared project."
                        print("message", message)
                    print("sending mail")
                    send_mail(subject, message, settings.DEFAULT_FROM_EMAIL, [shared_with_user.email], fail_silently=True)
                    print("mail sent")
                except Exception as e:
                    print("email not sent", e)
                    pass  # Do not block API on email failure
                return Response({
                    'success': True,
                    'message': 'Project shared successfully',
                    'share_id': share.id
                }, status=201)
                
        except Exception as e:
            return Response({
                'success': False,
                'message': f'Error sharing project: {str(e)}'
            }, status=500)


class RemoveProjectShare(APIView):
    """
    API endpoint to remove a project share.
    """
    def post(self, request):
        try:
            # Extract data from request
            share_id = request.data.get('share_id')
            project_id = request.data.get('project_id')
            shared_by_user_id = request.data.get('shared_by_user_id')
            shared_with_email = request.data.get('shared_with_email')
            
            # Validate required fields
            if not share_id and not all([project_id, shared_by_user_id, shared_with_email]):
                return Response({
                    'success': False,
                    'message': 'Either share_id or (project_id, shared_by_user_id, shared_with_email) are required'
                }, status=400)
            
            # Get the share to remove
            if share_id:
                try:
                    share = ProjectShare.objects.get(id=share_id, is_active=True)
                except ProjectShare.DoesNotExist:
                    return Response({
                        'success': False,
                        'message': 'Share not found'
                    }, status=404)
            else:
                try:
                    shared_by_user = User.objects.get(id=shared_by_user_id)
                    shared_with_user = User.objects.get(email=shared_with_email)
                    project = Projects.objects.get(id=project_id)
                except (User.DoesNotExist, Projects.DoesNotExist):
                    return Response({
                        'success': False,
                        'message': 'User or project not found'
                    }, status=404)
                
                try:
                    share = ProjectShare.objects.get(
                        project=project,
                        shared_by=shared_by_user,
                        shared_with=shared_with_user,
                        is_active=True
                    )
                except ProjectShare.DoesNotExist:
                    return Response({
                        'success': False,
                        'message': 'Share not found'
                    }, status=404)
            
            # Deactivate the share
            share.is_active = False
            share.updated_at = timezone.now()
            share.save()
            
            user = get_logging_user(request, getattr(share.project, 'user', None))
            ip = request.META.get('REMOTE_ADDR')
            log_user_action(user, "remove_project_share", details=f"Project share removed successfully", ip_address=ip)

            return Response({
                'success': True,
                'message': 'Project share removed successfully'
            }, status=200)
            
        except Exception as e:
            return Response({
                'success': False,
                'message': f'Error removing project share: {str(e)}'
            }, status=500)


class GetSharedProjects(APIView):
    """
    API endpoint to get all projects shared with a user.
    """
    def post(self, request):
        try:
            # Extract data from request
            user_id = request.data.get('user_id')
            
            if not user_id:
                return Response({
                    'success': False,
                    'message': 'user_id is required'
                }, status=400)
            
            # Get user
            try:
                user = User.objects.get(id=user_id)
            except User.DoesNotExist:
                return Response({
                    'success': False,
                    'message': 'User not found'
                }, status=404)
            
            # Get all active shares for this user
            shares = ProjectShare.objects.filter(
                shared_with=user,
                is_active=True
            ).select_related('project', 'shared_by').order_by('-created_at')
            
            shared_projects = []
            
            for share in shares:
                project_info = {
                    'share_id': share.id,
                    'project_id': share.project.id,
                    'project_name': share.project.name,
                    'share_type': share.share_type,
                    'permission_level': share.permission_level,
                    'shared_by': {
                        'id': share.shared_by.id,
                        'username': share.shared_by.username,
                        'email': share.shared_by.email
                    },
                    'created_at': share.created_at.isoformat(),
                    'updated_at': share.updated_at.isoformat()
                }
                
                # Add file-specific info if this is a file share
                if share.share_type == 'file':
                    project_info['shared_file'] = {
                        'file_type': share.file_type,
                        'file_name': share.file_name,
                        'sheet_name': share.sheet_name
                    }
                
                shared_projects.append(project_info)
            
            

            return Response({
                'success': True,
                'shared_projects': shared_projects,
                'total_count': len(shared_projects)
            }, status=200)
            
        except Exception as e:
            return Response({
                'success': False,
                'message': f'Error getting shared projects: {str(e)}'
            }, status=500)


class GetSharedProjectDetails(APIView):
    """
    API endpoint to get details of a shared project for a user.
    """
    def post(self, request):
        try:
            # Extract data from request
            user_id = request.data.get('user_id')
            project_id = request.data.get('project_id')
            share_id = request.data.get('share_id')
            
            if not user_id or not project_id:
                return Response({
                    'success': False,
                    'message': 'user_id and project_id are required'
                }, status=400)
            
            # Get user
            try:
                user = User.objects.get(id=user_id)
            except User.DoesNotExist:
                return Response({
                    'success': False,
                    'message': 'User not found'
                }, status=404)
            
            # Get the share
            share_filter = {
                'shared_with': user,
                'project_id': project_id,
                'is_active': True
            }
            
            if share_id:
                share_filter['id'] = share_id
            
            try:
                share = ProjectShare.objects.get(**share_filter)
            except ProjectShare.DoesNotExist:
                return Response({
                    'success': False,
                    'message': 'Project access not found or revoked'
                }, status=404)
            
            # Get project details
            project = share.project
            
            # Base project folder path
            project_folder = os.path.join(settings.MEDIA_ROOT, f"user_{project.user.id}/project_{project.id}")
            
            # Initialize response data structure
            project_details = {
                "project_id": project.id,
                "project_name": project.name,
                "share_info": {
                    "share_id": share.id,
                    "share_type": share.share_type,
                    "permission_level": share.permission_level,
                    "shared_by": {
                        "id": share.shared_by.id,
                        "username": share.shared_by.username,
                        "email": share.shared_by.email
                    },
                    "created_at": share.created_at.isoformat()
                },
                "files": {
                    "media": [],
                    "kpi": [],
                    "concatenated": []
                }
            }
            
            # Add file-specific info if this is a file share
            if share.share_type == 'file':
                project_details["shared_file"] = {
                    "file_type": share.file_type,
                    "file_name": share.file_name,
                    "sheet_name": share.sheet_name
                }
            
            # Get file details based on share type
            if share.share_type == 'project':
                # For project-level sharing, show all files
                for file_type in ["media", "kpi", "concatenated"]:
                    file_folder = os.path.join(project_folder, file_type)
                    if not os.path.exists(file_folder):
                        continue

                    # Get all files in the folder
                    for file_name in os.listdir(file_folder):
                        file_path = os.path.join(file_folder, file_name)
                        if os.path.isdir(file_path):  # Each file is a directory containing sheets
                            file_info = {
                                "name": file_name,
                                "sheets": []
                            }
                            
                            # Get all sheets (CSV files) in the file directory
                            for sheet_file in os.listdir(file_path):
                                if sheet_file.endswith('.csv'):
                                    sheet_path = os.path.join(file_path, sheet_file)
                                    sheet_info = {
                                        "name": sheet_file,
                                        "size": os.path.getsize(sheet_path),
                                        "last_modified": os.path.getmtime(sheet_path)
                                    }
                                    
                                    # Read complete sheet data
                                    try:
                                        df = pd.read_csv(sheet_path, dtype=str)
                                        df = df.replace([np.nan, np.inf, -np.inf], None)
                                        sheet_info["columns"] = df.columns.tolist()
                                        sheet_info["data"] = make_json_safe(df.values.tolist())
                                    except Exception as e:
                                        print(f"Error reading sheet {sheet_file}: {str(e)}")
                                        sheet_info["error"] = "Could not read sheet data"
                                    
                                    file_info["sheets"].append(sheet_info)
                            
                            project_details["files"][file_type].append(file_info)
            else:
                # For file-level sharing, show only the shared file
                file_type = share.file_type
                file_name = share.file_name
                
                file_folder = os.path.join(project_folder, file_type)
                file_path = os.path.join(file_folder, file_name)
                
                if os.path.exists(file_path) and os.path.isdir(file_path):
                    file_info = {
                        "name": file_name,
                        "sheets": []
                    }
                    
                    # Get all sheets (CSV files) in the file directory
                    for sheet_file in os.listdir(file_path):
                        if sheet_file.endswith('.csv'):
                            # For file-specific sharing, only show the shared sheet if specified
                            if share.sheet_name and sheet_file != share.sheet_name:
                                continue
                                
                            sheet_path = os.path.join(file_path, sheet_file)
                            sheet_info = {
                                "name": sheet_file,
                                "size": os.path.getsize(sheet_path),
                                "last_modified": os.path.getmtime(sheet_path)
                            }
                            
                            # Read complete sheet data
                            try:
                                df = pd.read_csv(sheet_path, dtype=str)
                                df = df.replace([np.nan, np.inf, -np.inf], None)
                                sheet_info["columns"] = df.columns.tolist()
                                sheet_info["data"] = make_json_safe(df.values.tolist())
                            except Exception as e:
                                print(f"Error reading sheet {sheet_file}: {str(e)}")
                                sheet_info["error"] = "Could not read sheet data"
                            
                            file_info["sheets"].append(sheet_info)
                    
                    project_details["files"][file_type].append(file_info)
            
            return Response({
                "success": True,
                "project_details": project_details
            }, status=200)
            
        except Exception as e:
            return Response({
                "success": False,
                "message": f"Error getting shared project details: {str(e)}"
            }, status=500)
    
    def get_sheet_names(self, file_path):
        """Helper method to get sheet names from Excel file"""
        try:
            import pandas as pd
            excel_file = pd.ExcelFile(file_path)
            return excel_file.sheet_names
        except Exception:
            return []


class GetSharedProjectPlots(APIView):
    """
    API endpoint to get plots/pivots for a shared project.
    """
    def post(self, request):
        try:
            # Extract data from request
            user_id = request.data.get('user_id')
            project_id = request.data.get('project_id')
            share_id = request.data.get('share_id')
            
            if not user_id or not project_id:
                return Response({
                    'success': False,
                    'message': 'user_id and project_id are required'
                }, status=400)
            
            # Get user
            try:
                user = User.objects.get(id=user_id)
            except User.DoesNotExist:
                return Response({
                    'success': False,
                    'message': 'User not found'
                }, status=404)
            
            # Get the share
            share_filter = {
                'shared_with': user,
                'project_id': project_id,
                'is_active': True
            }
            
            if share_id:
                share_filter['id'] = share_id
            
            try:
                share = ProjectShare.objects.get(**share_filter)
            except ProjectShare.DoesNotExist:
                return Response({
                    'success': False,
                    'message': 'Project access not found or revoked'
                }, status=404)
            
            # Get project
            project = share.project
            
            # Get plots and pivots based on share type
            plots = []
            pivots = []
            
            if share.share_type == 'project':
                # For project-level sharing, get all plots and pivots
                plots = SavedPlot.objects.filter(project=project).values()
                pivots = SavedPivot.objects.filter(project=project).values()
            else:
                # For file-level sharing, get only plots/pivots for the shared file
                plots = SavedPlot.objects.filter(
                    project=project,
                    file_type=share.file_type,
                    file_name=share.file_name
                ).values()
                
                if share.sheet_name:
                    plots = plots.filter(sheet_name=share.sheet_name)
                
                pivots = SavedPivot.objects.filter(
                    project=project,
                    file_type=share.file_type,
                    file_name=share.file_name
                ).values()
                
                if share.sheet_name:
                    pivots = pivots.filter(sheet_name=share.sheet_name)
            
            # Convert datetime fields to strings for JSON serialization
            for plot in plots:
                plot['created_at'] = plot['created_at'].isoformat()
                plot['updated_at'] = plot['updated_at'].isoformat()
            
            for pivot in pivots:
                pivot['created_at'] = pivot['created_at'].isoformat()
                pivot['updated_at'] = pivot['updated_at'].isoformat()
            
            return Response({
                'success': True,
                'plots': list(plots),
                'pivots': list(pivots),
                'share_info': {
                    'share_type': share.share_type,
                    'permission_level': share.permission_level
                }
            }, status=200)
            
        except Exception as e:
            return Response({
                'success': False,
                'message': f'Error getting shared project plots: {str(e)}'
            }, status=500)



class GetProjectSharedAccess(APIView):
    """
    API endpoint to get shared access information for a project that the user owns.
    """
    def post(self, request):
        try:
            # Extract data from request
            project_id = request.data.get('project_id')
            user_id = request.data.get('user_id')
            
            if not project_id or not user_id:
                return Response({
                    'success': False,
                    'message': 'project_id and user_id are required'
                }, status=400)
            
            # Get user
            try:
                user = User.objects.get(id=user_id)
            except User.DoesNotExist:
                return Response({
                    'success': False,
                    'message': 'User not found'
                }, status=404)
            
            # Get the project and verify ownership
            try:
                project = Projects.objects.get(id=project_id, user_id=user_id)
            except Projects.DoesNotExist:
                return Response({
                    'success': False,
                    'message': 'Project not found or access denied'
                }, status=404)
            
            # Get all active shares for this project
            shares = ProjectShare.objects.filter(
                project=project,
                is_active=True
            ).select_related('shared_with', 'shared_by')
            
            shared_access = []
            
            for share in shares:
                access_data = {
                    'share_id': share.id,
                    'shared_with': {
                        'id': share.shared_with.id,
                        'username': share.shared_with.username,
                        'email': share.shared_with.email
                    },
                    'shared_by': {
                        'id': share.shared_by.id,
                        'username': share.shared_by.username,
                        'email': share.shared_by.email
                    },
                    'share_type': share.share_type,
                    'permission_level': share.permission_level,
                    'created_at': share.created_at.isoformat(),
                    'updated_at': share.updated_at.isoformat()
                }
                
                if share.share_type == 'file':
                    access_data.update({
                        'file_type': share.file_type,
                        'file_name': share.file_name,
                        'sheet_name': share.sheet_name
                    })
                
                shared_access.append(access_data)
        

            return Response({
                'success': True,
                'shared_access': shared_access,
                'total_count': len(shared_access)
            }, status=200)
            
        except Exception as e:
            return Response({
                'success': False,
                'message': f'Error getting project shared access: {str(e)}'
            }, status=500)


class GetSheets(APIView):
    """
    API endpoint to get sheets from a specific file in a project.
    """
    def post(self, request):
        try:
            # Extract data from request
            file_type = request.data.get('file_type')
            file_name = request.data.get('file_name')
            file_id = request.data.get('file_id')
            project_id = request.data.get('project_id')
            user_id = request.data.get('user_id')
            permission_level = request.data.get('permission_level')

            # Validate required fields
            if not all([file_type, file_name, project_id, user_id]):
                missing_fields = []
                if not file_type: missing_fields.append('file_type')
                if not file_name: missing_fields.append('file_name')
                if not project_id: missing_fields.append('project_id')
                if not user_id: missing_fields.append('user_id')
                return Response({
                    'error': f'Missing required fields: {", ".join(missing_fields)}'
                }, status=400)

            # Get user
            try:
                user = User.objects.get(id=user_id)
            except User.DoesNotExist:
                return Response({'error': 'User not found'}, status=404)

            # Get project
            try:
                project = Projects.objects.get(id=project_id)
            except Projects.DoesNotExist:
                return Response({'error': 'Project not found'}, status=404)

            # Check if user has access to this project/file
            has_access, share_object, permission_level = check_project_access(
                user_id, project_id, file_type, file_name
            )
            
            if not has_access:
                return Response({
                    'error': 'Access denied. You don\'t have permission to access this project/file.'
                }, status=403)

            # Handle file name based on file type - CLEAN UP FIRST
            original_file_name = file_name
            if file_type == 'concatenated':
                # For concatenated files, use the folder name as the file name
                if '/' in file_name:
                    # If file_name includes the CSV name, extract just the folder name
                    file_name = file_name.split('/')[0]
                else:
                    file_name = os.path.basename(file_name)
            else:
                # For kpi and media files, extract just the file name from the path
                # Handle both forward and backward slashes
                if '/' in file_name or '\\' in file_name:
                    # Split by both forward and backward slashes and take the last part
                    file_name = file_name.replace('\\', '/').split('/')[-1]
                else:
                    file_name = os.path.basename(file_name)

            # Define folder structure
            project_folder = f"user_{project.user.id}/project_{project.id}"
            
            # Extract just the filenames from the database paths (handle both full paths and just filenames)
            def extract_filename(file_path):
                """Extract filename from either full path or just filename"""
                if '\\' in file_path or '/' in file_path:
                    return os.path.basename(file_path)
                return file_path
            
            last_name_kpi = [extract_filename(file) for file in project.kpi_file] if project.kpi_file else []
            last_name_media = [extract_filename(file) for file in project.media_file] if project.media_file else []
            last_name_concatenated = project.concatenated_file if hasattr(project, 'concatenated_file') and isinstance(project.concatenated_file, list) else []

            # Validate that the file exists in the project
            file_exists = False
            if file_name in last_name_kpi and file_type == 'kpi':
                file_exists = True
            elif file_name in last_name_media and file_type == 'media':
                file_exists = True
            elif file_type == 'concatenated':
                # For concatenated files, check if any folder contains this CSV file
                concatenated_folder = os.path.join(settings.MEDIA_ROOT, project_folder, "concatenated")
                if os.path.exists(concatenated_folder):
                    for folder_name in os.listdir(concatenated_folder):
                        folder_path = os.path.join(concatenated_folder, folder_name)
                        if os.path.isdir(folder_path):
                            # Check if this folder contains the CSV file
                            for csv_file in os.listdir(folder_path):
                                if csv_file == file_name or csv_file.endswith('.csv'):
                                    file_exists = True
                                    # Update file_name to the folder name for further processing
                                    file_name = folder_name
                                    break
                        if file_exists:
                            break
            elif file_name in last_name_concatenated and file_type == 'concatenated':
                file_exists = True
            
            if not file_exists:
                return Response({
                    'error': f'File "{file_name}" of type "{file_type}" not found in project. Available files: KPI={last_name_kpi}, Media={last_name_media}, Concatenated={last_name_concatenated}'
                }, status=404)

            # Build path based on file type
            if file_name in last_name_kpi and file_type == 'kpi':
                base_folder = os.path.join(settings.MEDIA_ROOT, project_folder, "kpi", file_name)
            elif file_name in last_name_media and file_type == 'media':
                base_folder = os.path.join(settings.MEDIA_ROOT, project_folder, "media", file_name)
            elif file_type == 'concatenated':
                # For concatenated files, use the folder name directly (which was updated during validation)
                base_folder = os.path.join(settings.MEDIA_ROOT, project_folder, "concatenated", file_name)
            elif file_name in last_name_concatenated and file_type == 'concatenated':
                base_folder = os.path.join(settings.MEDIA_ROOT, project_folder, "concatenated", file_name)
            else:
                return Response({'error': 'File not associated with the project'}, status=400)

            base_folder = os.path.normpath(base_folder)

            if not os.path.exists(base_folder):
                return Response({'error': 'Folder not found'}, status=404)

            # Collect all CSV files inside the folder
            csv_files = []
            for root, dirs, files in os.walk(base_folder):
                for file in files:
                    if file.endswith('.csv'):
                        csv_files.append(os.path.join(root, file))

            if not csv_files:
                return Response({'error': 'No CSV files found inside the folder'}, status=404)

            # Read and process each CSV file
            sheets_data = {}
            for csv_file in csv_files:
                try:
                    # Get relative path from base folder
                    relative_path = os.path.relpath(csv_file, base_folder)
                    sheet_name = os.path.splitext(relative_path)[0]  # Remove .csv extension
                    
                    # Read CSV file
                    try:
                        df = pd.read_csv(csv_file, encoding='utf-8')
                    except UnicodeDecodeError:
                        df = pd.read_csv(csv_file, encoding='latin1')
                    
                    # Replace infinite values with None
                    df = df.replace([np.inf, -np.inf], np.nan)
                    
                    # Convert DataFrame to JSON-safe format
                    safe_df = df.copy()
                    for col in safe_df.columns:
                        safe_df[col] = safe_df[col].astype(str)
                    
                    sheets_data[sheet_name] = {
                        'columns': safe_df.columns.tolist(),
                        'data': safe_df.values.tolist(),
                        'file_path': csv_file,
                        'size': os.path.getsize(csv_file),
                        'last_modified': os.path.getmtime(csv_file)
                    }
                    
                except Exception as e:
                    print(f"Error reading CSV file {csv_file}: {str(e)}")
                    continue

            if not sheets_data:
                return Response({'error': 'No valid sheets found in the file'}, status=404)

            return Response({
                'success': True,
                'file_info': {
                    'file_type': file_type,
                    'file_name': file_name,
                    'file_id': file_id,
                    'project_id': project_id,
                    'total_sheets': len(sheets_data),
                    'sheets': list(sheets_data.keys())
                },
                'sheets_data': sheets_data,
                'access_info': {
                    'permission_level': permission_level,
                    'is_owner': share_object is None
                }
            }, status=200)

        except Exception as e:
            return Response({
                'error': f'An unexpected error occurred: {str(e)}'
            }, status=500)


def get_logging_user(request, fallback_user=None):
    """
    Get user for logging purposes with priority order:
    1. user_id from request payload (highest priority)
    2. user_id from query parameters
    3. authenticated request.user
    4. fallback_user
    """
    # Priority 1: user_id from request payload
    user_id = None
    if hasattr(request, 'data') and isinstance(request.data, dict):
        user_id = request.data.get('user_id')
    
    # Priority 2: user_id from query parameters
    if not user_id and hasattr(request, 'query_params'):
        user_id = request.query_params.get('user_id')
    
    # If we found a user_id, try to get the user
    if user_id:
        try:
            return User.objects.get(id=user_id)
        except Exception:
            pass  # Continue to next priority if user not found
    
    # Priority 3: authenticated request.user
    if hasattr(request, 'user') and getattr(request.user, 'is_authenticated', False):
        return request.user
    
    # Priority 4: fallback_user
    if fallback_user is not None:
        return fallback_user
    
    return None

class GetUserEmails(APIView):
    def get(self, request):
        try:
            # Get all users and extract their email IDs
            users = User.objects.all()
            emails = [user.email for user in users]
            return Response({'emails': emails}, status=200)
        except Exception as e:
            return Response({'error': str(e)}, status=500)


class TestUserDetection(APIView):
    """Test endpoint to verify user detection in API logging"""
    def post(self, request):
        try:
            from .views import get_logging_user
            
            # Test user detection
            user = get_logging_user(request, None)
            
            # Test project-based user detection
            project_user = None
            if hasattr(request, 'data') and isinstance(request.data, dict):
                project_id = request.data.get('project_id')
                if project_id:
                    from .models import Projects
                    try:
                        project = Projects.objects.get(id=project_id)
                        project_user = project.user
                    except Projects.DoesNotExist:
                        pass
            
            return Response({
                'detected_user': {
                    'id': user.id if user else None,
                    'username': user.username if user else None,
                    'email': user.email if user else None
                } if user else None,
                'project_user': {
                    'id': project_user.id if project_user else None,
                    'username': project_user.username if project_user else None,
                    'email': project_user.email if project_user else None
                } if project_user else None,
                'request_data': request.data if hasattr(request, 'data') else {},
                'request_user': {
                    'id': request.user.id if hasattr(request, 'user') and request.user else None,
                    'username': request.user.username if hasattr(request, 'user') and request.user else None,
                    'email': request.user.email if hasattr(request, 'user') and request.user else None,
                    'is_authenticated': getattr(request.user, 'is_authenticated', False) if hasattr(request, 'user') and request.user else False
                } if hasattr(request, 'user') and request.user else None
            }, status=200)
            
        except Exception as e:
            return Response({'error': str(e)}, status=500)


class SparkSessionMonitor(APIView):
    """
    API endpoint to monitor and manage Spark sessions.
    Useful for debugging and monitoring Spark session health.
    """
    
    def get(self, request):
        """Get information about current Spark sessions."""
        try:
            from .spark_utils import get_spark_session_info
            session_info = get_spark_session_info()
            return Response(session_info, status=200)
        except Exception as e:
            return Response({'error': str(e)}, status=500)
    
    def post(self, request):
        """Test Spark session creation and validation."""
        try:
            from .spark_utils import get_spark_session, validate_spark_session, stop_all_spark_sessions
            
            action = request.data.get('action', 'test')
            
            if action == 'test':
                # Test session creation and validation
                spark = get_spark_session()
                is_valid = validate_spark_session(spark)
                
                return Response({
                    'status': 'success',
                    'session_created': True,
                    'session_valid': is_valid,
                    'spark_version': spark.version,
                    'message': 'Spark session test completed successfully'
                }, status=200)
            
            elif action == 'cleanup':
                # Stop all sessions
                stop_all_spark_sessions()
                
                return Response({
                    'status': 'success',
                    'message': 'All Spark sessions stopped successfully'
                }, status=200)
            
            else:
                return Response({
                    'error': 'Invalid action. Use "test" or "cleanup"'
                }, status=400)
                
        except Exception as e:
            return Response({
                'error': str(e),
                'status': 'failed'
            }, status=500)


class GetAPILogs(APIView):
    """
    API endpoint to retrieve API logs with filtering and pagination
    """
    def post(self, request):
        try:
            # Extract filter parameters
            user_id = request.data.get('user_id')
            project_id = request.data.get('project_id')
            endpoint = request.data.get('endpoint')
            method = request.data.get('method')
            response_status = request.data.get('response_status')
            start_date = request.data.get('start_date')
            end_date = request.data.get('end_date')
            page = request.data.get('page', 1)
            page_size = request.data.get('page_size', 50)
            
            # Import APILog model
            from .models import APILog
            from django.core.paginator import Paginator
            from django.db.models import Q
            from datetime import datetime
            
            # Build query
            queryset = APILog.objects.all()
            
            if user_id:
                # Get the username from User model and filter by user field
                try:
                    user = User.objects.get(id=user_id)
                    username = user.username or user.email
                    queryset = queryset.filter(user=username)
                except User.DoesNotExist:
                    # If user doesn't exist, return empty queryset
                    queryset = APILog.objects.none()
            
            if project_id:
                queryset = queryset.filter(project_id=project_id)
            
            if endpoint:
                queryset = queryset.filter(endpoint__icontains=endpoint)
            
            if method:
                queryset = queryset.filter(method__iexact=method)
            
            if response_status:
                queryset = queryset.filter(response_status=response_status)
            
            if start_date:
                try:
                    start_datetime = datetime.strptime(start_date, '%Y-%m-%d')
                    queryset = queryset.filter(request_timestamp__gte=start_datetime)
                except ValueError:
                    pass
            
            if end_date:
                try:
                    end_datetime = datetime.strptime(end_date, '%Y-%m-%d')
                    queryset = queryset.filter(request_timestamp__lte=end_datetime)
                except ValueError:
                    pass
            
            # Paginate results
            paginator = Paginator(queryset, page_size)
            page_obj = paginator.get_page(page)
            
            # Prepare response data
            logs_data = []
            for log in page_obj:
                logs_data.append({
                    'id': log.id,
                    'user': log.user,  # Now user is just a string (username)
                    'endpoint': log.endpoint,
                    'method': log.method,
                    'ip_address': log.ip_address,
                    'user_agent': log.user_agent,
                    'request_payload': log.request_payload,
                    'request_headers': log.request_headers,
                    'request_params': log.request_params,
                    'request_files': log.request_files,
                    'response_status': log.response_status,
                    'response_data': log.response_data,
                    'response_headers': log.response_headers,
                    'request_timestamp': log.request_timestamp.isoformat() if log.request_timestamp else None,
                    'response_timestamp': log.response_timestamp.isoformat() if log.response_timestamp else None,
                    'duration_ms': log.duration_ms,
                    'error_message': log.error_message,
                    'error_traceback': log.error_traceback,
                    'project_id': log.project_id,
                    'file_type': log.file_type,
                    'file_name': log.file_name,
                    'sheet_name': log.sheet_name,
                    'is_success': log.is_success,
                    'is_error': log.is_error
                })
            
            return Response({
                'logs': logs_data,
                'pagination': {
                    'current_page': page_obj.number,
                    'total_pages': paginator.num_pages,
                    'total_count': paginator.count,
                    'has_next': page_obj.has_next(),
                    'has_previous': page_obj.has_previous()
                }
            }, status=200)
            
        except Exception as e:
            return Response({'error': f'Failed to retrieve API logs: {str(e)}'}, status=500)


class GetAPILogStats(APIView):
    """
    API endpoint to get statistics about API logs
    """
    def post(self, request):
        try:
            from .models import APILog
            from django.db.models import Count, Avg, Q
            from datetime import datetime, timedelta
            
            # Get date range
            days = request.data.get('days', 7)
            end_date = datetime.now()
            start_date = end_date - timedelta(days=days)
            
            # Filter logs by date range
            queryset = APILog.objects.filter(
                request_timestamp__gte=start_date,
                request_timestamp__lte=end_date
            )
            
            # Calculate statistics
            total_requests = queryset.count()
            successful_requests = queryset.filter(response_status__gte=200, response_status__lt=300).count()
            error_requests = queryset.filter(Q(response_status__gte=400) | Q(error_message__isnull=False)).count()
            
            # Average response time
            avg_response_time = queryset.filter(duration_ms__isnull=False).aggregate(
                avg_duration=Avg('duration_ms')
            )['avg_duration'] or 0
            
            # Most common endpoints
            top_endpoints = queryset.values('endpoint').annotate(
                count=Count('id')
            ).order_by('-count')[:10]
            
            # Most common error status codes
            error_statuses = queryset.filter(response_status__gte=400).values('response_status').annotate(
                count=Count('id')
            ).order_by('-count')[:10]
            
            # Top users by request count
            top_users = queryset.filter(user__isnull=False).values(
                'user'
            ).annotate(
                count=Count('id')
            ).order_by('-count')[:10]
            
            # Method distribution
            method_distribution = queryset.values('method').annotate(
                count=Count('id')
            ).order_by('-count')
            
            return Response({
                'period': {
                    'start_date': start_date.isoformat(),
                    'end_date': end_date.isoformat(),
                    'days': days
                },
                'overview': {
                    'total_requests': total_requests,
                    'successful_requests': successful_requests,
                    'error_requests': error_requests,
                    'success_rate': (successful_requests / total_requests * 100) if total_requests > 0 else 0,
                    'avg_response_time_ms': round(avg_response_time, 2)
                },
                'top_endpoints': list(top_endpoints),
                'error_statuses': list(error_statuses),
                'top_users': list(top_users),
                'method_distribution': list(method_distribution)
            }, status=200)
            
        except Exception as e:
            return Response({'error': f'Failed to retrieve API stats: {str(e)}'}, status=500)


class DeleteAPILogs(APIView):
    """
    API endpoint to delete old API logs
    """
    def post(self, request):
        try:
            from .models import APILog
            from datetime import datetime, timedelta
            
            # Get parameters
            days_old = request.data.get('days_old', 30)
            delete_all = request.data.get('delete_all', False)
            
            if delete_all:
                # Delete all logs
                deleted_count = APILog.objects.count()
                APILog.objects.all().delete()
            else:
                # Delete logs older than specified days
                cutoff_date = datetime.now() - timedelta(days=days_old)
                deleted_count = APILog.objects.filter(
                    request_timestamp__lt=cutoff_date
                ).count()
                APILog.objects.filter(
                    request_timestamp__lt=cutoff_date
                ).delete()
            
            return Response({
                'message': f'Successfully deleted {deleted_count} API log entries',
                'deleted_count': deleted_count
            }, status=200)
            
        except Exception as e:
            return Response({'error': f'Failed to delete API logs: {str(e)}'}, status=500)


class TestAPILogging(APIView):
    """
    Test endpoint to verify API logging is working
    """
    def post(self, request):
        try:
            # Extract test data
            test_data = request.data.get('test_data', {})
            
            # Log the API call manually to ensure it works
            from .log_utils import log_api_request_from_view
            from django.utils import timezone
            
            response = Response({
                'message': 'API logging test successful',
                'received_data': test_data,
                'timestamp': timezone.now().isoformat()
            }, status=200)
            
            # Log the API call
            log_api_request_from_view(request, response, start_time=timezone.now())
            
            return response
            
        except Exception as e:
            return Response({'error': f'Test failed: {str(e)}'}, status=500)



